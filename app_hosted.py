import streamlit as st
import os
import datetime
import re
import streamlit.components.v1

st.set_page_config(
    page_title="PhysIQ — Science & English Tutor",
    page_icon="⚛️",
    layout="centered",
    initial_sidebar_state="expanded"
)

# ── Env vars ───────────────────────────────────────────────────
from dotenv import load_dotenv
load_dotenv()
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_ANON_KEY")
HF_TOKEN     = os.getenv("HF_TOKEN")

if not all([SUPABASE_URL, SUPABASE_KEY, HF_TOKEN]):
    st.error("❌ Missing environment variables. Check your .env file.")
    st.stop()

from supabase import create_client
supabase = create_client(SUPABASE_URL, SUPABASE_KEY)

def get_authed_client():
    client = create_client(SUPABASE_URL, SUPABASE_KEY)
    token = st.session_state.get("access_token")
    if token:
        client.auth.set_session(token, token)
    return client

# ── Session state ──────────────────────────────────────────────
defaults = {
    "user": None, "access_token": None, "messages": [],
    "pending_feedback": None, "backend_ready": False,
    "dark_mode": True, "show_landing": True,
    "simplify_target": None, "solver_result": None, "creative_mode": False, "coding_mode": False, "web_search_mode": False, "show_animator": False, "pdf_text": None, "pdf_name": None, "voice_mode": False, "voice_reply": None, "selected_voice": 0, "show_profile": False, "user_profile": None, "animation_data": None, "quiz_state": None, "quiz_active": False, "plugin_store_open": False, "voice_transcript": "", "_voice_question": None, "voice_clear_pending": False,
    "custom_plugins": {}, "tool_creator_html": None, "tool_creator_mode": None, "tool_creator_name": "",
    "visit_count": 0, "session_restored": False,
}
for k, v in defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v

# ── Theme CSS ──────────────────────────────────────────────────
def inject_theme():
    dark = st.session_state.dark_mode
    bg        = "#0d1117" if dark else "#f8f9fa"
    card_bg   = "#161b22" if dark else "#ffffff"
    text      = "#e6edf3" if dark else "#1a1a2e"
    subtext   = "#8b949e" if dark else "#6c757d"
    accent    = "#58a6ff" if dark else "#0066cc"
    border    = "#30363d" if dark else "#dee2e6"
    hero_grad = "linear-gradient(135deg,#0d1117 0%,#161b22 50%,#0d1117 100%)" if dark \
                else "linear-gradient(135deg,#e8f4f8 0%,#dbeafe 50%,#e8f4f8 100%)"
    btn_bg    = "#21262d" if dark else "#e9ecef"
    inp_bg    = "#0d1117" if dark else "#ffffff"

    st.markdown(f"""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=JetBrains+Mono:wght@400;600&display=swap');

    html, body, [class*="css"] {{
        font-family: 'Inter', sans-serif;
        background-color: {bg} !important;
        color: {text} !important;
    }}
    #MainMenu, footer, header {{ visibility: hidden; }}
    .block-container {{ padding-top: 1.5rem; max-width: 800px; }}

    /* ── Hero section ── */
    .hero {{
        background: {hero_grad};
        border: 1px solid {border};
        border-radius: 16px;
        padding: 48px 32px;
        text-align: center;
        margin-bottom: 28px;
    }}
    .hero-title {{
        font-size: 3rem;
        font-weight: 700;
        background: linear-gradient(90deg, #58a6ff, #79c0ff, #a5d6ff);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin: 0 0 8px 0;
        line-height: 1.1;
    }}
    .hero-sub {{
        font-size: 1.15rem;
        color: {subtext};
        margin-bottom: 28px;
    }}
    .feature-grid {{
        display: grid;
        grid-template-columns: repeat(2, 1fr);
        gap: 14px;
        margin: 24px 0;
    }}
    .feature-card {{
        background: {card_bg};
        border: 1px solid {border};
        border-radius: 12px;
        padding: 18px;
        text-align: left;
    }}
    .feature-icon {{ font-size: 1.6rem; margin-bottom: 8px; }}
    .feature-title {{ font-weight: 600; color: {text}; margin-bottom: 4px; }}
    .feature-desc {{ font-size: 0.85rem; color: {subtext}; }}

    /* ── Chat ── */
    .stChatMessage {{ background: {card_bg} !important; border: 1px solid {border} !important; border-radius: 12px !important; }}

    /* ── Solver box ── */
    .solver-box {{
        background: {card_bg};
        border: 1px solid {accent};
        border-radius: 12px;
        padding: 20px;
        margin: 12px 0;
    }}
    .solver-title {{ color: {accent}; font-weight: 600; font-size: 1rem; margin-bottom: 8px; }}

    /* ── Buttons ── */
    .stButton > button {{
        border-radius: 8px;
        border: 1px solid {border};
        background: {btn_bg};
        color: {text};
        font-family: 'Inter', sans-serif;
        transition: all 0.2s;
    }}
    .stButton > button:hover {{
        border-color: {accent};
        color: {accent};
    }}

    /* ── Inputs ── */
    .stTextInput > div > div > input,
    .stTextArea > div > div > textarea {{
        background: {inp_bg} !important;
        border: 1px solid {border} !important;
        color: {text} !important;
        border-radius: 8px !important;
    }}

    /* ── Tabs ── */
    .stTabs [data-baseweb="tab"] {{
        color: {subtext} !important;
    }}
    .stTabs [aria-selected="true"] {{
        color: {accent} !important;
    }}

    /* ── Sidebar ── */
    section[data-testid="stSidebar"] {{
        background: {card_bg} !important;
        border-right: 1px solid {border} !important;
    }}

    /* ── Confidence badge ── */
    .conf-high   {{ color: #3fb950; font-weight: 600; }}
    .conf-med    {{ color: #d29922; font-weight: 600; }}
    .conf-low    {{ color: #f85149; font-weight: 600; }}

    /* ── Divider ── */
    hr {{ border-color: {border} !important; }}

    /* ── Section headers ── */
    .section-label {{
        font-size: 0.72rem;
        font-weight: 600;
        letter-spacing: 1.5px;
        text-transform: uppercase;
        color: {subtext};
        margin-bottom: 10px;
    }}
    </style>
    """, unsafe_allow_html=True)

inject_theme()

# ══════════════════════════════════════════════════════════════
#  LANDING PAGE
# ══════════════════════════════════════════════════════════════
def show_landing():
    # Inject localStorage session restore on first load
    if not st.session_state.get("session_restored"):
        inject_session_persistence()
        st.session_state.session_restored = True
        # Check URL params for restored session
        try:
            params = st.query_params
            at = params.get("restore_at","")
            rt = params.get("restore_rt","")
            em = params.get("restore_em","")
            vc = int(params.get("restore_vc",0))
            if at and at != "":
                if restore_session_from_js(at, rt, em, "", "", vc):
                    try:
                        del st.query_params["restore_at"]
                        del st.query_params["restore_rt"]
                        del st.query_params["restore_em"]
                        del st.query_params["restore_vc"]
                    except: pass
                    st.rerun()
        except: pass

    # Theme toggle top-right
    col_sp, col_btn = st.columns([5,1])
    with col_btn:
        if st.button("🌙" if st.session_state.dark_mode else "☀️", help="Toggle dark/light mode"):
            st.session_state.dark_mode = not st.session_state.dark_mode
            st.rerun()

    st.markdown("""
    <div class="hero">
        <div class="hero-title">⚛️ PhysIQ</div>
        <div class="hero-sub">Your AI Science Tutor — Physics & Chemistry<br>
        Class 10 · Class 11 · Class 12 · College Level</div>
        <div class="feature-grid">
            <div class="feature-card">
                <div class="feature-icon">🧠</div>
                <div class="feature-title">Smart Answers</div>
                <div class="feature-desc">Powered by advanced AI with deep knowledge base</div>
            </div>
            <div class="feature-card">
                <div class="feature-icon">🔢</div>
                <div class="feature-title">Step-by-Step Solver</div>
                <div class="feature-desc">Solve any numerical problem with full working shown</div>
            </div>
            <div class="feature-card">
                <div class="feature-icon">🗣️</div>
                <div class="feature-title">Explain Simply</div>
                <div class="feature-desc">Any concept re-explained in simple everyday language</div>
            </div>
            <div class="feature-card">
                <div class="feature-icon">✍️</div>
                <div class="feature-title">Creative Writing</div>
                <div class="feature-desc">Essays, debates, letters, diary entries, stories &amp; more</div>
            </div>
            <div class="feature-card">
                <div class="feature-icon">📊</div>
                <div class="feature-title">Confidence Score</div>
                <div class="feature-desc">Know how confident the AI is in every answer</div>
            </div>
            <div class="feature-card">
                <div class="feature-icon">✍️</div>
                <div class="feature-title">English Writing</div>
                <div class="feature-desc">Essays, debates, letters, stories, poetry and more</div>
            </div>
            <div class="feature-card">
                <div class="feature-icon">🎨</div>
                <div class="feature-title">Creative Mode</div>
                <div class="feature-desc">Vivid, original creative writing with literary flair</div>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("---")

    tab_login, tab_signup = st.tabs(["🔑 Sign In", "📝 Create Account"])

    with tab_login:
        st.markdown('<div class="section-label">Sign in to your account</div>', unsafe_allow_html=True)
        email    = st.text_input("Email address", key="login_email", placeholder="you@example.com")
        password = st.text_input("Password", type="password", key="login_password", placeholder="••••••••")

        c1, c2 = st.columns(2)
        with c1:
            if st.button("Sign In →", use_container_width=True):
                if not email or not password:
                    st.error("Please enter email and password.")
                else:
                    try:
                        res = supabase.auth.sign_in_with_password({"email": email, "password": password})
                        st.session_state.user         = res.user
                        st.session_state.access_token = res.session.access_token
                        st.session_state.show_landing = False
                        st.toast("✅ Welcome back!")
                        # Save session to localStorage
                        save_session_to_js(res.user, res.session.access_token,
                            getattr(res.session,"refresh_token",""))
                        st.rerun()
                    except Exception as e:
                        st.error(f"❌ {e}")
        with c2:
            if st.button("Sign in with Google 🔵", use_container_width=True):
                try:
                    res = supabase.auth.sign_in_with_oauth({
                        "provider": "google",
                        "options": {"redirect_to": os.getenv("REDIRECT_URL", "http://localhost:8501")}
                    })
                    st.markdown(f"[Click here to sign in with Google →]({res.url})")
                except Exception as e:
                    st.error(f"❌ {e}")

    with tab_signup:
        st.markdown('<div class="section-label">Create your free account</div>', unsafe_allow_html=True)
        name       = st.text_input("Full Name", key="signup_name", placeholder="Your name")
        email_s    = st.text_input("Email", key="signup_email", placeholder="you@example.com")
        password_s = st.text_input("Password (min 6 chars)", type="password", key="signup_password")
        password_c = st.text_input("Confirm Password", type="password", key="signup_confirm")

        if st.button("Create Free Account →", use_container_width=True):
            if not all([name, email_s, password_s, password_c]):
                st.error("Please fill in all fields.")
            elif password_s != password_c:
                st.error("Passwords do not match.")
            elif len(password_s) < 6:
                st.error("Password must be at least 6 characters.")
            else:
                try:
                    res = supabase.auth.sign_up({
                        "email": email_s, "password": password_s,
                        "options": {"data": {"full_name": name}}
                    })
                    if res.user:
                        st.success("✅ Account created! Check your email to confirm, then sign in.")
                except Exception as e:
                    st.error(f"❌ {e}")

# ══════════════════════════════════════════════════════════════
#  DB HELPERS
# ══════════════════════════════════════════════════════════════
def get_user_id(): return st.session_state.user.id

def save_message(role, content, confidence=None):
    try:
        get_authed_client().table("messages").insert({
            "user_id": get_user_id(), "role": role,
            "content": content, "confidence": confidence,
            "created_at": datetime.datetime.utcnow().isoformat()
        }).execute()
    except Exception as e:
        st.warning(f"⚠️ Could not save: {e}")

def load_past_conversations():
    try:
        res = get_authed_client().table("messages").select("*")\
            .eq("user_id", get_user_id()).order("created_at").execute()
        return res.data or []
    except: return []

def save_learned(question, answer):
    try:
        get_authed_client().table("learned_answers").upsert({
            "user_id": get_user_id(), "question": question,
            "answer": answer, "created_at": datetime.datetime.utcnow().isoformat()
        }, on_conflict="user_id,question").execute()
    except Exception as e:
        st.warning(f"⚠️ Could not save: {e}")

def load_learned():
    try:
        res = get_authed_client().table("learned_answers").select("*")\
            .eq("user_id", get_user_id()).execute()
        return res.data or []
    except: return []

def delete_all_data():
    uid = get_user_id()
    try:
        get_authed_client().table("messages").delete().eq("user_id", uid).execute()
        get_authed_client().table("learned_answers").delete().eq("user_id", uid).execute()
    except Exception as e:
        st.error(f"❌ {e}")


# ══════════════════════════════════════════════════════════════
#  PERSISTENT LOGIN (localStorage-based, re-auth every 15 visits)
# ══════════════════════════════════════════════════════════════

VISITS_BEFORE_REAUTH = 15

def inject_session_persistence():
    """Inject JS to save/restore session from localStorage."""
    st.components.v1.html("""
<script>
(function(){
  // ── Restore saved session on load ──
  const saved = localStorage.getItem('physiq_session');
  if(saved){
    try{
      const d = JSON.parse(saved);
      // Send to Streamlit via query param hack
      if(d.access_token && d.refresh_token){
        window.parent.postMessage({
          type:'physiq_restore_session',
          access_token: d.access_token,
          refresh_token: d.refresh_token,
          user_email: d.user_email,
          user_id: d.user_id,
          user_name: d.user_name,
          visit_count: d.visit_count||0
        },'*');
      }
    }catch(e){}
  }

  // ── Listen for save-session from Streamlit ──
  window.addEventListener('message', function(e){
    if(e.data && e.data.type==='physiq_save_session'){
      const existing = JSON.parse(localStorage.getItem('physiq_session')||'{}');
      const vc = (existing.visit_count||0)+1;
      localStorage.setItem('physiq_session', JSON.stringify({
        access_token: e.data.access_token,
        refresh_token: e.data.refresh_token,
        user_email: e.data.user_email,
        user_id: e.data.user_id,
        user_name: e.data.user_name,
        visit_count: vc,
        saved_at: Date.now()
      }));
    }
    if(e.data && e.data.type==='physiq_clear_session'){
      localStorage.removeItem('physiq_session');
    }
    if(e.data && e.data.type==='physiq_restore_session'){
      // Already handled above on load
    }
  });
})();
</script>
""", height=0)

def restore_session_from_js(access_token, refresh_token, user_email, user_id, user_name, visit_count):
    """Try to restore session from saved token."""
    if st.session_state.get("user"):
        return True
    try:
        res = supabase.auth.set_session(access_token, refresh_token)
        if res and res.user:
            st.session_state.user = res.user
            st.session_state.access_token = access_token
            st.session_state.show_landing = False
            st.session_state.visit_count = visit_count
            return True
    except:
        pass
    return False

def save_session_to_js(user, access_token, refresh_token=""):
    """Tell JS to save session to localStorage."""
    name = ""
    try: name = user.user_metadata.get("full_name","")
    except: pass
    st.components.v1.html(f"""
<script>
window.parent.postMessage({{
  type:'physiq_save_session',
  access_token:{repr(access_token)},
  refresh_token:{repr(refresh_token or "")},
  user_email:{repr(user.email or "")},
  user_id:{repr(str(user.id))},
  user_name:{repr(name)},
}}, '*');
</script>
""", height=0)

def clear_session_from_js():
    """Tell JS to clear localStorage session."""
    st.components.v1.html("""
<script>window.parent.postMessage({type:'physiq_clear_session'},'*');</script>
""", height=0)


# ══════════════════════════════════════════════════════════════
#  BACKEND
# ══════════════════════════════════════════════════════════════
@st.cache_resource(show_spinner=False)
def load_backend():
    from langchain_community.vectorstores import FAISS
    from langchain_community.embeddings import HuggingFaceEmbeddings
    from langchain_community.docstore.document import Document
    emb = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2", model_kwargs={"device": "cpu"})
    vs  = FAISS.load_local("physics_index", emb, allow_dangerous_deserialization=True)
    return vs, Document

def call_hf(system_msg, user_msg, max_tokens=2000):
    """Core HuggingFace call — shared by all features."""
    from huggingface_hub import InferenceClient
    models = [
        "Qwen/Qwen2.5-7B-Instruct-Turbo",
        "Qwen/Qwen2.5-72B-Instruct-Turbo",
        "deepseek-ai/DeepSeek-V3",
    ]
    client = InferenceClient(base_url="https://router.huggingface.co/together/v1", api_key=HF_TOKEN)
    errors = []
    for model in models:
        try:
            resp = client.chat.completions.create(
                model=model,
                messages=[{"role":"system","content":system_msg},{"role":"user","content":user_msg}],
                max_tokens=max_tokens, temperature=0.3,
            )
            answer = resp.choices[0].message.content.strip()
            if answer:
                return answer
        except Exception as e:
            errors.append(f"- {model}: {str(e)[:100]}")
    st.error("All AI models failed:\n" + "\n".join(errors))
    return None

def get_confidence(results):
    if not results: return "🔴 Very Low", "conf-low"
    score = results[0][1]
    if score < 0.3:   return "🟢 High confidence",   "conf-high"
    elif score < 0.6: return "🟡 Medium confidence",  "conf-med"
    elif score < 1.0: return "🟠 Low confidence",     "conf-low"
    else:             return "🔴 Very Low",            "conf-low"

# ── Feature 1: Normal Ask ─────────────────────────────────────
def ask_question(question, vs, history_text):
    results = vs.similarity_search_with_score(question, k=4)
    conf_label, conf_class = get_confidence(results)
    context = "\n\n".join([r[0].page_content for r in results]) if results else ""

    # Detect if question is about English/writing
    english_keywords = ["essay", "debate", "letter", "diary", "story", "poem", "speech",
                        "report", "paragraph", "grammar", "vocabulary", "creative", "write",
                        "writing", "narrative", "describe", "argument", "persuasive", "formal",
                        "informal", "comprehension", "figure of speech", "metaphor", "simile",
                        "alliteration", "rebuttal", "thesis", "introduction", "conclusion"]
    is_english = any(w in question.lower() for w in english_keywords)

    if is_english:
        system = """You are an expert English Language and Literature teacher and creative writing coach.
You help students with:
- Essay writing (argumentative, descriptive, narrative, expository)
- Debate writing and speech writing
- Formal and informal letter writing
- Diary entries, creative stories, and poetry
- Grammar, vocabulary, and language analysis
- Comprehension and literary analysis

Always follow proper writing conventions and rules. When asked to write something:
1. Follow the correct FORMAT and STRUCTURE for that type of writing
2. Use appropriate LANGUAGE and TONE
3. Apply relevant LITERARY DEVICES where suitable
4. Show CREATIVITY while maintaining correctness

If asked to write an essay/letter/story etc., actually WRITE it — do not just explain how to write it.
Be creative, use vivid language, strong vocabulary, and varied sentence structures."""
    else:
        is_creative = st.session_state.get("creative_mode", False)
    is_coding = st.session_state.get("coding_mode", False)
    # Personalisation from user profile
    profile_ctx = build_personalisation_context(st.session_state.get("user_profile"))
    if is_creative:
        system = ("""You are a brilliant, creative English writing expert with the literary sensibility of a published author and the precision of an English professor.
You write with vivid imagery, original metaphors, varied sentence rhythms, and emotional intelligence.
When asked to write essays, debates, letters, stories, poems, diary entries, speeches, or any creative writing:
- Follow the correct format and rules for that genre meticulously
- Use sophisticated, precise vocabulary
- Vary sentence structure (short punchy sentences mixed with longer flowing ones)
- Include figurative language (original similes, metaphors, personification)
- Create atmosphere, mood, and emotional resonance
- Show don't tell
- Always explain the rules/structure you used after your writing
If asked about grammar, literature, or comprehension — explain clearly with examples.""") + profile_ctx
    else:
        system = """You are PhysIQ — an expert AI tutor covering Physics, Chemistry, Mathematics, English, Biology, and General Science from Class 10 to College level.

CRITICAL RULE — UNDERSTAND INTENT FIRST:
Before answering, decide what the user ACTUALLY wants:
1. EXPLANATION: "What is Newton's law?" / "Explain photosynthesis" / "How does DNA work?" → Give a clear educational explanation with examples and formulas.
2. CALCULATION: "Calculate the force when..." / "Find the pH of..." → Solve step-by-step with working shown.
3. DEFINITION: "What is entropy?" / "Define oxidation" → Give a precise definition with context.
4. COMPARISON: "Difference between..." / "Compare X and Y" → Structured comparison.
5. CONCEPT QUESTION: "Why does ice float?" / "How do planes fly?" → Intuitive explanation first, then technical detail.
6. CONVERSATIONAL: "Hello", "Thanks", "That was helpful" → Respond naturally and warmly.

NEVER write computer code unless the user EXPLICITLY asks for code (e.g. "write a Python program", "code for me").
If someone asks "how does Python handle memory?" — explain the concept, do NOT write Python code.
If someone asks "what is a for loop?" — explain with a simple example, do NOT write a full program.

FORMAT YOUR ANSWERS WELL:
- Use clear headings for complex topics
- Use bullet points for lists
- Show formulas clearly: F = ma, E = mc²
- Give real-world examples
- Keep answers focused and not too long unless detail is needed
- Be encouraging and friendly

If you don't know something or are unsure, say so honestly."""

    user = f"""Conversation so far:
{history_text if history_text else "(First question)"}

Relevant Knowledge:
{context}

Question: {question}

Answer:"""

    answer = call_hf(system, user)
    return answer or "Could not generate a response.", conf_label, conf_class

# ── Feature 2: Step-by-Step Numerical Solver ─────────────────
def solve_numerical(problem, vs):
    results = vs.similarity_search_with_score(problem, k=3) if vs else []
    context = "\n\n".join([r[0].page_content for r in results]) if results else ""

    system = """You are an expert Physics and Chemistry problem solver.
Solve the given numerical problem with COMPLETE step-by-step working.
Use this exact format:

**Given:**
(list all given values with units)

**To Find:**
(what is asked)

**Formula Used:**
(write the relevant formula)

**Solution:**
Step 1: (first step with calculation)
Step 2: (next step)
... (continue until answer)

**Answer:**
(final answer with correct units, in a box if possible)

**Key Concept:**
(one sentence explaining the concept tested)"""

    user = f"""Relevant formulas and concepts:
{context}

Problem to solve:
{problem}"""

    return call_hf(system, user, max_tokens=1500)

# ── Feature 3: Explain Simply ─────────────────────────────────
def explain_simply(original_answer, question):
    system = """You are a friendly science teacher explaining to a 12-year-old student.
Take the technical answer and re-explain it using:
- Simple everyday language (no jargon)
- A real-life analogy or example
- Short sentences
- Maybe an emoji or two to make it fun
Keep it under 150 words but make sure the key idea is clear."""

    user = f"""Original question: {question}

Technical answer:
{original_answer}

Now explain this simply, like you're talking to a curious 12-year-old:"""

    return call_hf(system, user, max_tokens=400)


# ── Feature 4: Animated Teaching ─────────────────────────────
def generate_animation_script(question, answer):
    """Ask AI to generate a JSON animation script for the canvas animator."""
    system = """You are a scientific animation director for an advanced educational platform.
Given a question and answer, produce a rich JSON animation scene using multiple objects.

CRITICAL: Return ONLY valid JSON. No markdown, no explanation, just JSON.

The JSON format uses a SCENE GRAPH — an array of objects you can freely combine and position:

{
  "subtitle": "Short description",
  "info_text": "Educational explanation shown at bottom (under 120 words)",
  "time_scale": 1.0,
  "bg_gradient": "#color or null",
  "legend": [{"color": "#hex", "label": "name"}],
  "steps": [
    {"title": "Step 1", "info": "What is happening in step 1..."},
    {"title": "Step 2", "info": "What is happening in step 2..."}
  ],
  "step_duration": 5,
  "connections": [
    {"from": 0, "to": 2, "type": "arrow", "color": "#color", "label": "label"}
  ],
  "objects": [
    { object1 },
    { object2 },
    ...many more objects as needed...
  ]
}

POSITIONING RULES — VERY IMPORTANT — read carefully:
- x goes from -1 (far left) to +1 (far right). Center is x=0.
- y goes from -0.9 (bottom) to +0.9 (top). Center is y=0.
- SPREAD objects out! Never put two objects at exactly the same position.
- Minimum distance between any two objects: 0.18 units.
- For atoms/molecules: put the main object at center (0,0), labels BELOW at y=-0.4
- For diagrams with left+right sides: use x=-0.55 and x=+0.55
- For top+bottom rows: use y=+0.3 and y=-0.3
- Formula boxes always go at y=-0.42 (near bottom, above info panel)
- Title labels always go at y=+0.42 (near top)
- Legends go to the right side, x=+0.7

LAYOUT TEMPLATES (use these patterns):

ATOM LAYOUT (example for Carbon):
- Atom at center: {"type":"atom","x":0,"y":0.05,"r":0.2,...}
- Title above: {"type":"label","x":0,"y":0.44,"text":"Carbon (C) — Z=6",...}
- Formula below: {"type":"formula_box","x":-0.3,"y":-0.43,"formula":"Shells: 2,4",...}
- Energy level (right): {"type":"formula_box","x":0.4,"y":0,"formula":"Valency = 4",...}

LEFT-RIGHT LAYOUT (two things to compare):
- Left object: x=-0.45, y=0
- Right object: x=+0.45, y=0  
- Label left: x=-0.45, y=-0.22
- Label right: x=+0.45, y=-0.22
- Arrow between: from x=-0.2 to x=+0.2

TOP-BOTTOM LAYOUT (cause and effect):
- Top object: x=0, y=+0.3
- Arrow down: ty=0 (center)
- Bottom object: x=0, y=-0.15

MULTI-OBJECT LAYOUT (4-6 things):
- Top left: x=-0.45, y=+0.28
- Top right: x=+0.45, y=+0.28
- Middle left: x=-0.45, y=-0.08
- Middle right: x=+0.45, y=-0.08
- Formula box bottom: x=0, y=-0.43

CLARITY RULES:
1. Each object must have a "caption" or nearby label so user knows what it is
2. Never put a label on top of another object — offset labels by at least 0.12 units
3. Use formula_box for equations — place at bottom (y=-0.42)
4. Use "steps" array to guide the user through the animation step by step
5. Keep legends minimal — only include if there are 3+ different colored things
6. The info_text at the bottom should explain WHAT IS HAPPENING right now

Return ONLY the JSON object."""

    user = f"""Question: {question}

Answer given: {answer[:500]}

Generate the animation JSON for this concept:"""

    import json
    raw = call_hf(system, user, max_tokens=800)
    if not raw:
        return None
    # Extract JSON from response
    raw = raw.strip()
    start = raw.find('{')
    end   = raw.rfind('}') + 1
    if start == -1 or end == 0:
        return None
    try:
        data = json.loads(raw[start:end])
        # Ensure it has objects array (new format)
        if 'objects' not in data and 'type' in data:
            # Convert old format to new scene graph format
            data = convert_old_to_new(data)
        return data
    except Exception as e:
        return None

def convert_old_to_new(old):
    """Convert legacy animation format to new scene-graph format."""
    t = old.get('type','atom')
    objects = []
    subtitle = old.get('subtitle','Animation')
    info = old.get('info_text','')
    
    if t=='atom':
        objects.append({"type":"atom","x":0,"y":0.05,"r":0.2,
            "protons":old.get('protons',1),"neutrons":old.get('neutrons',0),
            "shells":old.get('shells',[[1]])})
        objects.append({"type":"label","x":0,"y":-0.4,"text":subtitle,"color":"#58a6ff","size":14})
    elif t=='wave':
        for i,w in enumerate(old.get('waves',[])):
            objects.append({"type":"wave","x":0,"y":w.get('y_offset',0)*0.3,
                "width":0.85,"amp":w.get('amplitude',60)/400,
                "wavelength":w.get('wavelength',150)/400,"speed":w.get('speed',0.03)*2000,
                "color":w.get('color','#44ffaa'),"label_top":w.get('label','')})
    else:
        objects.append({"type":"label","x":0,"y":0,"text":subtitle,"color":"#58a6ff","size":16})
        objects.append({"type":"formula_box","x":0,"y":-0.25,"formula":info[:40] if info else t,
            "color":"#4da6ff","width":0.5,"height":0.12})
    
    return {"subtitle":subtitle,"info_text":info,"objects":objects,
            "legend":old.get('legend',[]),"bg_gradient":None}

def show_animator_panel(animation_data):
    """Render the HTML5 Canvas animator in Streamlit."""
    import json
    with open("animator.html", "r") as f:
        html_content = f.read()
    # Inject data directly into the HTML
    inject_script = f"""
<script>
window.addEventListener('load', function() {{
  try {{
    const data = {json.dumps(animation_data)};
    scene = data;
    if(scene._particles) delete scene._particles;
    document.getElementById('loading').style.display='none';
    const subtitleEl = document.getElementById('subtitle');
    const infoEl = document.getElementById('info-text');
    if (subtitleEl) subtitleEl.textContent = scene.subtitle || scene.type || 'Animation';
    if (infoEl) infoEl.textContent = scene.info_text || '';
    if(scene.legend) updateLegend(scene.legend);
  }} catch(e) {{ console.error(e); }}
}}, 500);
</script>"""
    html_content = html_content.replace('</body>', inject_script + '</body>')
    st.components.v1.html(html_content, height=520, scrolling=False)



# ══════════════════════════════════════════════════════════════
#  QUIZ CONNECTOR
# ══════════════════════════════════════════════════════════════

QUIZ_TRIGGER_WORDS = [
    "quiz", "test me", "question me", "mcq", "multiple choice",
    "fill in the blank", "practice questions", "exam questions",
    "test my knowledge", "ask me questions", "give me a quiz",
    "quiz me", "make a quiz", "create a quiz", "generate a quiz"
]

def is_quiz_request(message):
    """Detect if user is asking for a quiz."""
    msg = message.lower()
    return any(word in msg for word in QUIZ_TRIGGER_WORDS)

def extract_quiz_topic(message):
    """Pull topic and type from user message."""
    msg = message.lower()
    q_type = "mixed"
    if "mcq" in msg or "multiple choice" in msg:
        q_type = "mcq"
    elif "fill" in msg or "blank" in msg:
        q_type = "fill"
    # Try to extract number
    import re
    nums = re.findall(r"\b(\d+)\b", msg)
    num_q = int(nums[0]) if nums else 5
    num_q = min(max(num_q, 3), 15)
    return q_type, num_q

def generate_quiz(topic, q_type, num_q):
    """Generate quiz questions using the AI."""
    import json
    if q_type == "mcq":
        type_instruction = f"Generate exactly {num_q} Multiple Choice Questions (MCQ)."
        schema_note = 'For MCQ: "type":"mcq", "options":["A)...","B)...","C)...","D)..."], "correct":"A"'
    elif q_type == "fill":
        type_instruction = f"Generate exactly {num_q} Fill in the Blank questions."
        schema_note = 'For fill: "type":"fill", "answer":"exact answer"'
    else:
        half = num_q // 2
        type_instruction = f"Generate {half} MCQ and {num_q - half} Fill in the Blank questions."
        schema_note = 'For MCQ: "type":"mcq", "options":["A)...","B)...","C)...","D)..."], "correct":"A". For fill: "type":"fill", "answer":"exact answer"'

    system = f"""You are an expert science and English quiz maker.
{type_instruction}
Return ONLY a valid JSON array. No explanation, no markdown, just the JSON.

Each question object must have:
- "q": the question text (for fill-in-blank, use _______ where the blank is)
- "type": "mcq" or "fill"
- "correct": the correct answer (letter for MCQ like "A", or the word/phrase for fill)
- "explanation": 2-3 sentence explanation of why this is the answer
- "difficulty": "easy", "medium", or "hard"
{schema_note}

Make questions varied in difficulty. Cover the topic thoroughly.
Example fill: {{"q":"Newton's second law states F = m × _______","type":"fill","correct":"a","explanation":"F=ma where a is acceleration. Force equals mass times acceleration.","difficulty":"easy"}}
Example MCQ: {{"q":"What is the SI unit of force?","type":"mcq","options":["A) Joule","B) Newton","C) Watt","D) Pascal"],"correct":"B","explanation":"The Newton (N) is the SI unit of force. 1 N = 1 kg⋅m/s².","difficulty":"easy"}}
"""

    user = f"Generate a quiz about: {topic}"
    raw = call_hf(system, user, max_tokens=2000)
    if not raw:
        return None
    raw = raw.strip()
    start = raw.find('[')
    end = raw.rfind(']') + 1
    if start == -1 or end == 0:
        # Try single object wrapped in array
        start = raw.find('{')
        if start == -1:
            return None
        raw = '[' + raw[start:] + ']'
        end = len(raw)
    try:
        questions = json.loads(raw[start:end])
        return questions[:num_q]
    except:
        return None

def save_quiz_mistakes(mistakes_data):
    """Save quiz mistakes to Supabase for the logged-in user."""
    if not st.session_state.get("user"):
        return
    try:
        import json
        existing = []
        try:
            res = get_authed_client().table("quiz_mistakes").select("*").eq("user_id", get_user_id()).execute()
            existing = res.data or []
        except:
            pass
        all_mistakes = existing + mistakes_data
        # Keep last 50 mistakes only
        all_mistakes = all_mistakes[-50:]
        # Delete old and insert fresh
        try:
            get_authed_client().table("quiz_mistakes").delete().eq("user_id", get_user_id()).execute()
        except:
            pass
        for m in all_mistakes:
            try:
                get_authed_client().table("quiz_mistakes").insert({
                    "user_id": get_user_id(),
                    "question": m.get("q",""),
                    "user_answer": m.get("user_answer",""),
                    "correct_answer": m.get("correct",""),
                    "explanation": m.get("explanation",""),
                    "topic": m.get("topic",""),
                    "created_at": datetime.datetime.utcnow().isoformat()
                }).execute()
            except:
                pass
    except Exception as e:
        st.warning(f"Could not save mistakes: {e}")

def load_quiz_mistakes():
    """Load past mistakes from Supabase."""
    if not st.session_state.get("user"):
        return []
    try:
        res = get_authed_client().table("quiz_mistakes").select("*").eq("user_id", get_user_id()).order("created_at", desc=True).limit(30).execute()
        return res.data or []
    except:
        return []

def generate_mistake_lesson(mistakes):
    """AI generates a personalised lesson based on mistakes."""
    if not mistakes:
        return None
    topics = list(set([m.get("topic","general") for m in mistakes[:10]]))
    wrong_qs = [f"Q: {m.get('question','?')} | Your answer: {m.get('user_answer','?')} | Correct: {m.get('correct_answer','?')}" for m in mistakes[:10]]
    system = """You are a patient, encouraging science and English tutor.
Based on the student's quiz mistakes, write a personalised lesson that:
1. Identifies the pattern of mistakes (which concepts are weak)
2. Clearly explains each weak concept with examples
3. Gives memory tips and mnemonics
4. Ends with 2-3 practice tips
Use friendly, encouraging language. Format with clear sections."""
    user = f"""Topics tested: {', '.join(topics)}

Student's mistakes:
{chr(10).join(wrong_qs)}

Write a personalised lesson to fix these mistakes:"""
    return call_hf(system, user, max_tokens=1500)

def show_celebration():
    """Show celebration animation for correct answer."""
    st.markdown("""
    <div id="celebration" style="text-align:center;padding:8px;animation:pop 0.5s ease-out;">
        <span style="font-size:2rem;">🎉</span>
        <span style="color:#3fb950;font-weight:700;font-size:1.1rem;"> Correct! Well done!</span>
        <span style="font-size:2rem;">⭐</span>
    </div>
    <style>
    @keyframes pop {
        0%{transform:scale(0.5);opacity:0}
        70%{transform:scale(1.15)}
        100%{transform:scale(1);opacity:1}
    }
    </style>
    """, unsafe_allow_html=True)

def show_quiz_ui():
    """Render the interactive quiz panel."""
    qs = st.session_state.quiz_state
    if not qs:
        return

    questions   = qs["questions"]
    current_idx = qs.get("current_idx", 0)
    score       = qs.get("score", 0)
    answers     = qs.get("answers", {})
    skipped     = qs.get("skipped", [])
    phase       = qs.get("phase", "quiz")  # quiz | review | results | lesson

    topic = qs.get("topic","Science")
    total = len(questions)

    st.markdown("---")

    # ── PHASE: QUIZ ───────────────────────────────────────────
    if phase == "quiz" and current_idx < total:
        q = questions[current_idx]
        prog = current_idx / total
        st.progress(prog, text=f"Question {current_idx+1} of {total} | ⭐ Score: {score}/{current_idx}")

        # Difficulty badge
        diff = q.get("difficulty","medium")
        diff_color = "#3fb950" if diff=="easy" else "#d29922" if diff=="medium" else "#f85149"
        st.markdown(f'<span style="background:{diff_color};color:#fff;padding:2px 10px;border-radius:12px;font-size:11px;font-weight:600">{diff.upper()}</span>', unsafe_allow_html=True)

        st.markdown(f"### Q{current_idx+1}. {q['q']}")

        q_key = f"quiz_q_{current_idx}"

        if q["type"] == "mcq":
            options = q.get("options", [])
            user_ans = st.radio("Choose your answer:", options, key=q_key, label_visibility="collapsed")
            col1, col2 = st.columns([1,4])
            with col1:
                if st.button("Submit →", key=f"submit_{current_idx}", use_container_width=True):
                    # Extract letter from option like "A) Newton"
                    chosen_letter = user_ans[0] if user_ans else ""
                    correct_letter = q.get("correct","A")
                    is_correct = chosen_letter.upper() == correct_letter.upper()
                    answers[current_idx] = {"q": q["q"], "user_answer": user_ans,
                                            "correct": q["correct"], "is_correct": is_correct,
                                            "explanation": q.get("explanation",""),
                                            "topic": topic}
                    if is_correct:
                        qs["score"] = score + 1
                    else:
                        skipped.append(current_idx)
                    qs["answers"] = answers
                    qs["skipped"] = skipped
                    qs["current_idx"] = current_idx + 1
                    qs["last_correct"] = is_correct
                    st.session_state.quiz_state = qs
                    st.rerun()

        else:  # fill in blank
            user_ans = st.text_input("Type your answer:", key=q_key, placeholder="Your answer here...")
            col1, col2 = st.columns([1,4])
            with col1:
                if st.button("Submit →", key=f"submit_{current_idx}", use_container_width=True):
                    correct = q.get("correct","")
                    is_correct = user_ans.strip().lower() == correct.strip().lower()
                    # Also allow partial match for longer answers
                    if not is_correct and len(correct) > 3:
                        is_correct = correct.strip().lower() in user_ans.strip().lower()
                    answers[current_idx] = {"q": q["q"], "user_answer": user_ans,
                                            "correct": correct, "is_correct": is_correct,
                                            "explanation": q.get("explanation",""),
                                            "topic": topic}
                    if is_correct:
                        qs["score"] = score + 1
                    else:
                        skipped.append(current_idx)
                    qs["answers"] = answers
                    qs["skipped"] = skipped
                    qs["current_idx"] = current_idx + 1
                    qs["last_correct"] = is_correct
                    st.session_state.quiz_state = qs
                    st.rerun()

        # Show result of last answer
        if current_idx > 0 and qs.get("last_correct") is not None:
            if qs["last_correct"]:
                show_celebration()
            else:
                st.error("❌ Incorrect — question saved for review at the end")
        if st.button("⏭️ Skip this question", key=f"skip_{current_idx}"):
            answers[current_idx] = {"q": q["q"], "user_answer": "(skipped)",
                                    "correct": q.get("correct",""), "is_correct": False,
                                    "explanation": q.get("explanation",""), "topic": topic}
            skipped.append(current_idx)
            qs["answers"] = answers
            qs["skipped"] = skipped
            qs["current_idx"] = current_idx + 1
            qs["last_correct"] = False
            st.session_state.quiz_state = qs
            st.rerun()

    # ── PHASE: REVIEW WRONG ANSWERS ──────────────────────────
    elif phase == "quiz" and current_idx >= total:
        qs["phase"] = "review"
        st.session_state.quiz_state = qs
        st.rerun()

    elif phase == "review":
        score = qs.get("score",0)
        pct = round(score / total * 100)
        st.markdown(f"## 🏁 Quiz Complete!")

        # Score display
        score_color = "#3fb950" if pct >= 70 else "#d29922" if pct >= 40 else "#f85149"
        st.markdown(f"""
        <div style="text-align:center;padding:20px;background:rgba(0,0,0,0.3);border-radius:12px;margin:12px 0">
            <div style="font-size:3.5rem;font-weight:700;color:{score_color}">{pct}%</div>
            <div style="color:#8b949e;font-size:1rem">You got {score} out of {total} correct</div>
            <div style="font-size:1.5rem;margin-top:8px">{"🌟 Excellent!" if pct>=80 else "👍 Good effort!" if pct>=50 else "📚 Keep practising!"}</div>
        </div>
        """, unsafe_allow_html=True)

        # Wrong answers section
        wrong = [a for a in answers.values() if not a.get("is_correct")]
        if wrong:
            st.markdown(f"### ❌ {len(wrong)} Questions to Review")
            for i, ans in enumerate(wrong, 1):
                with st.expander(f"❌ Q{i}: {ans['q'][:60]}..."):
                    st.markdown(f"**Your answer:** {ans.get('user_answer','(skipped)')}")
                    st.markdown(f"**✅ Correct answer:** {ans.get('correct','')}")
                    st.markdown(f"**📖 Explanation:** {ans.get('explanation','')}")

            # Save mistakes to Supabase
            save_quiz_mistakes(wrong)
            st.success(f"✅ {len(wrong)} mistakes saved to your account for personalised learning!")

            col1, col2, col3 = st.columns(3)
            with col1:
                if st.button("🎓 Get Personalised Lesson", use_container_width=True):
                    qs["phase"] = "lesson"
                    st.session_state.quiz_state = qs
                    st.rerun()
            with col2:
                if st.button("🔄 Try Again", use_container_width=True):
                    st.session_state.quiz_active = False
                    st.session_state.quiz_state = None
                    st.rerun()
            with col3:
                if st.button("✕ Close Quiz", use_container_width=True):
                    st.session_state.quiz_active = False
                    st.session_state.quiz_state = None
                    st.rerun()
        else:
            st.balloons()
            st.success("🎉 Perfect score! You got everything right!")
            if st.button("✕ Close Quiz", use_container_width=True):
                st.session_state.quiz_active = False
                st.session_state.quiz_state = None
                st.rerun()

    # ── PHASE: PERSONALISED LESSON ────────────────────────────
    elif phase == "lesson":
        wrong = [a for a in answers.values() if not a.get("is_correct")]
        st.markdown("## 🎓 Your Personalised Lesson")
        st.caption("Based on your quiz mistakes, the AI has created a custom lesson just for you.")
        with st.spinner("📚 Generating your personalised lesson..."):
            lesson = generate_mistake_lesson(wrong)
        if lesson:
            st.markdown(lesson)
        else:
            st.warning("Could not generate lesson. Please try again.")
        col1, col2 = st.columns(2)
        with col1:
            if st.button("🔄 Take Another Quiz", use_container_width=True):
                st.session_state.quiz_active = False
                st.session_state.quiz_state = None
                st.rerun()
        with col2:
            if st.button("✕ Close", use_container_width=True):
                st.session_state.quiz_active = False
                st.session_state.quiz_state = None
                st.rerun()



# ══════════════════════════════════════════════════════════════
#  CODING CONNECTOR
# ══════════════════════════════════════════════════════════════

CODING_LANGUAGES = {
    "python":"Python","java":"Java","c++":"C++","cpp":"C++",
    "c#":"C#","csharp":"C#","javascript":"JavaScript","js":"JavaScript",
    "typescript":"TypeScript","ts":"TypeScript","roblox":"Roblox Lua",
    "lua":"Lua","rust":"Rust","go":"Go","golang":"Go","swift":"Swift",
    "kotlin":"Kotlin","sql":"SQL","html":"HTML","css":"CSS",
    "react":"React/JSX","php":"PHP","ruby":"Ruby","bash":"Bash/Shell",
    "r":"R","matlab":"MATLAB","dart":"Dart","flutter":"Flutter",
}

# Phrases that STRONGLY indicate user wants code written
CODING_WRITE_TRIGGERS = [
    "write code", "write a program", "write a script", "write a function",
    "write a class", "write a game", "write an app", "write a bot",
    "create a program", "create a script", "create a game", "create an app",
    "build a program", "build a game", "build an app", "build a website",
    "make a program", "make a game", "make a script", "make an app",
    "code for me", "generate code", "can you code", "write me code",
    "implement a", "develop a", "program a", "give me the code",
    "show me the code", "write the code", "full code", "complete code",
]

# Words that just mention coding but user wants EXPLANATION, not code
CODING_EXPLAIN_TRIGGERS = [
    "what is", "explain", "how does", "difference between", "what are",
    "define", "tell me about", "describe", "why does", "how do",
    "what does", "meaning of", "concept of", "theory", "example of",
    "when to use", "pros and cons", "advantages", "disadvantages",
]

def is_coding_request(message):
    """Only return True when user CLEARLY wants code written, not just asking a question."""
    if st.session_state.get("coding_mode", False):
        return True
    msg = message.lower().strip()

    # Strong action words that mean "write me code"
    strong_actions = [
        "write code", "write a program", "write a script", "write a function",
        "write a class", "create a program", "create an app", "build a program",
        "build an app", "build a game", "code for me", "generate code",
        "implement a", "develop a", "program that", "write me code",
        "give me code", "make a program", "make a game", "can you code",
        "write the code", "create the code", "build the code",
        "make a script", "write a bot", "create a bot",
    ]
    # Library names that only appear in coding contexts
    strong_libraries = [
        "ursina", "pygame", "tkinter", "kivy", "flask app", "fastapi",
        "django app", "tensorflow model", "pytorch model", "react component",
        "express server", "spring boot", "roblox script", "arduino code",
    ]

    if any(a in msg for a in strong_actions):
        return True
    if any(lib in msg for lib in strong_libraries):
        return True

    # Language + explicit code word combo
    lang_keys = list(CODING_LANGUAGES.keys())
    has_lang = any(f" {lang} " in f" {msg} " for lang in lang_keys)
    has_code_word = any(w in msg for w in ["code", "program", "script", "function", "class", "syntax"])
    return has_lang and has_code_word


def detect_language(message):
    msg = message.lower()
    for key, lang in CODING_LANGUAGES.items():
        if key in msg:
            return lang
    return "Python"  # default

def generate_code(question, vs, language="Python"):
    results = vs.similarity_search_with_score(question, k=3) if vs else []
    context = "\n\n".join([r[0].page_content for r in results]) if results else ""

    system = f"""You are an elite {language} programmer with 15+ years of experience.
You write clean, efficient, well-commented, production-quality code.
RULES:
1. Write COMPLETE, WORKING code — never truncate or use placeholder comments like "# add code here"
2. Include ALL necessary imports at the top
3. Add clear comments explaining complex logic
4. Follow best practices and design patterns for {language}
5. If the user asks for a game/app, write the FULL implementation
6. You can write up to 2000 lines if needed — never cut short
7. After the code, write a brief explanation of how it works
8. If using a library (ursina, pygame, pandas etc), use it correctly with real API calls
9. Make the code actually runnable with zero modifications needed
10. For Roblox Lua, use proper Roblox API (Services, Instances, Events)

LIBRARIES YOU KNOW PERFECTLY:
Python: numpy, pandas, matplotlib, seaborn, scikit-learn, tensorflow, pytorch,
        flask, fastapi, django, sqlalchemy, pygame, ursina, opencv, requests,
        beautifulsoup4, pillow, pydantic, asyncio, aiohttp, celery, redis,
        plotly, dash, streamlit, gradio, transformers, langchain, openai
Java: Spring Boot, Hibernate, Maven, Gradle, JUnit, Mockito
C++: STL, Boost, OpenGL, SDL2, SFML, CMake
JavaScript: React, Vue, Angular, Express, Next.js, Node.js, TypeScript
Roblox Lua: All Roblox services, DataStore, TweenService, RemoteEvents
Rust: tokio, serde, actix-web, reqwest
Go: gin, echo, gorm, cobra

Write the code now:"""

    user = f"""Relevant documentation:
{context}

Task: {question}

Write complete, working {language} code:"""

    return call_hf(system, user, max_tokens=30000)



# ══════════════════════════════════════════════════════════════
#  PDF READER
# ══════════════════════════════════════════════════════════════

def extract_pdf_text(pdf_file):
    """Extract text from uploaded PDF file."""
    try:
        import io
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(pdf_file.read()))
        text = ""
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text += f"\n--- Page {i+1} ---\n{page_text}"
        return text.strip()
    except Exception as e:
        return None

def answer_from_pdf(pdf_text, question):
    """Use AI to answer a question based on PDF content."""
    # Trim PDF text to avoid token overflow
    trimmed = pdf_text[:4000] if len(pdf_text) > 4000 else pdf_text
    system = """You are an expert tutor. A student has uploaded a PDF document and asked a question about it.
Read the document carefully and answer the question accurately and clearly.
If the answer is not in the document, say so honestly.
Format your answer with clear sections if needed."""
    user = f"""PDF Content:
{trimmed}

Student's question: {question}

Answer based on the PDF:"""
    return call_hf(system, user, max_tokens=1500)

def summarise_pdf(pdf_text):
    """Generate a structured summary of a PDF."""
    trimmed = pdf_text[:5000] if len(pdf_text) > 5000 else pdf_text
    system = """You are an expert academic summariser.
Create a clear, structured summary of this document with:
- Main topic and purpose
- Key concepts covered  
- Important facts, formulas, or findings
- Conclusions or takeaways
Use bullet points and headers for clarity."""
    user = f"""Document to summarise:\n{trimmed}\n\nProvide a structured summary:"""
    return call_hf(system, user, max_tokens=1200)



# ══════════════════════════════════════════════════════════════
#  VOICE CONNECTOR
# ══════════════════════════════════════════════════════════════

def generate_voice_answer(question, full_answer):
    """Generate a short, spoken-friendly version of the answer."""
    system = """You are a friendly science tutor speaking out loud.
Convert this answer into a SHORT spoken response (under 60 words).
Rules:
- Speak naturally, like talking to a student face-to-face
- No markdown, no bullet points, no symbols like * or #
- No formulas with special characters — say them in words (e.g. "F equals m times a")
- Get straight to the point
- End with one short key fact
- Sound warm and encouraging"""
    user = f"""Question: {question}
Full answer: {full_answer[:800]}

Write a short spoken version (under 60 words):"""
    return call_hf(system, user, max_tokens=200)

@st.cache_resource(show_spinner=False)
def load_local_asr():
    """Load a lightweight local Whisper pipeline for transcription."""
    import torch
    from transformers import pipeline

    device = 0 if torch.cuda.is_available() else -1
    return pipeline(
        "automatic-speech-recognition",
        model="openai/whisper-tiny.en",
        device=device,
    )

def transcribe_audio_hf(audio_bytes):
    """Transcribe recorded audio using a local Whisper pipeline."""
    import io
    import numpy as np
    from scipy.io import wavfile

    try:
        sample_rate, audio_array = wavfile.read(io.BytesIO(audio_bytes))
        if audio_array.ndim > 1:
            audio_array = audio_array.mean(axis=1)
        if np.issubdtype(audio_array.dtype, np.integer):
            max_val = np.iinfo(audio_array.dtype).max
            audio_array = audio_array.astype(np.float32) / max(max_val, 1)
        else:
            audio_array = audio_array.astype(np.float32)

        asr = load_local_asr()
        result = asr({"raw": audio_array, "sampling_rate": int(sample_rate)})
        text = result.get("text") if isinstance(result, dict) else None
        if text and str(text).strip():
            return str(text).strip(), None
        return None, "The audio was processed, but no speech was detected."
    except Exception as e:
        return None, f"Local transcription failed: {str(e)[:160]}"

def show_voice_ui():
    """Render a native Streamlit voice panel using microphone recording."""
    st.markdown("Record a short question, transcribe it, then send it into the chat.")

    if st.session_state.get("voice_clear_pending"):
        st.session_state.voice_transcript = ""
        st.session_state.voice_transcript_editor = ""
        st.session_state.voice_clear_pending = False

    audio = st.audio_input("🎙️ Record your question", key="voice_audio_input")
    if audio is None:
        st.caption("Tap record, speak clearly, and stop when you're done.")
        return

    st.audio(audio)

    col1, col2 = st.columns(2)
    with col1:
        if st.button("📝 Transcribe", use_container_width=True, key="voice_transcribe_btn"):
            with st.spinner("🎧 Converting your speech to text..."):
                transcript, error = transcribe_audio_hf(audio.getvalue())
            if transcript:
                st.session_state.voice_transcript = transcript
                st.session_state.voice_transcript_editor = transcript
                st.success("✅ Speech detected and transcribed.")
            else:
                st.error(f"❌ Could not transcribe this recording. {error or ''}".strip())
    with col2:
        if st.button("🧹 Clear Voice", use_container_width=True, key="voice_clear_btn"):
            st.session_state.voice_transcript = ""
            st.session_state.voice_transcript_editor = ""
            st.rerun()

    if "voice_transcript_editor" not in st.session_state:
        st.session_state.voice_transcript_editor = st.session_state.get("voice_transcript", "")

    transcript_value = st.text_area(
        "Transcript",
        placeholder="Your recorded question will appear here...",
        key="voice_transcript_editor",
        height=100,
    )
    st.session_state.voice_transcript = transcript_value

    if st.button("➤ Ask PhysIQ", use_container_width=True, key="voice_send_btn"):
        if transcript_value.strip():
            st.session_state._voice_question = transcript_value.strip()
            st.session_state._is_voice_question = True
            st.session_state.voice_transcript = ""
            st.session_state.voice_clear_pending = True
            st.rerun()
        else:
            st.warning("Record something first, then transcribe it before sending.")




# ══════════════════════════════════════════════════════════════
#  IMAGE GENERATOR CONNECTOR
# ══════════════════════════════════════════════════════════════

IMAGE_TRIGGER_WORDS = [
    "generate an image", "create an image", "make an image",
    "draw", "generate a picture", "create a picture",
    "show me a picture", "show me an image", "visualise",
    "generate a photo", "create a photo", "make a photo",
    "render", "illustrate", "generate a diagram",
    "create a 3d", "make a 3d", "generate a 3d",
    "show me a 3d", "3d model of", "2d image of",
    "image of", "picture of", "photo of",
]

def is_image_request(message):
    msg = message.lower()
    return any(t in msg for t in IMAGE_TRIGGER_WORDS)

def build_image_prompt(user_request, context_answer=""):
    """Use AI to craft a perfect, detailed image generation prompt."""
    system = """You are a world-class AI image prompt engineer.
Given a user request, create the PERFECT image generation prompt for FLUX.1.
Your prompt must produce hyperrealistic, stunning, award-winning images.

Rules:
1. Start with the most important subject
2. Include: lighting, style, quality keywords, camera/perspective details
3. For 3D: add "hyperrealistic 3D render, octane render, physically based rendering, ray tracing, 8K, ultra detailed"
4. For scientific: add "scientific illustration, photorealistic, detailed, accurate, educational"
5. For diagrams: add "clean vector illustration, white background, professional, labeled"
6. ALWAYS end with: "sharp focus, high resolution, 8K UHD, masterpiece, best quality"
7. Keep prompt under 200 words
8. Return ONLY the prompt text, nothing else"""

    user = f"""User wants: {user_request}
Context (if any): {context_answer[:200] if context_answer else "none"}

Write the perfect image generation prompt:"""

    prompt = call_hf(system, user, max_tokens=250)
    return prompt.strip() if prompt else user_request

def generate_image_hf(prompt, width=1024, height=1024, model="flux"):
    """Generate image using Hugging Face text-to-image."""
    import base64
    import io
    from huggingface_hub import InferenceClient

    # Model options — FLUX.1-schnell is free and excellent
    models = {
        "flux":   "black-forest-labs/FLUX.1-schnell",
        "flux_dev": "black-forest-labs/FLUX.1-dev",
        "sdxl":   "stabilityai/stable-diffusion-xl-base-1.0",
        "sdxl_turbo": "stabilityai/sdxl-turbo",
    }
    model_order = [models.get(model, models["flux"]), models["sdxl"], models["sdxl_turbo"]]

    try:
        errors = []
        for model_id in model_order:
            try:
                client = InferenceClient(provider="hf-inference", api_key=HF_TOKEN, timeout=180)
                image = client.text_to_image(
                    prompt,
                    model=model_id,
                    width=width if "sdxl" not in model_id else min(width, 1024),
                    height=height if "sdxl" not in model_id else min(height, 1024),
                    num_inference_steps=4 if "FLUX.1-schnell" in model_id else 20,
                    guidance_scale=0.0 if "FLUX.1-schnell" in model_id else 7.5,
                )
                buffer = io.BytesIO()
                image.save(buffer, format="PNG")
                b64 = base64.b64encode(buffer.getvalue()).decode("utf-8")
                return b64, None
            except Exception as model_error:
                errors.append(f"{model_id}: {str(model_error)[:140]}")
        return None, "Generation failed: " + " | ".join(errors[:3])
    except Exception as e:
        return None, f"Error: {str(e)[:100]}"

def detect_image_style(message):
    """Detect if user wants 2D, 3D, diagram, or photo."""
    msg = message.lower()
    if any(x in msg for x in ["3d","three d","three-d","render","hyperrealistic","photorealistic"]):
        return "3d", 1024, 1024
    elif any(x in msg for x in ["diagram","chart","schematic","blueprint","2d","flat","vector"]):
        return "diagram", 1024, 768
    elif any(x in msg for x in ["landscape","panorama","wide"]):
        return "landscape", 1280, 768
    elif any(x in msg for x in ["portrait","person","human","face"]):
        return "portrait", 768, 1024
    else:
        return "general", 1024, 1024

def enhance_prompt_for_style(prompt, style):
    """Add style-specific quality keywords."""
    style_suffixes = {
        "3d": ", hyperrealistic 3D render, octane render, physically based rendering, ray tracing, subsurface scattering, ultra detailed, photorealistic, 8K UHD, masterpiece, best quality, sharp focus, professional lighting, studio quality, highly detailed",
        "diagram": ", clean scientific illustration, precise, labeled, white background, professional diagram, vector style, educational, high resolution, sharp, clear",
        "landscape": ", epic landscape photography, golden hour lighting, ultra wide angle, 8K UHD, masterpiece, photorealistic, award winning photography, sharp focus, vivid colors",
        "portrait": ", portrait photography, professional lighting, bokeh background, 8K, photorealistic, ultra detailed skin texture, cinematic, high quality",
        "general": ", ultra detailed, photorealistic, 8K UHD, masterpiece, best quality, sharp focus, vivid colors, professional, award winning",
    }
    suffix = style_suffixes.get(style, style_suffixes["general"])
    # Avoid duplicate quality words
    if "8K" not in prompt and "masterpiece" not in prompt:
        return prompt + suffix
    return prompt

def show_generated_image(b64_image, prompt, style, user_question):
    """Display the generated image beautifully in Streamlit."""
    import base64
    try:
        img_data = base64.b64decode(b64_image)
    except Exception:
        st.error("❌ The generated image could not be decoded.")
        return

    st.image(img_data, use_container_width=True)

    # Action buttons below image
    col1, col2, col3 = st.columns(3)
    with col1:
        # Download button
        st.download_button(
            "⬇️ Download",
            data=img_data,
            file_name=f"physiq_{style}_{int(__import__('time').time())}.png",
            mime="image/png",
            use_container_width=True
        )
    with col2:
        if st.button("🎬 Animate This", key=f"anim_img_{hash(prompt)}", use_container_width=True):
            # Store image for animator
            st.session_state.animator_bg_image = b64_image
            st.session_state.animator_bg_prompt = prompt
            # Generate animation with image as background
            with st.spinner("🎬 Building animation..."):
                anim = generate_animation_script(user_question,
                    f"Create animation using this image as reference: {prompt[:100]}")
                if anim:
                    anim["bg_image_b64"] = b64_image
                    st.session_state.animation_data = anim
                    st.session_state.show_animator = True
                    st.rerun()
    with col3:
        if st.button("🔄 Regenerate", key=f"regen_{hash(prompt)}", use_container_width=True):
            st.session_state._regen_prompt = prompt
            st.session_state._regen_style = style
            st.session_state._regen_question = user_question
            st.rerun()

    # Show prompt used
    with st.expander("🔍 Prompt used"):
        st.code(prompt, language=None)

def handle_image_request(question, vs):
    """Full pipeline: detect → craft prompt → generate → display."""
    # Detect style and dimensions
    style, width, height = detect_image_style(question)

    # Get some context from knowledge base
    context = ""
    if vs:
        try:
            results = vs.similarity_search_with_score(question, k=2)
            context = results[0][0].page_content[:300] if results else ""
        except: pass

    # Build optimised prompt
    with st.spinner("🎨 Crafting the perfect prompt..."):
        raw_prompt = build_image_prompt(question, context)

    # Add quality enhancers
    final_prompt = enhance_prompt_for_style(raw_prompt, style)

    # Generate
    style_labels = {
        "3d": "⚡ Generating hyperrealistic 3D image with FLUX.1...",
        "diagram": "📐 Generating scientific diagram...",
        "landscape": "🌄 Generating landscape...",
        "portrait": "👤 Generating portrait...",
        "general": "🖼️ Generating image",
    }
    spinner_msg = style_labels.get(style, "🖼️ Generating image...")

    with st.spinner(spinner_msg):
        b64, error = generate_image_hf(final_prompt, width, height)

    if error:
        st.error(f"❌ {error}")
        st.info("💡 Tip: FLUX.1 model sometimes needs a moment to warm up. Try again in 30 seconds.")
        return None, final_prompt

    return b64, final_prompt


# ══════════════════════════════════════════════════════════════
#  USER PROFILE & PERSONALISATION SYSTEM
# ══════════════════════════════════════════════════════════════

def save_user_profile(profile_data):
    """Save analysed profile to Supabase."""
    try:
        import json
        get_authed_client().table("user_profiles").upsert({
            "user_id": get_user_id(),
            "profile_json": json.dumps(profile_data),
            "updated_at": datetime.datetime.utcnow().isoformat()
        }, on_conflict="user_id").execute()
    except Exception as e:
        pass  # table may not exist yet, fail silently

def load_user_profile():
    """Load saved profile from Supabase."""
    try:
        import json
        res = get_authed_client().table("user_profiles").select("*").eq(
            "user_id", get_user_id()).execute()
        if res.data:
            return json.loads(res.data[0]["profile_json"])
    except:
        pass
    return None

def analyse_user_profile(messages, past_convos, quiz_mistakes):
    """Use AI to deeply analyse all user data and build a profile."""
    # Build a summary of all interactions
    recent_questions = [m["content"] for m in messages if m["role"]=="user"][-30:]
    past_questions = [p["content"] for p in past_convos if p.get("role")=="user"][-50:]
    all_questions = list(set(recent_questions + past_questions))

    mistake_topics = [m.get("topic","") for m in quiz_mistakes] if quiz_mistakes else []
    wrong_questions = [m.get("question","") for m in quiz_mistakes] if quiz_mistakes else []

    if len(all_questions) < 2:
        return None

    system = """You are an expert educational psychologist and learning analyst.
Analyse the student's conversation history, quiz mistakes, and question patterns.
Return ONLY a valid JSON object with this exact structure — no extra text:

{
  "weak_concepts": ["concept1", "concept2", "concept3"],
  "strong_concepts": ["concept1", "concept2", "concept3"],
  "favourite_topics": ["topic1", "topic2"],
  "learning_style": "Visual / Reading / Problem-solving / Conceptual",
  "preferred_detail_level": "Simple / Balanced / Detailed / Expert",
  "asking_pattern": "e.g. Asks short direct questions / Asks deep theoretical questions",
  "personality": "e.g. Curious and explorative / Goal-oriented / Creative",
  "how_to_reply": "e.g. Use analogies, keep it short, avoid jargon",
  "improvement_plan": ["tip1", "tip2", "tip3"],
  "strengths_summary": "One sentence about what they are best at",
  "weakness_summary": "One sentence about their main gap",
  "engagement_level": "High / Medium / Low",
  "topics_explored": ["topic1", "topic2", "topic3", "topic4", "topic5"],
  "recommended_next": ["topic1", "topic2", "topic3"]
}"""

    user = f"""Student's questions (most recent first):
{chr(10).join([f'- {q}' for q in all_questions[:40]])}

Quiz mistakes (topics where errors occurred):
{chr(10).join([f'- {t}: {q}' for t,q in zip(mistake_topics[:15], wrong_questions[:15])])}

Analyse this student deeply and return the JSON profile:"""

    import json
    raw = call_hf(system, user, max_tokens=800)
    if not raw:
        return None
    raw = raw.strip()
    start = raw.find("{")
    end = raw.rfind("}") + 1
    if start == -1:
        return None
    try:
        return json.loads(raw[start:end])
    except:
        return None

def build_personalisation_context(profile):
    """Build a short context string to inject into every AI response."""
    if not profile:
        return ""
    style = profile.get("how_to_reply", "")
    level = profile.get("preferred_detail_level", "Balanced")
    personality = profile.get("personality", "")
    weak = profile.get("weak_concepts", [])
    strong = profile.get("strong_concepts", [])
    context = f"""\n\n[STUDENT PROFILE — adjust your response accordingly:
- Detail level: {level}
- How they like replies: {style}
- Personality: {personality}
- Strong in: {", ".join(strong[:3]) if strong else "general"}
- Needs help with: {", ".join(weak[:3]) if weak else "varies"}
Always match your tone and depth to this profile.]"""
    return context

def show_profile_page():
    """Render the full user profile dashboard."""
    profile = st.session_state.get("user_profile")
    user_name = st.session_state.user.user_metadata.get("full_name", "Student")
    dark = st.session_state.dark_mode
    bg2  = "#161b22" if dark else "#ffffff"
    acc  = "#58a6ff"
    grn  = "#3fb950"
    red  = "#f85149"
    yel  = "#d29922"

    st.markdown(f"## 👤 {user_name}'s Learning Profile")
    st.caption("Automatically generated from your conversations and quiz history")

    col_refresh, col_back = st.columns([1,4])
    with col_refresh:
        if st.button("🔄 Refresh Profile", use_container_width=True):
            with st.spinner("🧠 Analysing your learning history..."):
                msgs = st.session_state.get("messages",[])
                past = load_past_conversations()
                mistakes = load_quiz_mistakes()
                new_profile = analyse_user_profile(msgs, past, mistakes)
                if new_profile:
                    st.session_state.user_profile = new_profile
                    save_user_profile(new_profile)
                    st.success("✅ Profile updated!")
                    st.rerun()
                else:
                    st.warning("Not enough conversation data yet. Ask more questions first!")
    with col_back:
        if st.button("← Back to Chat", use_container_width=True):
            st.session_state.show_profile = False
            st.rerun()

    if not profile:
        # Try loading from Supabase
        saved = load_user_profile()
        if saved:
            st.session_state.user_profile = saved
            profile = saved
        else:
            st.info("📊 No profile yet! Ask some questions in the chat first, then click 🔄 Refresh Profile")
            st.markdown("""
            **The profile will show:**
            - 🔴 Topics where you need improvement
            - 🟢 Topics where you are strong
            - ❤️ Your favourite subjects
            - 🎯 How the AI should talk to you
            - 📈 A personalised improvement plan
            """)
            return

    # ── Score cards ───────────────────────────────────────────
    st.markdown("---")
    c1, c2, c3, c4 = st.columns(4)
    metrics = [
        (c1, "🏆 Strengths", len(profile.get("strong_concepts",[])), grn),
        (c2, "📚 Topics Done", len(profile.get("topics_explored",[])), acc),
        (c3, "⚠️ Weak Areas", len(profile.get("weak_concepts",[])), red),
        (c4, "🎯 Engagement", profile.get("engagement_level","?"), yel),
    ]
    for col, label, val, color in metrics:
        with col:
            st.markdown(f"""
            <div style="background:{bg2};border:1px solid {color}30;border-left:4px solid {color};
            border-radius:10px;padding:14px 16px;text-align:center;margin-bottom:8px">
                <div style="color:{color};font-size:0.8rem;font-weight:600;margin-bottom:4px">{label}</div>
                <div style="color:#e6edf3;font-size:1.6rem;font-weight:700">{val}</div>
            </div>""", unsafe_allow_html=True)

    st.markdown("---")

    # ── Two column layout ─────────────────────────────────────
    left, right = st.columns(2)

    with left:
        # Strong concepts
        strong = profile.get("strong_concepts", [])
        if strong:
            st.markdown(f"### 🟢 You're Great At")
            for s in strong:
                st.markdown(f"""<div style="background:{bg2};border:1px solid {grn}40;border-radius:8px;
                padding:8px 14px;margin:4px 0;color:#e6edf3;font-size:13px">
                ✅ {s}</div>""", unsafe_allow_html=True)
        st.markdown("")

        # Favourite topics
        favs = profile.get("favourite_topics", [])
        if favs:
            st.markdown(f"### ❤️ Favourite Topics")
            for f2 in favs:
                st.markdown(f"""<div style="background:{bg2};border:1px solid {acc}40;border-radius:8px;
                padding:8px 14px;margin:4px 0;color:#e6edf3;font-size:13px">
                ❤️ {f2}</div>""", unsafe_allow_html=True)

    with right:
        # Weak concepts
        weak = profile.get("weak_concepts", [])
        if weak:
            st.markdown(f"### 🔴 Needs Improvement")
            for w in weak:
                st.markdown(f"""<div style="background:{bg2};border:1px solid {red}40;border-radius:8px;
                padding:8px 14px;margin:4px 0;color:#e6edf3;font-size:13px">
                ⚠️ {w}</div>""", unsafe_allow_html=True)
        st.markdown("")

        # Topics explored
        explored = profile.get("topics_explored", [])
        if explored:
            st.markdown(f"### 🗺️ Topics Explored")
            for t in explored[:6]:
                st.markdown(f"""<div style="background:{bg2};border:1px solid #30363d;border-radius:8px;
                padding:8px 14px;margin:4px 0;color:#8b949e;font-size:12px">
                📖 {t}</div>""", unsafe_allow_html=True)

    st.markdown("---")

    # ── Learning style card ───────────────────────────────────
    st.markdown("### 🧠 Your Learning DNA")
    dna1, dna2 = st.columns(2)
    with dna1:
        style_items = [
            ("Learning Style", profile.get("learning_style","?"), "🎓"),
            ("Detail Level", profile.get("preferred_detail_level","?"), "📏"),
            ("Asking Pattern", profile.get("asking_pattern","?"), "💬"),
            ("Personality", profile.get("personality","?"), "🌟"),
            ("Engagement", profile.get("engagement_level","?"), "⚡"),
        ]
        for label, val, icon in style_items:
            st.markdown(f"""<div style="background:{bg2};border:1px solid #30363d;border-radius:10px;
            padding:12px 16px;margin:6px 0;display:flex;align-items:center;gap:12px">
            <span style="font-size:1.3rem">{icon}</span>
            <div><div style="color:#8b949e;font-size:10px;font-weight:600;letter-spacing:1px">{label.upper()}</div>
            <div style="color:#e6edf3;font-size:13px;font-weight:500">{val}</div></div>
            </div>""", unsafe_allow_html=True)

    with dna2:
        # How AI talks to this user
        how = profile.get("how_to_reply","")
        if how:
            st.markdown(f"""<div style="background:linear-gradient(135deg,#1f6feb20,#388bfd10);
            border:1px solid {acc}40;border-radius:12px;padding:16px;margin-bottom:10px">
            <div style="color:{acc};font-weight:700;margin-bottom:8px;font-size:13px">🤖 How PhysIQ Talks to You</div>
            <div style="color:#e6edf3;font-size:13px;line-height:1.6">{how}</div>
            </div>""", unsafe_allow_html=True)

        # Summaries
        s_sum = profile.get("strengths_summary","")
        w_sum = profile.get("weakness_summary","")
        if s_sum:
            st.markdown(f"""<div style="background:{bg2};border:1px solid {grn}40;border-radius:10px;
            padding:12px 16px;margin:6px 0">
            <div style="color:{grn};font-size:11px;font-weight:700;margin-bottom:4px">YOUR STRENGTH</div>
            <div style="color:#e6edf3;font-size:12px">{s_sum}</div>
            </div>""", unsafe_allow_html=True)
        if w_sum:
            st.markdown(f"""<div style="background:{bg2};border:1px solid {red}40;border-radius:10px;
            padding:12px 16px;margin:6px 0">
            <div style="color:{red};font-size:11px;font-weight:700;margin-bottom:4px">YOUR MAIN GAP</div>
            <div style="color:#e6edf3;font-size:12px">{w_sum}</div>
            </div>""", unsafe_allow_html=True)

    st.markdown("---")

    # ── Improvement plan ─────────────────────────────────────
    plan = profile.get("improvement_plan", [])
    rec  = profile.get("recommended_next", [])
    if plan:
        st.markdown("### 📈 Your Personalised Improvement Plan")
        for i, tip in enumerate(plan, 1):
            st.markdown(f"""<div style="background:{bg2};border:1px solid #30363d;border-left:4px solid {acc};
            border-radius:10px;padding:12px 16px;margin:6px 0;color:#e6edf3;font-size:13px">
            <span style="color:{acc};font-weight:700">Step {i}:</span> {tip}
            </div>""", unsafe_allow_html=True)

    if rec:
        st.markdown("### 🎯 What to Study Next")
        cols_rec = st.columns(min(len(rec), 3))
        for i, topic in enumerate(rec[:3]):
            with cols_rec[i]:
                if st.button(f"📖 {topic}", use_container_width=True, key=f"rec_{i}"):
                    st.session_state.show_profile = False
                    st.session_state.messages.append({
                        "role":"user","content":f"Teach me about {topic}"
                    })
                    st.rerun()

    st.markdown("---")
    st.caption("Profile is rebuilt whenever you click 🔄 Refresh Profile. The more you use PhysIQ, the more accurate it gets!")



# ══════════════════════════════════════════════════════════════
#  WEB SEARCH CONNECTOR
#  Free APIs only: Wikipedia, arXiv, OpenLibrary, DuckDuckGo
#  All free for commercial use, no API keys required
# ══════════════════════════════════════════════════════════════

import re as _re

WEB_SEARCH_TRIGGERS = [
    "search for","look up","find information","what is the latest",
    "current","recent","news about","search the web","look online",
    "find out","google","web search","search online","what happened",
    "today","2024","2025","latest research","recent study","new paper",
    "who is","where is","when did","how many","what country",
]

def needs_web_search(question, confidence_class):
    """Decide if web search would help."""
    msg = question.lower()
    low_conf = confidence_class in ["conf-low","conf-med"]
    has_trigger = any(t in msg for t in WEB_SEARCH_TRIGGERS)
    is_factual = any(x in msg for x in ["who","what","when","where","how many","which","list of"])
    return (low_conf and is_factual) or has_trigger

def search_wikipedia(query, sentences=4):
    """Search Wikipedia — completely free, no key needed."""
    try:
        import urllib.request, urllib.parse, json
        q = urllib.parse.quote(query)
        url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{q}"
        req = urllib.request.Request(url, headers={"User-Agent":"PhysIQ/1.0 (educational)"})
        with urllib.request.urlopen(req, timeout=8) as r:
            data = json.loads(r.read())
        if data.get("type") == "disambiguation":
            return None
        extract = data.get("extract","")
        # Return first N sentences
        sents = extract.split(". ")
        return ". ".join(sents[:sentences]) + ("." if len(sents)>sentences else "")
    except:
        return None

def search_arxiv(query, max_results=3):
    """Search arXiv for scientific papers — completely free."""
    try:
        import urllib.request, urllib.parse
        q = urllib.parse.quote(query)
        url = f"http://export.arxiv.org/api/query?search_query=all:{q}&start=0&max_results={max_results}&sortBy=relevance"
        req = urllib.request.Request(url, headers={"User-Agent":"PhysIQ/1.0"})
        with urllib.request.urlopen(req, timeout=10) as r:
            xml = r.read().decode("utf-8")
        # Parse entries
        titles = _re.findall(r'<title>(.*?)</title>', xml)[1:]  # skip feed title
        summaries = _re.findall(r'<summary>(.*?)</summary>', xml, _re.DOTALL)
        results = []
        for i, (t, s) in enumerate(zip(titles[:max_results], summaries[:max_results])):
            t = t.strip().replace("\n"," ")
            s = s.strip().replace("\n"," ")[:200]
            results.append(f"📄 **{t}**: {s}...")
        return results if results else None
    except:
        return None

def search_duckduckgo(query):
    """DuckDuckGo Instant Answer API — free, no key, good for facts."""
    try:
        import urllib.request, urllib.parse, json
        q = urllib.parse.quote(query)
        url = f"https://api.duckduckgo.com/?q={q}&format=json&no_redirect=1&no_html=1&skip_disambig=1"
        req = urllib.request.Request(url, headers={"User-Agent":"PhysIQ/1.0"})
        with urllib.request.urlopen(req, timeout=8) as r:
            data = json.loads(r.read())
        abstract = data.get("AbstractText","").strip()
        answer = data.get("Answer","").strip()
        definition = data.get("Definition","").strip()
        related = [t.get("Text","") for t in data.get("RelatedTopics",[])[:2] if t.get("Text")]
        result = answer or abstract or definition
        if result:
            return result + (" | " + " | ".join(related) if related else "")
        return None
    except:
        return None

def search_openlibrary(query):
    """Open Library search — free, great for book/author queries."""
    try:
        import urllib.request, urllib.parse, json
        q = urllib.parse.quote(query)
        url = f"https://openlibrary.org/search.json?q={q}&limit=3&fields=title,author_name,first_publish_year,subject"
        req = urllib.request.Request(url, headers={"User-Agent":"PhysIQ/1.0"})
        with urllib.request.urlopen(req, timeout=8) as r:
            data = json.loads(r.read())
        docs = data.get("docs",[])[:3]
        results = []
        for d in docs:
            t = d.get("title","?")
            a = ", ".join(d.get("author_name",["Unknown"])[:2])
            y = d.get("first_publish_year","")
            results.append(f"📚 *{t}* by {a} ({y})")
        return results if results else None
    except:
        return None

def web_search_all(query):
    """Run all free searches and compile results."""
    results = {}
    # Wikipedia (most reliable for facts)
    wiki = search_wikipedia(query)
    if wiki: results["Wikipedia"] = wiki
    # DuckDuckGo instant answers
    ddg = search_duckduckgo(query)
    if ddg: results["DuckDuckGo"] = ddg
    # arXiv for science topics
    sci_words = ["physics","chemistry","quantum","atom","molecule","theorem",
                 "research","paper","study","formula","equation","theory"]
    if any(w in query.lower() for w in sci_words):
        arxiv = search_arxiv(query)
        if arxiv: results["arXiv"] = arxiv
    # Open Library for books/authors
    book_words = ["book","author","novel","written by","publication","published"]
    if any(w in query.lower() for w in book_words):
        books = search_openlibrary(query)
        if books: results["Open Library"] = books
    return results

def format_web_results(results):
    """Format search results into readable text for AI context."""
    if not results:
        return ""
    lines = ["\n[WEB SEARCH RESULTS — use these to enrich your answer:]"]
    for source, data in results.items():
        lines.append(f"\n📡 {source}:")
        if isinstance(data, list):
            for item in data:
                lines.append(f"  • {item}")
        else:
            lines.append(f"  {data}")
    lines.append("[End of web results]\n")
    return "\n".join(lines)

def answer_with_web_search(question, vs, history_text):
    """Ask question WITH web search context added."""
    # Run web searches
    web_results = web_search_all(question)
    web_context = format_web_results(web_results)
    # Also get RAG context
    rag_context = ""
    if vs:
        try:
            results = vs.similarity_search_with_score(question, k=3)
            rag_context = "\n\n".join([r[0].page_content for r in results])
        except: pass
    system = """You are an expert tutor with access to live web search results.
Use BOTH your knowledge AND the web search results to give the most accurate, up-to-date answer.
Always mention if information comes from a web source.
Be accurate, clear and educational."""
    user = f"""Context from knowledge base:
{rag_context}

{web_context}

Conversation:
{history_text}

Question: {question}

Answer (use web results where relevant):"""
    answer = call_hf(system, user, max_tokens=1500)
    sources = list(web_results.keys()) if web_results else []
    return answer, sources



# ══════════════════════════════════════════════════════════════
#  PLUGIN ENGINE
# ══════════════════════════════════════════════════════════════
# ══════════════════════════════════════════════════════════════
#  PLUGIN ENGINE — PhysIQ
#  AI auto-selects plugins. Each plugin is a self-contained HTML
#  component that runs in the browser.
# ══════════════════════════════════════════════════════════════

import json as _json

# ── Plugin registry ───────────────────────────────────────────
PLUGINS = {
    "calculator": {
        "name": "Calculator",
        "icon": "🧮",
        "description": "Scientific calculator with step-by-step working",
        "triggers": ["calculate","compute","what is","=","solve","evaluate","how much is",
                     "square root","sin","cos","tan","log","derivative","integral","+","-","×","÷"],
        "color": "#06d6a0",
    },
    "code_runner": {
        "name": "Code Runner",
        "icon": "▶️",
        "description": "Run Python/JS code in browser with live output",
        "triggers": ["run this","execute","output of","what does this code","test this code",
                     "debug","run the following","print result","does this work"],
        "color": "#58a6ff",
    },
    "document_analyser": {
        "name": "Document Analyser",
        "icon": "📋",
        "description": "Analyse tables, extract data, summarise documents",
        "triggers": ["analyse this","extract data","summarise","what does this say",
                     "find in document","table","spreadsheet","compare","statistics from"],
        "color": "#a78bfa",
    },
    "multi_agent": {
        "name": "Multi-Agent",
        "icon": "🤖",
        "description": "Spawn multiple AI agents to solve one problem from different angles",
        "triggers": ["compare","pros and cons","multiple perspectives","debate","argue both sides",
                     "what are all the ways","different approaches","solve it multiple ways",
                     "all possible","brainstorm"],
        "color": "#f471b5",
    },
    "decision_maker": {
        "name": "Decision Maker",
        "icon": "⚖️",
        "description": "Structured decision tree with weighted factors",
        "triggers": ["should i","which is better","help me decide","compare options","best choice",
                     "which one","recommend","advice on choosing","decision","trade-off"],
        "color": "#ffd166",
    },
    "tool_creator": {
        "name": "Tool Creator",
        "icon": "🔧",
        "description": "AI generates a custom tool/plugin for your specific need",
        "triggers": ["create a tool","make a tool","build a plugin","i need a widget",
                     "generate a calculator","make an app for","create a form","build me a"],
        "color": "#ff8c42",
    },
    "idea_evolution": {
        "name": "Idea Evolution",
        "icon": "🧠",
        "description": "Turns rough ideas into stronger versions through staged refinement",
        "triggers": ["improve this idea","evolve this idea","refine my idea","brainstorm this idea",
                     "make this concept better","develop this idea","expand this idea","iterate on this"],
        "color": "#22c55e",
    },
    "ai_mentor": {
        "name": "AI Mentor",
        "icon": "🌱",
        "description": "Tracks long-term growth and creates a coaching plan over weeks or months",
        "triggers": ["mentor me","long term growth","study plan","weekly plan","monthly plan",
                     "track my progress","guide me over time","improvement roadmap"],
        "color": "#38bdf8",
    },
    "video_reader": {
        "name": "Video Reader",
        "icon": "🎬",
        "description": "Extract and analyse content from YouTube videos",
        "triggers": ["youtube","video","watch","transcript","summarise this video",
                     "what is this video about","analyse the video","video link"],
        "color": "#ff5555",
    },
    "hypothesis_generator": {
        "name": "Hypothesis Generator",
        "icon": "🔬",
        "description": "Generate testable scientific hypotheses from any topic",
        "triggers": ["hypothesis","hypotheses","research question","scientific theory",
                     "generate hypothesis","null hypothesis","alternative hypothesis",
                     "testable prediction","experiment design","variables"],
        "color": "#06d6a0",
    },
    "context_rewriter": {
        "name": "Context Rewriter",
        "icon": "✏️",
        "description": "Rewrite your query 4 ways for better AI understanding",
        "triggers": ["rephrase","rewrite my question","better way to ask","improve my query",
                     "reword","clarify my question","rewrite this","context rewriter"],
        "color": "#a78bfa",
    },
    "multi_step_planner": {
        "name": "Multi-Step Planner",
        "icon": "📋",
        "description": "Break any goal into a clear actionable step-by-step plan",
        "triggers": ["plan for","how to achieve","step by step plan","roadmap",
                     "action plan","learning plan","study plan","project plan","goal"],
        "color": "#ffd166",
    },
    "idea_generator": {
        "name": "Idea Generator",
        "icon": "💡",
        "description": "Generate startup, project and research ideas",
        "triggers": ["startup idea","business idea","project idea","app idea","research idea",
                     "generate ideas","brainstorm ideas","innovation","entrepreneurship"],
        "color": "#f471b5",
    },
}

def detect_plugins(question: str) -> list:
    """AI-driven plugin detection — returns list of plugin keys to invoke."""
    q = question.lower()
    matched = []
    for key, plugin in PLUGINS.items():
        score = sum(1 for t in plugin["triggers"] if t in q)
        if score > 0:
            matched.append((key, score))
    matched.sort(key=lambda x: -x[1])
    return [k for k, _ in matched[:3]]  # max 3 plugins at once

def ai_select_plugins(question: str, vs=None) -> list:
    """Use AI to intelligently decide which plugins to invoke."""
    plugin_list = "\n".join([f"- {k}: {v['description']} (triggers: {', '.join(v['triggers'][:4])})"
                              for k, v in PLUGINS.items()])
    system = """You are a plugin router for an educational AI app.
Given a user question, decide which plugins (if any) to invoke.
Return a JSON array of plugin keys to use, in order of use.
Available plugins:
""" + plugin_list + """

Rules:
- Only suggest plugins that genuinely help answer the question
- Return empty array [] if no plugin needed
- Max 3 plugins
- For "calculate + explain": ["calculator"]
- For "compare options": ["decision_maker"]
- For "run this code": ["code_runner"]
- For "multiple angles": ["multi_agent"]
Return ONLY a JSON array like: ["calculator"] or ["multi_agent","calculator"] or []"""

    raw = call_hf(system, f"Question: {question}\nPlugins to invoke:", max_tokens=60)
    if not raw:
        return detect_plugins(question)
    try:
        raw = raw.strip()
        start = raw.find("[")
        end = raw.rfind("]") + 1
        if start >= 0:
            plugins = _json.loads(raw[start:end])
            return [p for p in plugins if p in PLUGINS]
    except:
        pass
    return detect_plugins(question)


# ═══════════════════════════════════════════════════════════════
# PLUGIN HTML COMPONENTS
# ═══════════════════════════════════════════════════════════════

def plugin_calculator(expression=""):
    """Scientific calculator with step-by-step working."""
    return f"""<!DOCTYPE html><html><head>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;font-family:'Segoe UI',sans-serif;padding:14px;color:#e6edf3}}
.calc-wrap{{background:#161b22;border:1px solid #06d6a0;border-radius:14px;padding:16px;max-width:520px}}
.title{{color:#06d6a0;font-weight:700;font-size:13px;margin-bottom:12px;letter-spacing:.5px}}
.display{{background:#0d1117;border:1px solid #30363d;border-radius:10px;padding:14px;margin-bottom:12px;min-height:60px}}
.expr{{color:#8b949e;font-size:12px;min-height:16px}}
.result{{color:#06d6a0;font-size:24px;font-weight:700;text-align:right;word-break:break-all}}
.steps{{color:#c9d1d9;font-size:11px;margin-top:8px;padding:8px;background:#0d1117;border-radius:6px;display:none}}
.steps.show{{display:block}}
.grid{{display:grid;grid-template-columns:repeat(5,1fr);gap:6px;margin-bottom:8px}}
.btn{{padding:10px 6px;border:1px solid #21262d;border-radius:8px;background:#21262d;
  color:#e6edf3;cursor:pointer;font-size:13px;font-weight:600;transition:all .15s;text-align:center}}
.btn:hover{{background:#30363d;border-color:#06d6a0}}
.btn.op{{color:#06d6a0}}
.btn.fn{{color:#58a6ff;font-size:11px}}
.btn.eq{{background:#06d6a0;color:#000;border-color:#06d6a0}}
.btn.eq:hover{{background:#05c48f}}
.btn.clear{{color:#f85149}}
.input-row{{display:flex;gap:8px;margin-bottom:10px}}
.expr-input{{flex:1;background:#0d1117;border:1px solid #30363d;border-radius:8px;
  padding:8px 12px;color:#e6edf3;font-size:13px;font-family:'JetBrains Mono',monospace}}
.calc-btn{{background:#06d6a0;color:#000;border:none;border-radius:8px;padding:8px 14px;cursor:pointer;font-weight:700}}
.step-toggle{{background:none;border:1px solid #30363d;color:#8b949e;border-radius:6px;
  padding:4px 10px;cursor:pointer;font-size:11px;margin-top:6px}}
</style></head><body>
<div class="calc-wrap">
<div class="title">🧮 Scientific Calculator</div>
<div class="display">
  <div class="expr" id="expr"></div>
  <div class="result" id="result">0</div>
  <div class="steps" id="steps"></div>
  <button class="step-toggle" id="stBtn" onclick="toggleSteps()" style="display:none">📋 Show Steps</button>
</div>
<div class="input-row">
  <input class="expr-input" id="exprIn" placeholder="Type expression e.g. sin(30)*2 + sqrt(144)" value="{expression}" onkeydown="if(event.key==='Enter')evalExpr()">
  <button class="calc-btn" onclick="evalExpr()">= Calculate</button>
</div>
<div class="grid">
  <button class="btn fn" onclick="ins('sin(')">sin</button><button class="btn fn" onclick="ins('cos(')">cos</button><button class="btn fn" onclick="ins('tan(')">tan</button><button class="btn fn" onclick="ins('log(')">log</button><button class="btn fn" onclick="ins('sqrt(')">sqrt</button>
  <button class="btn fn" onclick="ins('asin(')">asin</button><button class="btn fn" onclick="ins('acos(')">acos</button><button class="btn fn" onclick="ins('atan(')">atan</button><button class="btn fn" onclick="ins('ln(')">ln</button><button class="btn fn" onclick="ins('abs(')">abs</button>
  <button class="btn op" onclick="ins('π')">π</button><button class="btn op" onclick="ins('e')">e</button><button class="btn op" onclick="ins('(')">(</button><button class="btn op" onclick="ins(')')">)</button><button class="btn op" onclick="ins('**')">**</button>
  <button class="btn" onclick="ins('7')">7</button><button class="btn" onclick="ins('8')">8</button><button class="btn" onclick="ins('9')">9</button><button class="btn op" onclick="ins('÷')">÷</button><button class="btn op" onclick="ins('%')">%</button>
  <button class="btn" onclick="ins('4')">4</button><button class="btn" onclick="ins('5')">5</button><button class="btn" onclick="ins('6')">6</button><button class="btn op" onclick="ins('×')">×</button><button class="btn op" onclick="ins('^')">^</button>
  <button class="btn" onclick="ins('1')">1</button><button class="btn" onclick="ins('2')">2</button><button class="btn" onclick="ins('3')">3</button><button class="btn op" onclick="ins('+')">+</button><button class="btn fn" onclick="ins('√')">√</button>
  <button class="btn" onclick="ins('0')">0</button>
  <button class="btn" onclick="ins('.')">.</button>
  <button class="btn clear" onclick="clr()">C</button>
  <button class="btn op" onclick="ins('-')">−</button>
  <button class="btn eq" onclick="evalExpr()">=</button>
</div>
</div>
<script>
const inp=document.getElementById('exprIn');
const res=document.getElementById('result');
const exprD=document.getElementById('expr');
const stepsD=document.getElementById('steps');
const stBtn=document.getElementById('stBtn');
let lastSteps=[];

function ins(v){{
  const m={{'÷':'/','×':'*','π':'Math.PI','e':'Math.E','√':'sqrt(','ln':'log('}};
  inp.value+=m[v]||v;
  inp.focus();
}}

function clr(){{inp.value='';res.textContent='0';exprD.textContent='';stepsD.classList.remove('show');stBtn.style.display='none';lastSteps=[];}}

function evalExpr(){{
  let expr=inp.value.trim();
  if(!expr)return;
  exprD.textContent=expr;
  lastSteps=[];
  try{{
    // Replace friendly notation
    let e=expr
      .replace(/÷/g,'/')
      .replace(/×/g,'*')
      .replace(/√/g,'Math.sqrt')
      .replace(/π/g,'Math.PI')
      .replace(/\\bsin\\(/g,'Math.sin(Math.PI/180*')
      .replace(/\\bcos\\(/g,'Math.cos(Math.PI/180*')
      .replace(/\\btan\\(/g,'Math.tan(Math.PI/180*')
      .replace(/\\basin\\(/g,'(180/Math.PI)*Math.asin(')
      .replace(/\\bacos\\(/g,'(180/Math.PI)*Math.acos(')
      .replace(/\\batan\\(/g,'(180/Math.PI)*Math.atan(')
      .replace(/\\bsqrt\\(/g,'Math.sqrt(')
      .replace(/\\babs\\(/g,'Math.abs(')
      .replace(/\\blog\\(/g,'Math.log10(')
      .replace(/\\bln\\(/g,'Math.log(')
      .replace(/\\*\\*/g,'**')
      .replace(/\\^/g,'**');
    const result=eval(e);
    res.textContent=Number.isInteger(result)?result:parseFloat(result.toFixed(10));
    // Build steps
    lastSteps=[
      `📝 Expression: ${{expr}}`,
      `🔄 Parsed: ${{e.replace(/Math\\./g,'')}}`,
      `✅ Result: ${{result}}`,
    ];
    if(expr.includes('sin')||expr.includes('cos')||expr.includes('tan'))
      lastSteps.splice(1,0,'ℹ️ Trig functions use degrees');
    stBtn.style.display='block';
    stepsD.classList.remove('show');
    stBtn.textContent='📋 Show Steps';
    // Send result back to parent
    window.parent.postMessage({{type:'plugin_result',plugin:'calculator',result:String(result),expr:expr}},'*');
  }}catch(err){{
    res.textContent='Error: '+err.message;
  }}
}}

function toggleSteps(){{
  if(stepsD.classList.contains('show')){{stepsD.classList.remove('show');stBtn.textContent='📋 Show Steps';}}
  else{{
    stepsD.innerHTML=lastSteps.map(s=>`<div style="margin:3px 0">${{s}}</div>`).join('');
    stepsD.classList.add('show');stBtn.textContent='🔼 Hide Steps';
  }}
}}

// Auto-evaluate if expression passed
if(inp.value)setTimeout(evalExpr,400);
</script></body></html>"""


def plugin_code_runner(code="", language="python"):
    """In-browser code runner using Pyodide (Python) or native JS."""
    return f"""<!DOCTYPE html><html><head>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;font-family:'Segoe UI',sans-serif;padding:14px}}
.wrap{{background:#161b22;border:1px solid #58a6ff;border-radius:14px;padding:16px;max-width:600px}}
.title{{color:#58a6ff;font-weight:700;font-size:13px;margin-bottom:12px}}
.lang-tabs{{display:flex;gap:6px;margin-bottom:10px}}
.tab{{padding:5px 14px;border-radius:6px;border:1px solid #30363d;background:#21262d;
  color:#8b949e;cursor:pointer;font-size:12px}}
.tab.active{{background:#58a6ff;color:#000;border-color:#58a6ff}}
.editor{{width:100%;height:180px;background:#0d1117;border:1px solid #30363d;border-radius:8px;
  padding:10px;color:#e6edf3;font-family:'JetBrains Mono',monospace;font-size:12px;resize:vertical}}
.run-btn{{margin-top:8px;background:#238636;color:#fff;border:none;border-radius:8px;
  padding:8px 20px;cursor:pointer;font-size:13px;font-weight:600}}
.run-btn:hover{{background:#2ea043}}
.output{{margin-top:10px;background:#0d1117;border:1px solid #30363d;border-radius:8px;
  padding:12px;min-height:60px;font-family:'JetBrains Mono',monospace;font-size:12px;
  color:#06d6a0;white-space:pre-wrap;word-break:break-all}}
.err{{color:#f85149}}
.status{{font-size:11px;color:#8b949e;margin-top:6px}}
.load-bar{{height:3px;background:#58a6ff;border-radius:2px;animation:load 1.5s infinite;display:none}}
@keyframes load{{0%{{width:0%}}100%{{width:100%}}}}
</style></head><body>
<div class="wrap">
<div class="title">▶️ Code Runner</div>
<div class="lang-tabs">
  <div class="tab {'active' if language=='python' else ''}" onclick="setLang('python')">🐍 Python</div>
  <div class="tab {'active' if language=='javascript' else ''}" onclick="setLang('javascript')">📜 JavaScript</div>
  <div class="tab" onclick="setLang('html')">🌐 HTML Preview</div>
</div>
<textarea class="editor" id="code" placeholder="Write your code here...">{code}</textarea>
<div class="load-bar" id="lb"></div>
<button class="run-btn" id="runBtn" onclick="runCode()">▶ Run</button>
<div class="status" id="status"></div>
<div class="output" id="output">Output will appear here...</div>
</div>
<script>
let lang='{language}';
let pyodide=null, pyLoading=false, pyLoaded=false;

function setLang(l){{
  lang=l;
  document.querySelectorAll('.tab').forEach((t,i)=>{{
    t.classList.toggle('active',['python','javascript','html'][i]===l);
  }});
  if(l==='python'&&!pyLoaded&&!pyLoading)loadPyodide();
}}

async function loadPyodide(){{
  pyLoading=true;
  document.getElementById('status').textContent='⏳ Loading Python runtime...';
  document.getElementById('lb').style.display='block';
  try{{
    const s=document.createElement('script');
    s.src='https://cdn.jsdelivr.net/pyodide/v0.24.1/full/pyodide.js';
    document.head.appendChild(s);
    await new Promise(r=>s.onload=r);
    pyodide=await window.loadPyodide();
    pyLoaded=true;
    document.getElementById('status').textContent='✅ Python ready!';
    document.getElementById('lb').style.display='none';
  }}catch(e){{
    document.getElementById('status').textContent='❌ Python load failed. Use JS mode.';
    document.getElementById('lb').style.display='none';
  }}
  pyLoading=false;
}}

async function runCode(){{
  const code=document.getElementById('code').value;
  const out=document.getElementById('output');
  const btn=document.getElementById('runBtn');
  btn.textContent='⏳ Running...';btn.disabled=true;
  out.className='output';out.textContent='';

  try{{
    if(lang==='python'){{
      if(!pyLoaded){{await loadPyodide();}}
      let captured='';
      pyodide.globals.set('__captured__','');
      await pyodide.runPythonAsync(`
import sys, traceback
from io import StringIO
_buf=StringIO()
sys.stdout=_buf
sys.stderr=_buf
try:
    exec(code_to_run)
except Exception as _e:
    print("Error:", _e)
finally:
    sys.stdout=sys.__stdout__
    sys.stderr=sys.__stderr__
import js
js.window.__captured__=_buf.getvalue()
`, {{code_to_run: code}});
      const result=window.__captured__;
      out.textContent=result||'(No output)';
      window.parent.postMessage({{type:'plugin_result',plugin:'code_runner',result:result,lang:'python'}},'*');
    }} else if(lang==='javascript'){{
      let logs=[];
      const orig=console.log;
      console.log=(...a)=>{{logs.push(a.map(x=>typeof x==='object'?JSON.stringify(x,null,2):String(x)).join(' '));orig(...a);}};
      try{{ eval(code); }}catch(e){{logs.push('Error: '+e.message);out.classList.add('err');}}
      console.log=orig;
      out.textContent=logs.join('\\n')||'(No console output)';
      window.parent.postMessage({{type:'plugin_result',plugin:'code_runner',result:logs.join('\\n'),lang:'js'}},'*');
    }} else {{
      // HTML preview
      const iframe=document.createElement('iframe');
      iframe.style.cssText='width:100%;height:200px;border:1px solid #30363d;border-radius:8px;background:#fff';
      out.innerHTML='';out.appendChild(iframe);
      iframe.contentDocument.open();
      iframe.contentDocument.write(code);
      iframe.contentDocument.close();
    }}
  }}catch(e){{out.textContent='Error: '+e.message;out.classList.add('err');}}
  btn.textContent='▶ Run';btn.disabled=false;
}}

// Auto-load Python if needed
if(lang==='python')loadPyodide();
// Auto-run if code provided
if(document.getElementById('code').value.trim())setTimeout(runCode,1500);
</script></body></html>"""


def plugin_multi_agent(question=""):
    """Spawn 3 AI agents to solve from different angles."""
    return f"""<!DOCTYPE html><html><head>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;font-family:'Segoe UI',sans-serif;padding:14px}}
.wrap{{background:#161b22;border:1px solid #f471b5;border-radius:14px;padding:16px}}
.title{{color:#f471b5;font-weight:700;font-size:13px;margin-bottom:4px}}
.subtitle{{color:#8b949e;font-size:11px;margin-bottom:14px}}
.question-row{{display:flex;gap:8px;margin-bottom:14px}}
.q-input{{flex:1;background:#0d1117;border:1px solid #30363d;border-radius:8px;
  padding:8px 12px;color:#e6edf3;font-size:12px}}
.go-btn{{background:#f471b5;color:#000;border:none;border-radius:8px;padding:8px 16px;
  cursor:pointer;font-weight:700;font-size:12px}}
.agents-grid{{display:grid;grid-template-columns:repeat(3,1fr);gap:10px}}
.agent{{background:#0d1117;border:1px solid #30363d;border-radius:10px;padding:12px}}
.agent-name{{font-size:11px;font-weight:700;margin-bottom:8px;padding-bottom:6px;border-bottom:1px solid #21262d}}
.agent-body{{font-size:11px;color:#c9d1d9;line-height:1.6;min-height:80px;white-space:pre-wrap}}
.loading{{color:#8b949e;font-style:italic}}
.synthesis{{margin-top:12px;background:#0d1117;border:1px solid #f471b5;border-radius:10px;padding:12px}}
.syn-title{{color:#f471b5;font-size:12px;font-weight:700;margin-bottom:8px}}
.syn-body{{font-size:12px;color:#e6edf3;line-height:1.6}}
.status{{margin-top:8px;font-size:11px;color:#8b949e;text-align:center}}
</style></head><body>
<div class="wrap">
<div class="title">🤖 Multi-Agent Solver</div>
<div class="subtitle">3 specialised AI agents attack the problem simultaneously</div>
<div class="question-row">
  <input class="q-input" id="qIn" placeholder="Enter problem for agents to solve..." value="{question}">
  <button class="go-btn" onclick="runAgents()">🚀 Deploy Agents</button>
</div>
<div class="agents-grid">
  <div class="agent">
    <div class="agent-name" style="color:#4da6ff">🔬 Agent 1: Analytical</div>
    <div class="agent-body loading" id="a1">Waiting...</div>
  </div>
  <div class="agent">
    <div class="agent-name" style="color:#06d6a0">🎨 Agent 2: Creative</div>
    <div class="agent-body loading" id="a2">Waiting...</div>
  </div>
  <div class="agent">
    <div class="agent-name" style="color:#ffd166">📚 Agent 3: Practical</div>
    <div class="agent-body loading" id="a3">Waiting...</div>
  </div>
</div>
<div class="synthesis" id="syn" style="display:none">
  <div class="syn-title">⚡ Synthesised Answer</div>
  <div class="syn-body" id="synBody"></div>
</div>
<div class="status" id="status"></div>
</div>
<script>
async function callAI(sysPrompt, userPrompt){{
  // Use the parent window's HF token via postMessage
  return new Promise((resolve)=>{{
    const id='req_'+Date.now()+'_'+Math.random();
    window.parent.postMessage({{type:'plugin_ai_call',id,system:sysPrompt,user:userPrompt,max_tokens:300}},'*');
    function handler(e){{
      if(e.data&&e.data.type==='plugin_ai_response'&&e.data.id===id){{
        window.removeEventListener('message',handler);
        resolve(e.data.result||'No response');
      }}
    }}
    window.addEventListener('message',handler);
    setTimeout(()=>resolve('Timeout — please try again.'),25000);
  }});
}}

async function runAgents(){{
  const q=document.getElementById('qIn').value.trim();
  if(!q)return;
  const a1El=document.getElementById('a1');
  const a2El=document.getElementById('a2');
  const a3El=document.getElementById('a3');
  const syn=document.getElementById('syn');
  const synBody=document.getElementById('synBody');
  const status=document.getElementById('status');
  [a1El,a2El,a3El].forEach(el=>{{el.textContent='⏳ Thinking...';el.className='agent-body loading';}});
  syn.style.display='none';
  status.textContent='🤖 All 3 agents working...';
  const systems={{
    a1:'You are an analytical expert. Answer this question with precise logic, formulas, and step-by-step reasoning. Be thorough but concise.',
    a2:'You are a creative teacher. Explain this using vivid analogies, real-world examples, and an intuitive approach. Make it memorable.',
    a3:'You are a practical expert. Focus on how this applies in real life, what to actually do, common mistakes, and practical tips.',
  }};
  const [r1,r2,r3]=await Promise.all([
    callAI(systems.a1,q),
    callAI(systems.a2,q),
    callAI(systems.a3,q),
  ]);
  a1El.textContent=r1;a1El.className='agent-body';
  a2El.textContent=r2;a2El.className='agent-body';
  a3El.textContent=r3;a3El.className='agent-body';
  status.textContent='✅ Synthesising...';
  const synResult=await callAI(
    'You are a synthesis engine. Given 3 different expert answers to the same question, create the single best combined answer that takes the strongest points from each.',
    `Question: ${{q}}\\n\\nAnalytical: ${{r1}}\\n\\nCreative: ${{r2}}\\n\\nPractical: ${{r3}}\\n\\nSynthesise the best answer:`
  );
  synBody.textContent=synResult;
  syn.style.display='block';
  status.textContent='✅ All agents done!';
  window.parent.postMessage({{type:'plugin_result',plugin:'multi_agent',result:synResult,question:q}},'*');
}}
if(document.getElementById('qIn').value)setTimeout(runAgents,300);
</script></body></html>"""


def plugin_decision_maker(question=""):
    """Structured decision tree with weighted scoring."""
    return f"""<!DOCTYPE html><html><head>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;font-family:'Segoe UI',sans-serif;padding:14px}}
.wrap{{background:#161b22;border:1px solid #ffd166;border-radius:14px;padding:16px;max-width:600px}}
.title{{color:#ffd166;font-weight:700;font-size:13px;margin-bottom:4px}}
.subtitle{{color:#8b949e;font-size:11px;margin-bottom:12px}}
.row{{display:flex;gap:8px;margin-bottom:10px}}
.inp{{flex:1;background:#0d1117;border:1px solid #30363d;border-radius:8px;padding:8px 12px;color:#e6edf3;font-size:12px}}
.btn{{background:#ffd166;color:#000;border:none;border-radius:8px;padding:8px 16px;cursor:pointer;font-weight:700;font-size:12px}}
.options-area{{margin-bottom:10px}}
.option-row{{display:flex;align-items:center;gap:8px;margin-bottom:6px}}
.opt-input{{flex:1;background:#0d1117;border:1px solid #30363d;border-radius:6px;padding:6px 10px;color:#e6edf3;font-size:12px}}
.add-btn{{background:#21262d;color:#ffd166;border:1px solid #30363d;border-radius:6px;padding:6px 12px;cursor:pointer;font-size:12px}}
.remove-btn{{background:none;border:none;color:#f85149;cursor:pointer;font-size:14px}}
.criteria-title{{color:#8b949e;font-size:11px;font-weight:600;margin:10px 0 6px;letter-spacing:1px}}
.criteria-row{{display:flex;align-items:center;gap:8px;margin-bottom:6px}}
.crit-input{{flex:1;background:#0d1117;border:1px solid #30363d;border-radius:6px;padding:6px 10px;color:#e6edf3;font-size:12px}}
.weight-input{{width:60px;background:#0d1117;border:1px solid #30363d;border-radius:6px;padding:6px;color:#ffd166;font-size:12px;text-align:center}}
.result-area{{margin-top:12px;background:#0d1117;border:1px solid #30363d;border-radius:10px;padding:12px;display:none}}
.winner{{color:#ffd166;font-size:16px;font-weight:700;margin-bottom:8px}}
.score-bar{{margin-bottom:6px}}
.score-label{{display:flex;justify-content:space-between;font-size:11px;margin-bottom:2px}}
.bar-bg{{height:8px;background:#21262d;border-radius:4px;overflow:hidden}}
.bar-fill{{height:100%;background:linear-gradient(90deg,#ffd166,#ff8c42);border-radius:4px;transition:width .5s}}
.reasoning{{font-size:11px;color:#c9d1d9;margin-top:10px;line-height:1.6}}
</style></head><body>
<div class="wrap">
<div class="title">⚖️ Decision Maker</div>
<div class="subtitle">Multi-criteria weighted decision analysis</div>
<div class="row">
  <input class="inp" id="decisionQ" placeholder="What are you deciding? e.g. Which programming language to learn?" value="{question}">
</div>
<div class="options-area">
  <div class="criteria-title">OPTIONS</div>
  <div id="options">
    <div class="option-row"><input class="opt-input" placeholder="Option 1"><button class="remove-btn" onclick="this.parentElement.remove()">✕</button></div>
    <div class="option-row"><input class="opt-input" placeholder="Option 2"><button class="remove-btn" onclick="this.parentElement.remove()">✕</button></div>
  </div>
  <button class="add-btn" onclick="addOption()">+ Add Option</button>
</div>
<div class="criteria-title">CRITERIA & WEIGHTS (0-10)</div>
<div id="criteria">
  <div class="criteria-row"><input class="crit-input" placeholder="Criterion e.g. Cost" value="Cost"><input class="weight-input" type="number" value="8" min="1" max="10"><button class="remove-btn" onclick="this.parentElement.remove()">✕</button></div>
  <div class="criteria-row"><input class="crit-input" placeholder="Criterion e.g. Ease of Use" value="Ease of Use"><input class="weight-input" type="number" value="7" min="1" max="10"><button class="remove-btn" onclick="this.parentElement.remove()">✕</button></div>
  <div class="criteria-row"><input class="crit-input" placeholder="Criterion" value="Future Value"><input class="weight-input" type="number" value="9" min="1" max="10"><button class="remove-btn" onclick="this.parentElement.remove()">✕</button></div>
</div>
<button class="add-btn" style="margin-top:6px" onclick="addCriteria()">+ Add Criterion</button>
<div class="row" style="margin-top:12px">
  <button class="btn" style="width:100%" onclick="decide()">⚖️ Analyse Decision</button>
</div>
<div class="result-area" id="results"></div>
</div>
<script>
function addOption(){{
  const d=document.createElement('div');d.className='option-row';
  d.innerHTML=`<input class="opt-input" placeholder="New option"><button class="remove-btn" onclick="this.parentElement.remove()">✕</button>`;
  document.getElementById('options').appendChild(d);
}}
function addCriteria(){{
  const d=document.createElement('div');d.className='criteria-row';
  d.innerHTML=`<input class="crit-input" placeholder="Criterion"><input class="weight-input" type="number" value="5" min="1" max="10"><button class="remove-btn" onclick="this.parentElement.remove()">✕</button>`;
  document.getElementById('criteria').appendChild(d);
}}

async function decide(){{
  const q=document.getElementById('decisionQ').value;
  const opts=[...document.querySelectorAll('#options .opt-input')].map(i=>i.value).filter(v=>v.trim());
  const crits=[...document.querySelectorAll('#criteria .criteria-row')].map(r=>{{
    const inputs=r.querySelectorAll('input');
    return {{name:inputs[0].value,weight:parseFloat(inputs[1].value)||5}};
  }}).filter(c=>c.name.trim());
  if(opts.length<2||crits.length<1){{alert('Add at least 2 options and 1 criterion');return;}}

  const resultsEl=document.getElementById('results');
  resultsEl.style.display='block';
  resultsEl.innerHTML='<div style="color:#8b949e;font-size:12px">⚖️ Analysing...</div>';

  // Ask AI to score each option on each criterion
  const prompt=`Decision: ${{q}}
Options: ${{opts.join(', ')}}
Criteria with weights: ${{crits.map(c=>c.name+'(weight:'+c.weight+')').join(', ')}}

Score each option on each criterion from 1-10. Return ONLY a JSON object:
{{"scores":{{"option1":{{"criterion1":score,...}},...}},"reasoning":"brief overall analysis"}}`;

  window.parent.postMessage({{type:'plugin_ai_call',id:'decision',system:'You are an expert decision analyst. Return only valid JSON.',user:prompt,max_tokens:400}},'*');
}}

window.addEventListener('message',e=>{{
  if(e.data?.type==='plugin_ai_response'&&e.data?.id==='decision'){{
    try{{
      const raw=e.data.result;
      const start=raw.indexOf('{{');const end=raw.lastIndexOf('}}')+1;
      const data=JSON.parse(raw.slice(start,end));
      const opts=[...document.querySelectorAll('#options .opt-input')].map(i=>i.value).filter(v=>v.trim());
      const crits=[...document.querySelectorAll('#criteria .criteria-row')].map(r=>{{
        const inputs=r.querySelectorAll('input');
        return {{name:inputs[0].value,weight:parseFloat(inputs[1].value)||5}};
      }}).filter(c=>c.name.trim());
      // Calculate weighted scores
      const totals={{}};
      opts.forEach(opt=>{{
        let total=0,maxPossible=0;
        crits.forEach(crit=>{{
          const score=(data.scores?.[opt]?.[crit.name]||data.scores?.[opt.toLowerCase()]?.[crit.name.toLowerCase()]||5);
          total+=score*crit.weight;
          maxPossible+=10*crit.weight;
        }});
        totals[opt]={{raw:total,max:maxPossible,pct:Math.round(total/maxPossible*100)}};
      }});
      const sorted=Object.entries(totals).sort((a,b)=>b[1].pct-a[1].pct);
      const winner=sorted[0][0];
      const maxPct=sorted[0][1].pct;
      let html=`<div class="winner">🏆 Recommended: ${{winner}}</div>`;
      sorted.forEach(([name,scores])=>{{
        const pct=scores.pct;
        const col=name===winner?'#ffd166':'#58a6ff';
        html+=`<div class="score-bar"><div class="score-label"><span style="color:${{col}}">${{name}}</span><span style="color:${{col}}">${{pct}}%</span></div>
          <div class="bar-bg"><div class="bar-fill" style="width:${{pct}}%;background:${{col}}20;background:linear-gradient(90deg,${{col}},${{col}}88)"></div></div></div>`;
      }});
      if(data.reasoning)html+=`<div class="reasoning">📝 ${{data.reasoning}}</div>`;
      document.getElementById('results').innerHTML=html;
      window.parent.postMessage({{type:'plugin_result',plugin:'decision_maker',result:`Winner: ${{winner}} (${{maxPct}}%). ${{data.reasoning||''}}`,question:document.getElementById('decisionQ').value}},'*');
    }}catch(err){{document.getElementById('results').innerHTML='<div style="color:#f85149;font-size:12px">Parse error — try again</div>';}}
  }}
}});
</script></body></html>"""


def plugin_tool_creator(request=""):
    """AI generates a new plugin/tool on demand."""
    return f"""<!DOCTYPE html><html><head>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;font-family:'Segoe UI',sans-serif;padding:14px}}
.wrap{{background:#161b22;border:1px solid #ff8c42;border-radius:14px;padding:16px;max-width:600px}}
.title{{color:#ff8c42;font-weight:700;font-size:13px;margin-bottom:4px}}
.subtitle{{color:#8b949e;font-size:11px;margin-bottom:12px}}
.row{{display:flex;gap:8px;margin-bottom:10px}}
.inp{{flex:1;background:#0d1117;border:1px solid #30363d;border-radius:8px;padding:8px 12px;color:#e6edf3;font-size:12px}}
.btn{{background:#ff8c42;color:#000;border:none;border-radius:8px;padding:8px 16px;cursor:pointer;font-weight:700;font-size:12px}}
.preview{{margin-top:12px;border:1px solid #30363d;border-radius:10px;overflow:hidden;display:none}}
.preview-title{{background:#21262d;padding:8px 12px;font-size:11px;color:#8b949e;display:flex;justify-content:space-between;align-items:center}}
.download-btn{{background:#238636;color:#fff;border:none;border-radius:6px;padding:4px 10px;cursor:pointer;font-size:11px}}
.status{{font-size:11px;color:#ff8c42;margin-top:6px;text-align:center}}
.examples{{display:flex;flex-wrap:wrap;gap:6px;margin-bottom:12px}}
.ex-chip{{background:#21262d;border:1px solid #30363d;border-radius:12px;padding:4px 10px;
  font-size:11px;color:#8b949e;cursor:pointer}}
.ex-chip:hover{{border-color:#ff8c42;color:#ff8c42}}
</style></head><body>
<div class="wrap">
<div class="title">🔧 Tool Creator</div>
<div class="subtitle">Describe any tool and AI builds it instantly</div>
<div class="examples">
  <div class="ex-chip" onclick="setEx('A unit converter for physics')">⚗️ Unit converter</div>
  <div class="ex-chip" onclick="setEx('A periodic table lookup tool')">🔬 Periodic table</div>
  <div class="ex-chip" onclick="setEx('A chemical equation balancer')">🧪 Equation balancer</div>
  <div class="ex-chip" onclick="setEx('A physics formula finder')">📐 Formula finder</div>
  <div class="ex-chip" onclick="setEx('A colour mixing tool')">🎨 Colour mixer</div>
  <div class="ex-chip" onclick="setEx('A grade calculator')">📊 Grade calc</div>
</div>
<div class="row">
  <input class="inp" id="toolReq" placeholder="Describe the tool you want... e.g. 'A tool to calculate projectile motion with a graph'" value="{request}">
  <button class="btn" onclick="createTool()">🔧 Build It</button>
</div>
<div class="status" id="status"></div>
<div class="preview" id="preview">
  <div class="preview-title">
    <span id="toolName">Generated Tool</span>
    <button class="download-btn" onclick="downloadTool()">⬇️ Save Tool</button>
  </div>
  <iframe id="toolFrame" style="width:100%;height:350px;border:none;background:#fff"></iframe>
</div>
</div>
<script>
let generatedHTML='';
function setEx(v){{document.getElementById('toolReq').value=v;}}

async function createTool(){{
  const req=document.getElementById('toolReq').value.trim();
  if(!req)return;
  document.getElementById('status').textContent='🤖 AI is building your tool...';
  document.getElementById('preview').style.display='none';

  const sysPrompt=`You are an expert web developer. Build a complete, self-contained HTML tool based on the user's request.
Requirements:
- Single HTML file with embedded CSS and JavaScript
- Beautiful dark-themed UI (#0d1117 background, colored accents)
- Fully functional, no external dependencies except CDN fonts/math
- Educational and accurate
- Responsive design
- Return ONLY the complete HTML code, nothing else`;

  window.parent.postMessage({{type:'plugin_ai_call',id:'toolcreate',system:sysPrompt,user:`Build this tool: ${{req}}`,max_tokens:2000}},'*');
}}

window.addEventListener('message',e=>{{
  if(e.data?.type==='plugin_ai_response'&&e.data?.id==='toolcreate'){{
    const raw=e.data.result||'';
    const start=raw.indexOf('<!DOCTYPE')<0?raw.indexOf('<html'):raw.indexOf('<!DOCTYPE');
    const end=raw.lastIndexOf('</html>')+7;
    if(start>=0&&end>start){{
      generatedHTML=raw.slice(start,end);
    }}else{{
      generatedHTML=raw;
    }}
    const frame=document.getElementById('toolFrame');
    frame.srcdoc=generatedHTML;
    document.getElementById('preview').style.display='block';
    document.getElementById('toolName').textContent=document.getElementById('toolReq').value.slice(0,40);
    document.getElementById('status').textContent='✅ Tool created!';
    window.parent.postMessage({{type:'plugin_result',plugin:'tool_creator',result:'Tool generated: '+document.getElementById('toolReq').value}},'*');
  }}
}});

function downloadTool(){{
  if(!generatedHTML)return;
  const blob=new Blob([generatedHTML],{{type:'text/html'}});
  const a=document.createElement('a');
  a.href=URL.createObjectURL(blob);
  a.download=document.getElementById('toolReq').value.slice(0,30).replace(/[^a-z0-9]/gi,'_')+'.html';
  a.click();
}}
</script></body></html>"""

def build_fast_tool_template(request: str):
    """Return an instant local HTML tool for common requests."""
    req = (request or "").lower()

    if any(word in req for word in ["unit converter", "convert units", "physics units"]):
        return """<!DOCTYPE html><html><head><meta charset="utf-8"><style>
body{margin:0;background:#0d1117;color:#e6edf3;font-family:Segoe UI,sans-serif;padding:18px}
.wrap{max-width:520px;margin:0 auto;background:#161b22;border:1px solid #ff8c42;border-radius:14px;padding:18px}
h2{margin:0 0 12px;color:#ff8c42;font-size:18px}.row{display:grid;grid-template-columns:1fr 1fr;gap:10px;margin-bottom:10px}
input,select,button{width:100%;padding:10px;border-radius:8px;border:1px solid #30363d;background:#0d1117;color:#e6edf3}
button{background:#ff8c42;color:#111;font-weight:700;cursor:pointer}.out{margin-top:12px;padding:12px;background:#0d1117;border-radius:10px;border:1px solid #30363d}
</style></head><body><div class="wrap"><h2>Unit Converter</h2>
<div class="row"><input id="val" type="number" placeholder="Value"><select id="kind"><option value="length">Length</option><option value="mass">Mass</option><option value="temp">Temperature</option></select></div>
<div class="row"><select id="from"></select><select id="to"></select></div><button onclick="convert()">Convert</button><div class="out" id="out">Result will appear here.</div></div>
<script>
const units={length:['m','cm','km','inch','ft'],mass:['kg','g','lb'],temp:['C','F','K']};
const from=document.getElementById('from'),to=document.getElementById('to'),kind=document.getElementById('kind');
function fill(){const arr=units[kind.value];from.innerHTML=arr.map(x=>`<option>${x}</option>`).join('');to.innerHTML=from.innerHTML;} kind.onchange=fill; fill();
function convert(){let v=parseFloat(document.getElementById('val').value); if(Number.isNaN(v)) return;
 let f=from.value,t=to.value,r=v;
 if(kind.value==='length'){const m={m:1,cm:0.01,km:1000,inch:0.0254,ft:0.3048}; r=v*m[f]/m[t];}
 if(kind.value==='mass'){const kg={kg:1,g:0.001,lb:0.45359237}; r=v*kg[f]/kg[t];}
 if(kind.value==='temp'){if(f==='C') r=v; if(f==='F') r=(v-32)*5/9; if(f==='K') r=v-273.15; if(t==='F') r=r*9/5+32; else if(t==='K') r=r+273.15;}
 document.getElementById('out').textContent=`${v} ${f} = ${r.toFixed(4)} ${t}`;}
</script></body></html>"""

    if any(word in req for word in ["grade calculator", "marks calculator", "percentage calculator"]):
        return """<!DOCTYPE html><html><head><meta charset="utf-8"><style>
body{margin:0;background:#0d1117;color:#e6edf3;font-family:Segoe UI,sans-serif;padding:18px}
.wrap{max-width:520px;margin:0 auto;background:#161b22;border:1px solid #ff8c42;border-radius:14px;padding:18px}
h2{margin:0 0 12px;color:#ff8c42}.row{display:flex;gap:8px;margin-bottom:10px} input,button{padding:10px;border-radius:8px;border:1px solid #30363d;background:#0d1117;color:#e6edf3}
input{flex:1} button{background:#ff8c42;color:#111;font-weight:700;cursor:pointer}.item{font-size:14px;padding:6px 0;border-bottom:1px solid #21262d}.out{margin-top:12px;padding:12px;background:#0d1117;border-radius:10px}
</style></head><body><div class="wrap"><h2>Grade Calculator</h2><div class="row"><input id="score" type="number" placeholder="Score"><input id="total" type="number" placeholder="Total"></div>
<button onclick="addMark()">Add Subject</button><div id="list"></div><div class="out" id="out">Add marks to see percentage and grade.</div></div>
<script>
const items=[]; function render(){document.getElementById('list').innerHTML=items.map((x,i)=>`<div class="item">Subject ${i+1}: ${x.score}/${x.total}</div>`).join('');
const sum=items.reduce((a,b)=>a+b.score,0), total=items.reduce((a,b)=>a+b.total,0); if(!total) return;
const pct=(sum/total)*100; let grade='F'; if(pct>=90) grade='A+'; else if(pct>=80) grade='A'; else if(pct>=70) grade='B'; else if(pct>=60) grade='C'; else if(pct>=50) grade='D';
document.getElementById('out').textContent=`Percentage: ${pct.toFixed(2)}% | Grade: ${grade}`;}
function addMark(){const s=parseFloat(score.value), t=parseFloat(total.value); if(Number.isNaN(s)||Number.isNaN(t)||t<=0)return; items.push({score:s,total:t}); score.value=''; total.value=''; render();}
</script></body></html>"""

    if any(word in req for word in ["formula finder", "physics formula", "chemistry formula"]):
        return """<!DOCTYPE html><html><head><meta charset="utf-8"><style>
body{margin:0;background:#0d1117;color:#e6edf3;font-family:Segoe UI,sans-serif;padding:18px}
.wrap{max-width:620px;margin:0 auto;background:#161b22;border:1px solid #ff8c42;border-radius:14px;padding:18px}
h2{margin:0 0 12px;color:#ff8c42} input{width:100%;padding:10px;border-radius:8px;border:1px solid #30363d;background:#0d1117;color:#e6edf3;margin-bottom:10px}
.card{padding:12px;border:1px solid #30363d;border-radius:10px;background:#0d1117;margin-top:8px}
.name{color:#ff8c42;font-weight:700}.meta{color:#8b949e;font-size:12px;margin-top:4px}
</style></head><body><div class="wrap"><h2>Formula Finder</h2><input id="q" placeholder="Search velocity, force, ohm, molarity..." oninput="search()"><div id="results"></div></div>
<script>
const formulas=[
{name:'Force', eq:'F = m × a', tag:'motion newton dynamics'},
{name:'Velocity', eq:'v = s / t', tag:'speed motion'},
{name:'Density', eq:'density = mass / volume', tag:'matter'},
{name:'Ohm Law', eq:'V = I × R', tag:'electricity current voltage resistance'},
{name:'Pressure', eq:'P = F / A', tag:'fluids force area'},
{name:'Molarity', eq:'M = moles / litres', tag:'chemistry solution concentration'}
];
function search(){const q=document.getElementById('q').value.toLowerCase().trim(); const matches=formulas.filter(f=>!q||(`${f.name} ${f.eq} ${f.tag}`).toLowerCase().includes(q));
document.getElementById('results').innerHTML=matches.map(f=>`<div class="card"><div class="name">${f.name}</div><div>${f.eq}</div><div class="meta">${f.tag}</div></div>`).join('')||'<div class="card">No matching formula found.</div>';}
search();
</script></body></html>"""

    return None

def needs_research_for_tool_request(request: str) -> bool:
    """Decide whether a tool request should gather real topic knowledge first."""
    req = (request or "").lower()
    research_terms = [
        "finder", "lookup", "encyclopedia", "reference", "formula", "equation",
        "periodic", "theory", "law", "concept", "glossary", "cheat sheet",
        "database", "catalog", "facts", "learn", "study",
    ]
    return any(term in req for term in research_terms)

def build_tool_research_context(request: str) -> str:
    """Gather topic knowledge from Wikipedia and DuckDuckGo for tool generation."""
    queries = [request.strip()]
    lowered = request.lower()
    removals = [
        "build", "make", "create", "tool", "plugin", "app", "finder", "lookup",
        "for", "that", "can", "with", "a", "an", "the"
    ]
    simplified = " ".join(word for word in re.sub(r"[^a-zA-Z0-9 ]+", " ", lowered).split() if word not in removals)
    simplified = " ".join(simplified.split()[:6]).strip()
    if simplified and simplified not in queries:
        queries.append(simplified)

    lines = []
    seen = set()
    for query in queries[:2]:
        wiki = search_wikipedia(query, sentences=5)
        ddg = search_duckduckgo(query)
        if wiki and wiki not in seen:
            lines.append(f"Wikipedia for '{query}': {wiki}")
            seen.add(wiki)
        if ddg and ddg not in seen:
            lines.append(f"DuckDuckGo for '{query}': {ddg}")
            seen.add(ddg)
    return "\n\n".join(lines[:4])

def build_tool_creator_prompt(request: str, style: str, quality: str, must_have: str, research_context: str = "") -> tuple[str, str]:
    """Build a stronger prompt for high-quality plugin generation."""
    style_guides = {
        "Studio": "polished product UI with strong hierarchy, thoughtful spacing, rich cards, and confident typography",
        "Lab": "clean scientific interface with precise controls, high readability, and structured output panels",
        "Playful": "bright, energetic interface with delightful labels, strong contrast, and a more expressive layout",
    }
    quality_guides = {
        "Fast MVP": "prefer compact implementation, fewer controls, and rapid usefulness over depth",
        "Balanced": "balance polish, clarity, and implementation simplicity",
        "World-Class": "aim for premium interaction design, robust validation, polished empty states, and high perceived quality",
    }
    system = f"""You are an elite product designer and frontend engineer building a plugin for a premium AI marketplace.
Build a single-file HTML plugin that feels {style_guides.get(style, style_guides['Lab'])}.
Quality target: {quality_guides.get(quality, quality_guides['Balanced'])}.

Requirements:
- Return ONLY complete HTML
- Inline CSS and JavaScript only
- No external dependencies
- Fully functional and responsive
- Dark UI using #0d1117, #161b22, and one strong accent color
- Clear title, subtitle, input area, primary action, and output/result area
- Strong UX: validation, empty states, helpful placeholder text, and polished result presentation
- Include at least one premium touch: tabbed sections, rich result cards, smart presets, sticky action bar, or guided empty state
- Make mobile layout feel intentional, not just shrunken desktop
- Use semantic labels and clear affordances
- Avoid broken states and partial interactions
- Avoid toy layouts and generic boilerplate
- Keep code self-contained and production-like
- If research context is provided, use it to make the plugin's content, labels, presets, examples, and built-in data more accurate
- Make sure the final tool includes everything the user explicitly asked for
"""
    user = f"""Build a plugin for this request:
{request}

Must-have features:
{must_have or "Choose the best features yourself based on the request."}

Research context:
{research_context or "No external research context was provided."}

Return the final HTML only."""
    return system, user

def refine_generated_tool(raw_html: str, request: str, style: str, quality: str, research_context: str = "") -> str:
    """Run a second-pass quality review on generated HTML."""
    if not raw_html:
        return raw_html
    system = """You are a senior frontend reviewer.
Improve the provided HTML tool while keeping it single-file and self-contained.
Fix weak UX, broken structure, missing labels, poor spacing, incomplete interactions, bland styling, and weak result presentation.
Upgrade it toward a premium plugin-marketplace standard.
If research context exists, use it to improve correctness and built-in content.
Return ONLY the improved complete HTML."""
    user = f"""Original request: {request}
Style direction: {style}
Quality target: {quality}

Research context:
{research_context or "None"}

Current HTML:
{raw_html}

Improved HTML:"""
    improved = call_hf(system, user, max_tokens=1800)
    return improved or raw_html

def suggest_tool_triggers(request: str) -> list[str]:
    """Create simple trigger phrases so custom plugins can be auto-selected later."""
    cleaned = re.sub(r"[^a-zA-Z0-9 ]+", " ", (request or "").lower())
    words = [w for w in cleaned.split() if len(w) > 2]
    base = " ".join(words[:4]).strip()
    triggers = []
    if base:
        triggers.extend([base, f"build {base}", f"make {base}"])
    if len(words) >= 2:
        triggers.append(" ".join(words[:2]))
    seen = []
    for item in triggers:
        if item and item not in seen:
            seen.append(item)
    return seen[:6]

def generate_tool_metadata(request: str, html: str) -> dict:
    """Generate plugin metadata so custom plugins feel first-class in the shop."""
    fallback = {
        "name": infer_tool_name(request),
        "description": f"Custom tool built from: {request[:70]}",
        "category": "Custom",
        "triggers": suggest_tool_triggers(request),
        "quality_note": "Built with PhysIQ Tool Creator",
    }
    raw = call_hf(
        "You are a plugin marketplace editor. Return only valid JSON with keys name, description, category, triggers, quality_note.",
        f"User request: {request}\n\nHTML preview:\n{html[:2200]}\n\nReturn compact metadata JSON:",
        max_tokens=220,
    )
    if not raw:
        return fallback
    try:
        start = raw.find("{")
        end = raw.rfind("}") + 1
        data = _json.loads(raw[start:end])
        return {
            "name": data.get("name") or fallback["name"],
            "description": data.get("description") or fallback["description"],
            "category": data.get("category") or fallback["category"],
            "triggers": [t for t in data.get("triggers", fallback["triggers"]) if isinstance(t, str)][:8] or fallback["triggers"],
            "quality_note": data.get("quality_note") or fallback["quality_note"],
        }
    except Exception:
        return fallback

def generate_tool_fast(request: str, style: str = "Lab", quality: str = "Balanced", must_have: str = ""):
    """Generate a stronger tool with templates first, then a refined AI pass."""
    research_context = build_tool_research_context(request) if needs_research_for_tool_request(request) else ""
    use_template = not must_have.strip() and not research_context
    local_html = build_fast_tool_template(request) if use_template else None
    if local_html:
        return local_html, "instant"

    system, user = build_tool_creator_prompt(request, style, quality, must_have, research_context)
    html = call_hf(system, user, max_tokens=1800)
    if html:
        html = refine_generated_tool(html, request, style, quality, research_context)
    return html, "ai"

def infer_tool_name(request: str) -> str:
    """Create a short human-friendly plugin name from the request."""
    cleaned = re.sub(r"[^a-zA-Z0-9 ]+", " ", (request or "").strip())
    cleaned = " ".join(cleaned.split())
    if not cleaned:
        return "Custom Tool"
    return " ".join(word.capitalize() for word in cleaned.split()[:4])

def normalize_tool_html(raw_html: str, request: str) -> str:
    """Extract a usable HTML document or wrap partial output in a safe shell."""
    raw = (raw_html or "").strip()
    if not raw:
        return ""

    start = raw.find("<!DOCTYPE")
    if start < 0:
        start = raw.find("<html")
    if start >= 0:
        raw = raw[start:]
        end = raw.lower().rfind("</html>")
        if end >= 0:
            raw = raw[: end + 7]
        return raw

    title = infer_tool_name(request)
    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><title>{title}</title><style>
body{{margin:0;background:#0d1117;color:#e6edf3;font-family:'Segoe UI',sans-serif;padding:18px}}
.wrap{{max-width:760px;margin:0 auto;background:#161b22;border:1px solid #ff8c42;border-radius:14px;padding:18px}}
h2{{margin:0 0 12px;color:#ff8c42}}
</style></head><body><div class="wrap"><h2>{title}</h2>{raw}</div></body></html>"""

def build_generic_tool_shell(request: str) -> str:
    """Guaranteed fallback plugin when AI output is missing or malformed."""
    title = infer_tool_name(request)
    safe_request = (request or "Custom tool").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><title>{title}</title><style>
body{{margin:0;background:#0d1117;color:#e6edf3;font-family:'Segoe UI',sans-serif;padding:18px}}
.wrap{{max-width:760px;margin:0 auto;background:#161b22;border:1px solid #ff8c42;border-radius:14px;padding:18px}}
h2{{margin:0 0 10px;color:#ff8c42}} .sub{{color:#8b949e;font-size:13px;margin-bottom:16px}}
.row{{display:flex;gap:8px;margin-bottom:10px}} input,textarea,button{{width:100%;padding:10px;border-radius:8px;border:1px solid #30363d;background:#0d1117;color:#e6edf3}}
button{{background:#ff8c42;color:#111;font-weight:700;cursor:pointer}} .card{{padding:12px;border:1px solid #30363d;border-radius:10px;background:#0d1117;margin-top:10px}}
</style></head><body><div class="wrap">
<h2>{title}</h2>
<div class="sub">Guaranteed fallback plugin for: {safe_request}</div>
<div class="row"><input id="mainInput" placeholder="Type your input here"></div>
<div class="row"><textarea id="notes" rows="6" placeholder="Add details, values, or steps here"></textarea></div>
<button onclick="runTool()">Run Tool</button>
<div class="card" id="out">This plugin is ready. Enter values and press Run Tool.</div>
</div>
<script>
function runTool(){{
  const a=document.getElementById('mainInput').value.trim();
  const b=document.getElementById('notes').value.trim();
  document.getElementById('out').innerHTML=`<strong>Input:</strong> ${{a||'(empty)'}}<br><br><strong>Details:</strong><br>${{(b||'No extra details provided.').replace(/\\n/g,'<br>')}}`;
}}
</script></body></html>"""

def register_custom_tool_plugin(request: str, html: str, mode: str) -> str:
    """Save a generated tool as a reusable custom plugin."""
    custom_plugins = dict(st.session_state.get("custom_plugins", {}))
    plugin_id = f"custom_tool_{len(custom_plugins) + 1}"
    metadata = generate_tool_metadata(request, html)
    custom_plugins[plugin_id] = {
        "name": metadata["name"],
        "icon": "🧩",
        "description": metadata["description"],
        "category": metadata["category"],
        "color": "#ff8c42",
        "html": html,
        "request": request,
        "mode": mode,
        "triggers": metadata["triggers"],
        "quality_note": metadata["quality_note"],
    }
    st.session_state.custom_plugins = custom_plugins
    return plugin_id

def render_tool_creator_plugin(context: str = "") -> None:
    """Native fast tool creator with local templates and lighter AI fallback."""
    pending_request = st.session_state.pop("tool_creator_request_pending", None)
    if pending_request is not None:
        st.session_state.tool_creator_request = pending_request

    if "tool_creator_request" not in st.session_state:
        st.session_state.tool_creator_request = context or ""
    if context and not st.session_state.tool_creator_request.strip():
        st.session_state.tool_creator_request = context

    st.caption("Create marketplace-grade plugins with a stronger design brief, generated metadata, and a refinement pass.")
    request = st.text_input(
        "Describe the tool you want",
        key="tool_creator_request",
        placeholder="A unit converter for physics",
        label_visibility="collapsed",
    )

    cfg1, cfg2 = st.columns(2)
    with cfg1:
        style = st.selectbox(
            "Visual style",
            ["Lab", "Studio", "Playful"],
            key="tool_creator_style",
        )
    with cfg2:
        quality = st.selectbox(
            "Build quality",
            ["World-Class", "Balanced", "Fast MVP"],
            key="tool_creator_quality",
        )

    must_have = st.text_area(
        "Must-have features",
        key="tool_creator_features",
        height=90,
        placeholder="Examples: export CSV, live preview, mobile-friendly layout, error validation, reset button",
    )

    ex1, ex2, ex3 = st.columns(3)
    with ex1:
        if st.button("⚗️ Unit Converter", key="tool_ex_unit", use_container_width=True):
            st.session_state.tool_creator_request_pending = "A unit converter for physics"
            st.rerun()
    with ex2:
        if st.button("📊 Grade Calc", key="tool_ex_grade", use_container_width=True):
            st.session_state.tool_creator_request_pending = "A grade calculator"
            st.rerun()
    with ex3:
        if st.button("📐 Formula Finder", key="tool_ex_formula", use_container_width=True):
            st.session_state.tool_creator_request_pending = "A physics formula finder"
            st.rerun()

    if st.button("🔧 Build Fast Tool", key="tool_creator_build_btn", use_container_width=True):
        if request.strip():
            with st.spinner("🔧 Building your tool..."):
                html, mode = generate_tool_fast(
                    request.strip(),
                    style=style,
                    quality=quality,
                    must_have=must_have.strip(),
                )
                html = normalize_tool_html(html, request.strip()) if html else ""
                if not html:
                    html = build_generic_tool_shell(request.strip())
                    mode = "fallback"
            if html:
                st.session_state.tool_creator_html = html
                st.session_state.tool_creator_mode = mode
                st.session_state.tool_creator_name = infer_tool_name(request.strip())
                st.session_state.tool_creator_last_plugin = register_custom_tool_plugin(request.strip(), html, mode)
            else:
                st.error("Could not generate the tool. Please try a shorter description.")
            st.rerun()
        else:
            st.warning("Describe the tool first.")

    html = st.session_state.get("tool_creator_html")
    if not html:
        return

    mode = st.session_state.get("tool_creator_mode", "ai")
    if mode == "instant":
        st.caption("Built instantly from a local template and added as a reusable custom plugin.")
    elif mode == "fallback":
        st.caption("Built from the guaranteed fallback shell and added as a reusable custom plugin.")
    else:
        st.caption("Built with the fast AI fallback and added as a reusable custom plugin.")
    st.components.v1.html(html, height=420, scrolling=True)
    if st.session_state.get("tool_creator_last_plugin"):
        plugin_id = st.session_state.get("tool_creator_last_plugin")
        plugin_meta = st.session_state.get("custom_plugins", {}).get(plugin_id, {})
        badges = []
        if plugin_meta.get("category"):
            badges.append(plugin_meta["category"])
        if plugin_meta.get("quality_note"):
            badges.append(plugin_meta["quality_note"])
        if badges:
            st.caption(" • ".join(badges))
        st.success(f"{st.session_state.get('tool_creator_name', 'Custom Tool')} is now available in the plugin shop.")
    st.download_button(
        "⬇️ Download Tool",
        data=html,
        file_name="physiq_tool.html",
        mime="text/html",
        key="download_fast_tool_btn",
        use_container_width=True,
    )

def run_multi_agent_solver(question: str) -> dict:
    """Solve one problem from multiple AI perspectives and synthesise the result."""
    prompts = {
        "analytical": (
            "🔬 Analytical",
            "You are an analytical expert. Answer with precise logic, formulas, and clear step-by-step reasoning. Be thorough but concise.",
        ),
        "creative": (
            "🎨 Creative",
            "You are a creative teacher. Explain using vivid analogies, intuition, and real-world examples. Make it memorable and easy to grasp.",
        ),
        "practical": (
            "📚 Practical",
            "You are a practical expert. Focus on real-life use, common mistakes, shortcuts, and what a student should actually remember.",
        ),
    }

    results = {}
    for key, (label, system_prompt) in prompts.items():
        answer = call_hf(system_prompt, f"Question: {question}\nAnswer:", max_tokens=500)
        results[key] = {
            "label": label,
            "answer": answer or "No response returned for this agent.",
        }

    synthesis_prompt = (
        "You are a synthesis engine. Combine the strongest points from the analytical, creative, and practical answers into one final answer. "
        "Keep it clear, useful, and well-structured."
    )
    synthesis_user = (
        f"Question: {question}\n\n"
        f"Analytical answer:\n{results['analytical']['answer']}\n\n"
        f"Creative answer:\n{results['creative']['answer']}\n\n"
        f"Practical answer:\n{results['practical']['answer']}\n\n"
        "Write the best combined final answer:"
    )
    synthesis = call_hf(synthesis_prompt, synthesis_user, max_tokens=800)

    return {
        "question": question,
        "agents": results,
        "synthesis": synthesis or "No synthesis was returned.",
    }

def render_multi_agent_plugin(context: str = "") -> None:
    """Native Streamlit renderer for the multi-agent solver."""
    if "multi_agent_input" not in st.session_state:
        st.session_state.multi_agent_input = context or ""

    if context and not st.session_state.multi_agent_input.strip():
        st.session_state.multi_agent_input = context

    st.caption("Three specialised AI agents solve the same problem from different angles, then PhysIQ combines their best ideas.")
    question = st.text_area(
        "Problem",
        key="multi_agent_input",
        height=110,
        placeholder="Enter a physics, chemistry, or reasoning problem...",
        label_visibility="collapsed",
    )

    if st.button("🚀 Deploy Agents", key="multi_agent_run_btn", use_container_width=True):
        if question.strip():
            with st.spinner("🤖 Analytical, creative, and practical agents are working..."):
                st.session_state.multi_agent_result = run_multi_agent_solver(question.strip())
            st.rerun()
        else:
            st.warning("Enter a problem first.")

    result = st.session_state.get("multi_agent_result")
    if not result:
        return

    st.markdown("**🤖 Agent Views**")
    cols = st.columns(3)
    for col, key in zip(cols, ["analytical", "creative", "practical"]):
        with col:
            agent = result["agents"][key]
            st.markdown(f"**{agent['label']}**")
            st.write(agent["answer"])

    st.markdown("**⚡ Combined Answer**")
    st.info(result["synthesis"])

def run_idea_evolution(idea: str) -> dict:
    """Evolve a rough idea through multiple improvement stages."""
    stages = {
        "core": (
            "🌱 Core Idea",
            "You are an innovation coach. Extract the core idea, the real goal, and the strongest seed inside the user's idea. Be concise and constructive.",
        ),
        "upgrade": (
            "🚀 Stronger Version",
            "You are a product and strategy thinker. Improve the idea so it becomes more useful, original, and practical while preserving the user's intent.",
        ),
        "execution": (
            "🛠️ Execution Plan",
            "You are an execution mentor. Turn the improved idea into concrete next steps, likely obstacles, and fast first actions.",
        ),
    }

    results = {}
    for key, (label, system_prompt) in stages.items():
        answer = call_hf(system_prompt, f"Idea: {idea}\nResponse:", max_tokens=450)
        results[key] = {"label": label, "answer": answer or "No response returned."}

    final_answer = call_hf(
        "You are an idea synthesiser. Combine the core idea, improved version, and execution guidance into one clear and motivating final evolution summary.",
        f"Idea: {idea}\n\nCore:\n{results['core']['answer']}\n\nUpgrade:\n{results['upgrade']['answer']}\n\nExecution:\n{results['execution']['answer']}\n\nFinal evolved version:",
        max_tokens=700,
    )
    return {"idea": idea, "stages": results, "final": final_answer or "No final summary was returned."}

def render_idea_evolution_plugin(context: str = "") -> None:
    """Native renderer for idea refinement and staged evolution."""
    if "idea_evolution_input" not in st.session_state:
        st.session_state.idea_evolution_input = context or ""
    if context and not st.session_state.idea_evolution_input.strip():
        st.session_state.idea_evolution_input = context

    st.caption("Start with a rough idea and PhysIQ will strengthen it, then turn it into an actionable version.")
    idea = st.text_area(
        "Idea",
        key="idea_evolution_input",
        height=110,
        placeholder="Example: I want to build an app that helps students understand physics visually.",
        label_visibility="collapsed",
    )

    if st.button("🧠 Evolve Idea", key="idea_evolution_run_btn", use_container_width=True):
        if idea.strip():
            with st.spinner("🧠 Strengthening and refining your idea..."):
                st.session_state.idea_evolution_result = run_idea_evolution(idea.strip())
            st.rerun()
        else:
            st.warning("Enter an idea first.")

    result = st.session_state.get("idea_evolution_result")
    if not result:
        return

    cols = st.columns(3)
    for col, key in zip(cols, ["core", "upgrade", "execution"]):
        with col:
            stage = result["stages"][key]
            st.markdown(f"**{stage['label']}**")
            st.write(stage["answer"])

    st.markdown("**✨ Evolved Idea**")
    st.success(result["final"])

def run_ai_mentor(topic: str, horizon: str, current_level: str, goal: str) -> dict:
    """Create a long-term mentoring plan with milestones and habits."""
    summary = call_hf(
        "You are a thoughtful AI mentor. Summarise the learner's current state, strengths, and biggest gap from their inputs.",
        f"Topic: {topic}\nCurrent level: {current_level}\nGoal: {goal}\nMentor summary:",
        max_tokens=260,
    )
    roadmap = call_hf(
        "You are a learning strategist. Create a structured growth roadmap with phases, milestones, and weekly focus areas.",
        f"Topic: {topic}\nTime horizon: {horizon}\nCurrent level: {current_level}\nGoal: {goal}\nRoadmap:",
        max_tokens=650,
    )
    habits = call_hf(
        "You are an accountability coach. Give weekly habits, check-ins, and motivation tips to sustain progress over time.",
        f"Topic: {topic}\nTime horizon: {horizon}\nCurrent level: {current_level}\nGoal: {goal}\nCoaching plan:",
        max_tokens=420,
    )
    return {
        "summary": summary or "No mentor summary returned.",
        "roadmap": roadmap or "No roadmap returned.",
        "habits": habits or "No coaching plan returned.",
    }

def render_ai_mentor_plugin(context: str = "") -> None:
    """Native renderer for long-term mentoring and growth planning."""
    if "mentor_topic" not in st.session_state:
        st.session_state.mentor_topic = context or ""
    if context and not st.session_state.mentor_topic.strip():
        st.session_state.mentor_topic = context
    if "mentor_goal" not in st.session_state:
        st.session_state.mentor_goal = ""
    if "mentor_level" not in st.session_state:
        st.session_state.mentor_level = "Beginner"
    if "mentor_horizon" not in st.session_state:
        st.session_state.mentor_horizon = "8 weeks"

    st.caption("Build a mentoring plan that tracks growth over weeks or months instead of answering only one question.")
    topic = st.text_input("Focus area", key="mentor_topic", placeholder="Physics problem solving", label_visibility="collapsed")
    goal = st.text_area("Goal", key="mentor_goal", height=90, placeholder="I want to become confident at solving Class 12 electrostatics questions.", label_visibility="collapsed")
    col1, col2 = st.columns(2)
    with col1:
        level = st.selectbox("Current level", ["Beginner", "Intermediate", "Advanced"], key="mentor_level")
    with col2:
        horizon = st.selectbox("Time horizon", ["4 weeks", "8 weeks", "12 weeks", "6 months"], key="mentor_horizon")

    if st.button("🌱 Build Mentor Plan", key="mentor_run_btn", use_container_width=True):
        if topic.strip() and goal.strip():
            with st.spinner("🌱 Building your long-term mentoring plan..."):
                st.session_state.ai_mentor_result = run_ai_mentor(topic.strip(), horizon, level, goal.strip())
            st.rerun()
        else:
            st.warning("Add both a focus area and a goal.")

    result = st.session_state.get("ai_mentor_result")
    if not result:
        return

    st.markdown("**🧭 Mentor Snapshot**")
    st.info(result["summary"])
    st.markdown("**📈 Growth Roadmap**")
    st.write(result["roadmap"])
    st.markdown("**✅ Weekly Coaching**")
    st.write(result["habits"])

def render_document_analyser_plugin(context: str = "") -> None:
    """Native renderer for document analysis."""
    if "doc_analyser_prompt" not in st.session_state:
        st.session_state.doc_analyser_prompt = context or ""
    if context and not st.session_state.doc_analyser_prompt.strip():
        st.session_state.doc_analyser_prompt = context

    st.caption("Upload a PDF or paste text, then ask for a summary, extraction, or analysis.")
    uploaded_doc = st.file_uploader(
        "Upload document",
        type=["pdf", "txt"],
        key="doc_analyser_upload",
        label_visibility="collapsed",
    )
    pasted_text = st.text_area(
        "Paste document text",
        key="doc_analyser_text",
        height=120,
        placeholder="Paste article, notes, report text, or document content here...",
        label_visibility="collapsed",
    )
    prompt = st.text_input(
        "Analysis request",
        key="doc_analyser_prompt",
        placeholder="Summarise the main points and extract important data.",
        label_visibility="collapsed",
    )

    if st.button("📋 Analyse Document", key="doc_analyser_run_btn", use_container_width=True):
        doc_text = pasted_text.strip()
        if uploaded_doc is not None:
            if uploaded_doc.name.lower().endswith(".pdf"):
                doc_text = extract_pdf_text(uploaded_doc) or ""
            else:
                try:
                    doc_text = uploaded_doc.getvalue().decode("utf-8", errors="ignore")
                except Exception:
                    doc_text = ""
        if doc_text and prompt.strip():
            with st.spinner("📋 Analysing document..."):
                result = call_hf(
                    "You are an expert document analyst. Summarise, extract key data, identify structure, and answer the user's request clearly.",
                    f"Document:\n{doc_text[:6000]}\n\nUser request: {prompt}\n\nAnalysis:",
                    max_tokens=1400,
                )
            st.session_state.doc_analyser_result = result or "No analysis returned."
            st.rerun()
        else:
            st.warning("Add a document and an analysis request.")

    if st.session_state.get("doc_analyser_result"):
        st.markdown("**📋 Document Analysis**")
        st.write(st.session_state.doc_analyser_result)

def render_video_reader_plugin(context: str = "") -> None:
    """Native renderer for transcript-based video analysis."""
    if "video_reader_prompt" not in st.session_state:
        st.session_state.video_reader_prompt = context or ""
    if context and not st.session_state.video_reader_prompt.strip():
        st.session_state.video_reader_prompt = context

    st.caption("Paste a video transcript or notes, optionally include the link, and get a structured breakdown.")
    video_url = st.text_input(
        "Video URL",
        key="video_reader_url",
        placeholder="https://www.youtube.com/watch?v=...",
        label_visibility="collapsed",
    )
    transcript = st.text_area(
        "Transcript or notes",
        key="video_reader_text",
        height=140,
        placeholder="Paste the transcript or detailed notes here...",
        label_visibility="collapsed",
    )
    prompt = st.text_input(
        "What should PhysIQ do with it?",
        key="video_reader_prompt",
        placeholder="Summarise the video and list key takeaways.",
        label_visibility="collapsed",
    )

    if st.button("🎬 Analyse Video", key="video_reader_run_btn", use_container_width=True):
        if transcript.strip() and prompt.strip():
            with st.spinner("🎬 Reading the video transcript..."):
                result = call_hf(
                    "You are an expert video analyst. Use the transcript and link context to summarise, extract insights, and answer the user's request.",
                    f"Video URL: {video_url or 'Not provided'}\n\nTranscript/notes:\n{transcript[:6000]}\n\nUser request: {prompt}\n\nVideo analysis:",
                    max_tokens=1400,
                )
            st.session_state.video_reader_result = result or "No video analysis returned."
            st.rerun()
        else:
            st.warning("Add a transcript or notes and an analysis request.")

    if st.session_state.get("video_reader_result"):
        st.markdown("**🎬 Video Analysis**")
        st.write(st.session_state.video_reader_result)

def render_custom_tool_plugin(plugin_key: str) -> None:
    """Render a generated custom tool plugin."""
    plugin = st.session_state.get("custom_plugins", {}).get(plugin_key)
    if not plugin:
        st.warning("This custom plugin is no longer available.")
        return
    st.caption(f"Built from: {plugin.get('request', 'Custom request')}")
    st.components.v1.html(plugin.get("html", ""), height=420, scrolling=True)
    st.download_button(
        "⬇️ Download Tool",
        data=plugin.get("html", ""),
        file_name=f"{plugin.get('name', 'custom_tool').lower().replace(' ', '_')}.html",
        mime="text/html",
        key=f"download_{plugin_key}",
        use_container_width=True,
    )

def render_plugin(plugin_key: str, context: str = "") -> None:
    """Render a plugin in Streamlit with AI message bridge."""
    custom_plugins = st.session_state.get("custom_plugins", {})
    if plugin_key in custom_plugins:
        p = custom_plugins[plugin_key]
    elif plugin_key in PLUGINS:
        p = PLUGINS[plugin_key]
    else:
        return
    col_title, col_min = st.columns([5, 1])
    with col_title:
        st.markdown(f"""<div style="color:{p['color']};font-weight:700;font-size:13px;margin-bottom:6px">
        {p['icon']} {p['name']} Plugin</div>""", unsafe_allow_html=True)

    if plugin_key in custom_plugins:
        render_custom_tool_plugin(plugin_key)
        return
    if plugin_key == "multi_agent":
        render_multi_agent_plugin(context)
        return
    if plugin_key == "idea_evolution":
        render_idea_evolution_plugin(context)
        return
    if plugin_key == "ai_mentor":
        render_ai_mentor_plugin(context)
        return
    if plugin_key == "document_analyser":
        render_document_analyser_plugin(context)
        return
    if plugin_key == "video_reader":
        render_video_reader_plugin(context)
        return
    if plugin_key == "tool_creator":
        render_tool_creator_plugin(context)
        return

    html_generators = {
        "calculator": lambda: plugin_calculator(context),
        "code_runner": lambda: plugin_code_runner(context),
        "decision_maker": lambda: plugin_decision_maker(context),
    }

    heights = {
        "calculator": 440, "code_runner": 520, "multi_agent": 550,
        "decision_maker": 480, "tool_creator": 500,
        "document_analyser": 200, "video_reader": 200,
        "hypothesis_generator": 480, "context_rewriter": 440,
        "multi_step_planner": 520, "idea_generator": 500,
    }

    if plugin_key in html_generators:
        # AI message bridge — allows plugins to call AI
        bridge = f"""
<script>
// Bridge: relay AI calls from plugin iframes to Streamlit
window.addEventListener('message', function(e) {{
    if(e.data && e.data.type === 'plugin_ai_call') {{
        // Forward to Streamlit via query params
        const p = new URL(window.location.href);
        p.searchParams.set('plugin_call', encodeURIComponent(JSON.stringify(e.data)));
        // Notify all child iframes of the response when it comes
        window.__pendingCalls = window.__pendingCalls || {{}};
        window.__pendingCalls[e.data.id] = e.source;
    }}
    if(e.data && e.data.type === 'plugin_result') {{
        // Bubble result up to Streamlit
        window.parent && window.parent.postMessage(e.data, '*');
    }}
}});
</script>"""

        html_content = html_generators[plugin_key]()
        # Replace </body> to inject bridge
        html_content = html_content.replace("</body>", bridge + "</body>")
        st.components.v1.html(html_content, height=heights.get(plugin_key, 400), scrolling=True)
    else:
        st.info(f"Plugin **{p['name']}** is available through the native renderer.")


def show_plugin_store():
    """Show the plugin store/manager."""
    st.markdown("### 🔌 Plugin Store")
    st.caption("Plugins are auto-used by the AI when relevant, or launch manually below")
    plugin_catalog = dict(PLUGINS)
    plugin_catalog.update(st.session_state.get("custom_plugins", {}))
    cols = st.columns(3)
    for i, (key, plugin) in enumerate(plugin_catalog.items()):
        with cols[i % 3]:
            st.markdown(f"""<div style="background:#161b22;border:1px solid {plugin['color']}40;
            border-top:3px solid {plugin['color']};border-radius:10px;padding:14px;margin-bottom:10px">
            <div style="font-size:1.4rem;margin-bottom:6px">{plugin['icon']}</div>
            <div style="color:{plugin['color']};font-weight:700;font-size:13px">{plugin['name']}</div>
            <div style="color:#8b949e;font-size:11px;margin:4px 0 10px">{plugin['description']}</div>
            </div>""", unsafe_allow_html=True)
            if st.button(f"▶ Launch {plugin['name']}", key=f"launch_{key}", use_container_width=True):
                st.session_state[f"plugin_open_{key}"] = True
                st.rerun()


# ══════════════════════════════════════════════════════════════
#  NEW FEATURES MODULE
#  1. Advanced Voice (Web Speech + Wit.AI fallback)
#  2. Code → PDF download
#  3. Live Data (News, Weather, Stocks)
#  4. App Integrations (Word/Excel/PPT/VS Code)
#  5. New Plugins (Hypothesis, Rewriter, Planner, Ideas)
# ══════════════════════════════════════════════════════════════

# ──────────────────────────────────────────────────────────────
# SECTION 1: ADVANCED VOICE MIC
# Uses Web Speech API (free, built into browser) + Wit.AI as
# a cloud fallback for better accuracy. AssemblyAI/Dialogflow
# are enterprise paid — we use their free alternatives.
# ──────────────────────────────────────────────────────────────

def show_advanced_voice_button():
    """
    Inline mic button above chat bar.
    - Primary:  Web Speech API (Chrome/Edge/Safari — free, no key)
    - Fallback: Wit.AI free tier (5 req/sec, no daily limit for basic)
    - TTS:      Web Speech Synthesis (free, many voices)
    Auto-sends after 2.5s silence. Speaks AI reply aloud.
    """
    WIT_TOKEN = os.getenv("WIT_TOKEN", "")   # Optional — falls back to Web Speech
    voice_html = r"""
<!DOCTYPE html><html><head>
<style>
*{margin:0;padding:0;box-sizing:border-box}
body{background:transparent;padding:3px 0;font-family:'Segoe UI',sans-serif}
.vrow{display:flex;align-items:center;gap:8px;justify-content:flex-end;padding:0 6px}
.mb{width:42px;height:42px;border-radius:50%;border:2px solid #30363d;cursor:pointer;
  background:linear-gradient(135deg,#1f6feb,#388bfd);color:#fff;
  font-size:18px;display:flex;align-items:center;justify-content:center;
  box-shadow:0 2px 10px rgba(88,166,255,0.3);transition:all .2s;flex-shrink:0}
.mb:hover{transform:scale(1.1);box-shadow:0 4px 16px rgba(88,166,255,0.5)}
.mb.L{background:linear-gradient(135deg,#b91c1c,#ef4444);border-color:#ef4444;
  animation:pu 1s infinite;box-shadow:0 0 0 0 rgba(239,68,68,.5)}
.mb.S{background:linear-gradient(135deg,#15803d,#22c55e);border-color:#22c55e}
.mb.P{background:linear-gradient(135deg,#7c3aed,#8b5cf6);border-color:#8b5cf6;animation:pu2 1.4s infinite}
@keyframes pu{0%{box-shadow:0 0 0 0 rgba(239,68,68,.5)}70%{box-shadow:0 0 0 14px rgba(239,68,68,0)}100%{box-shadow:0 0 0 0 rgba(239,68,68,0)}}
@keyframes pu2{0%{box-shadow:0 0 0 0 rgba(139,92,246,.5)}70%{box-shadow:0 0 0 14px rgba(139,92,246,0)}100%{box-shadow:0 0 0 0 rgba(139,92,246,0)}}
.pill{background:#161b22;border:1px solid #30363d;border-radius:20px;padding:6px 14px;
  color:#e6edf3;font-size:12px;max-width:300px;white-space:nowrap;overflow:hidden;
  text-overflow:ellipsis;display:none;flex-shrink:1}
.pill.v{display:block}
.sb{background:#238636;color:#fff;border:none;border-radius:8px;padding:6px 14px;
  font-size:12px;cursor:pointer;display:none;white-space:nowrap;flex-shrink:0}
.sb.v{display:block}
.sb:hover{background:#2ea043}
.tbar{position:absolute;bottom:-3px;left:0;right:0;height:3px;
  background:rgba(88,166,255,.12);border-radius:2px;display:none;overflow:hidden}
.tfill{height:100%;background:linear-gradient(90deg,#58a6ff,#22c55e);border-radius:2px;
  transition:width .08s linear}
.wrap{position:relative}
.vosel{background:#0d1117;border:1px solid #30363d;border-radius:6px;
  padding:3px 8px;color:#8b949e;font-size:10px;max-width:120px;cursor:pointer;display:none}
.vosel.v{display:block}
</style></head><body>
<div class="vrow">
  <select class="vosel" id="voSel" title="Voice"></select>
  <div class="pill" id="pi"></div>
  <button class="sb" id="sb">➤ Send</button>
  <div class="wrap">
    <button class="mb" id="mb" title="Click to speak · auto-sends after silence">🎙️</button>
    <div class="tbar" id="tb"><div class="tfill" id="tf" style="width:100%"></div></div>
  </div>
</div>
<script>
const mb=document.getElementById('mb'),pi=document.getElementById('pi'),
      sb=document.getElementById('sb'),tb=document.getElementById('tb'),
      tf=document.getElementById('tf'),voSel=document.getElementById('voSel');
let R=null,on=false,txt='',sT=null,iT=null,voices=[];
const SILS=2500; // ms silence before auto-send

// ── VOICE LIST ──────────────────────────────────────────────
const syn=window.speechSynthesis;
function loadVoices(){
  voices=syn.getVoices().filter(v=>v.lang.startsWith('en')||v.lang.startsWith('hi'));
  voSel.innerHTML='';
  const pref=['Google UK English Female','Microsoft Libby','Google US English','Samantha','Karen'];
  const sorted=[...voices].sort((a,b)=>{
    const ai=pref.indexOf(a.name),bi=pref.indexOf(b.name);
    if(ai>=0&&bi<0)return -1;if(bi>=0&&ai<0)return 1;
    if(a.localService&&!b.localService)return -1;if(!a.localService&&b.localService)return 1;
    return 0;
  });
  sorted.forEach(v=>{
    const o=document.createElement('option');
    o.value=v.name;o.textContent=v.name.replace('Microsoft ','').replace('Google ','');
    voSel.appendChild(o);
  });
  if(sorted.length){voSel.classList.add('v');}
}
syn.onvoiceschanged=loadVoices;loadVoices();

// ── SPEECH RECOGNITION ──────────────────────────────────────
function go(){
  const SR=window.SpeechRecognition||window.webkitSpeechRecognition;
  if(!SR){
    pi.textContent='⚠️ Use Chrome or Edge for voice';pi.classList.add('v');
    setTimeout(()=>pi.classList.remove('v'),3000);return;
  }
  R=new SR();R.continuous=true;R.interimResults=true;R.lang='en-US';R.maxAlternatives=1;
  R.onstart=()=>{
    on=true;txt='';mb.className='mb L';mb.textContent='🔴';
    pi.textContent='🎙️ Listening…';pi.classList.add('v');
    sb.classList.remove('v');tb.style.display='none';
  };
  R.onresult=(e)=>{
    let f='',i2='';
    for(let k=e.resultIndex;k<e.results.length;k++){
      if(e.results[k].isFinal)f+=e.results[k][0].transcript;
      else i2+=e.results[k][0].transcript;
    }
    if(f)txt+=f;
    pi.textContent='"'+(txt||i2)+'"';
    // Reset silence countdown
    clearTimeout(sT);clearInterval(iT);
    let p=100;tf.style.width='100%';tb.style.display='block';
    iT=setInterval(()=>{p-=(100/SILS)*80;tf.style.width=Math.max(0,p)+'%';if(p<=0)clearInterval(iT);},80);
    sT=setTimeout(()=>{clearInterval(iT);if(txt.trim())send();else R.stop();},SILS);
  };
  R.onend=()=>{
    on=false;tb.style.display='none';
    if(txt.trim()){mb.className='mb S';mb.textContent='✅';sb.classList.add('v');}
    else{mb.className='mb';mb.textContent='🎙️';pi.classList.remove('v');}
  };
  R.onerror=(e)=>{
    on=false;mb.className='mb';mb.textContent='🎙️';
    const m={'not-allowed':'🔒 Allow mic access in browser','no-speech':'No speech — try again',
             'audio-capture':'No mic found','network':'Network error'};
    pi.textContent=m[e.error]||'Error: '+e.error;pi.classList.add('v');
    setTimeout(()=>pi.classList.remove('v'),3000);
  };
  try{R.start();}catch(e){console.error(e);}
}

function send(){
  clearTimeout(sT);clearInterval(iT);
  if(R&&on){try{R.stop();}catch(e){}}
  const q=txt.trim();if(!q)return;
  pi.textContent='⏳ Sending to PhysIQ…';sb.classList.remove('v');
  mb.className='mb';mb.textContent='🎙️';txt='';
  // Write to URL param → triggers Streamlit rerun
  const u=new URL(window.parent.location.href);
  u.searchParams.set('vq',encodeURIComponent(q));
  window.parent.history.replaceState({},'',u.toString());
  // Force navigation to trigger rerun
  setTimeout(()=>{window.parent.location.href=u.toString();},180);
}

mb.onclick=()=>{if(on){clearTimeout(sT);send();}else go();};
sb.onclick=()=>send();

// ── TTS REPLY ───────────────────────────────────────────────
window.addEventListener('message',e=>{
  if(e.data?.type==='physiq_tts'){
    syn.cancel();
    const u=new SpeechSynthesisUtterance(
      (e.data.text||'').replace(/[#*_`>[\]]/g,'').replace(/\n/g,' ').substring(0,500)
    );
    u.rate=parseFloat(e.data.rate||1.0);u.pitch=parseFloat(e.data.pitch||1.0);
    const sel=voSel.value;
    const v=voices.find(x=>x.name===sel);
    if(v)u.voice=v;
    u.onstart=()=>{mb.className='mb P';mb.textContent='🔊';};
    u.onend=()=>{mb.className='mb';mb.textContent='🎙️';};
    syn.speak(u);
  }
  if(e.data?.type==='physiq_tts_stop')syn.cancel();
});
</script></body></html>"""
    st.components.v1.html(voice_html, height=56, scrolling=False)


# ──────────────────────────────────────────────────────────────
# SECTION 2: CODE → PDF DOWNLOAD
# ──────────────────────────────────────────────────────────────

def code_to_pdf_bytes(code_text: str, language: str = "code") -> bytes:
    """Convert code string to a downloadable PDF using reportlab."""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
        from reportlab.lib.units import cm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Preformatted
        from reportlab.lib.enums import TA_LEFT
        import io

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4,
            leftMargin=1.5*cm, rightMargin=1.5*cm,
            topMargin=2*cm, bottomMargin=2*cm)

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle('Title',
            fontName='Helvetica-Bold', fontSize=16,
            textColor=colors.HexColor('#1f6feb'),
            spaceAfter=12)
        code_style = ParagraphStyle('Code',
            fontName='Courier', fontSize=8,
            leading=11, textColor=colors.HexColor('#e6edf3'),
            backColor=colors.HexColor('#0d1117'),
            leftIndent=6, rightIndent=6, spaceBefore=2, spaceAfter=2)
        meta_style = ParagraphStyle('Meta',
            fontName='Helvetica', fontSize=9,
            textColor=colors.HexColor('#8b949e'), spaceAfter=16)

        story = []
        story.append(Paragraph(f"⚛️ PhysIQ — Generated {language} Code", title_style))
        story.append(Paragraph(
            f"Generated: {datetime.datetime.now().strftime('%d %b %Y %H:%M')} | Language: {language}",
            meta_style))

        # Split into lines and add as preformatted
        lines = code_text.split('\n')
        for i in range(0, len(lines), 60):   # chunk to avoid memory issues
            chunk = '\n'.join(lines[i:i+60])
            # Escape XML chars
            chunk = chunk.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')
            story.append(Preformatted(chunk, code_style))

        doc.build(story)
        return buf.getvalue()
    except Exception as e:
        return None

def extract_code_blocks(text: str):
    """Extract all code blocks from markdown text."""
    import re
    # Match ```lang\n...\n``` blocks
    pattern = r'```(\w*)\n([\s\S]*?)```'
    blocks = re.findall(pattern, text)
    if not blocks:
        # Try without language tag
        pattern2 = r'```([\s\S]*?)```'
        raw = re.findall(pattern2, text)
        blocks = [('code', b) for b in raw]
    return blocks  # list of (language, code)

def show_code_download_button(answer_text: str, lang: str = "code"):
    """Show PDF download button if answer contains code blocks."""
    blocks = extract_code_blocks(answer_text)
    if not blocks:
        return
    all_code = "\n\n".join([f"# --- {b[0] or lang} ---\n{b[1]}" for b in blocks])
    pdf_bytes = code_to_pdf_bytes(all_code, lang)
    if pdf_bytes:
        fname = f"physiq_{lang.lower().replace(' ','_')}_{datetime.datetime.now().strftime('%H%M%S')}.pdf"
        st.download_button(
            label="📄 Download Code as PDF",
            data=pdf_bytes,
            file_name=fname,
            mime="application/pdf",
            help="Download all code blocks from this answer as a formatted PDF",
            use_container_width=False,
        )


# ──────────────────────────────────────────────────────────────
# SECTION 3: LIVE DATA (News, Weather, Stocks)
# All using free APIs — no keys for basic usage
# ──────────────────────────────────────────────────────────────

LIVE_DATA_TRIGGERS = {
    "news": ["news","headlines","latest news","what happened","today in","breaking","current events"],
    "weather": ["weather","temperature","forecast","rain","sunny","cold","hot","climate today"],
    "stock": ["stock","share price","market","nasdaq","nse","bse","sensex","nifty","crypto",
              "bitcoin","ethereum","dow jones","s&p"],
    "cricket": ["cricket","ipl","test match","score","odi","t20","match result","live score"],
}

def needs_live_data(question: str):
    q = question.lower()
    for dtype, triggers in LIVE_DATA_TRIGGERS.items():
        if any(t in q for t in triggers):
            return dtype
    return None

def fetch_news(query="top headlines", country="in"):
    """Free news via GNews API (100 req/day free) or Google RSS fallback."""
    try:
        import urllib.request, urllib.parse, json
        GNEWS = os.getenv("GNEWS_TOKEN","")
        if GNEWS:
            # GNews paid tier for better results
            gnews_url = f"https://gnews.io/api/v4/search?q={urllib.parse.quote(query)}&lang=en&country={country}&max=5&apikey={GNEWS}"
            req0 = urllib.request.Request(gnews_url, headers={"User-Agent":"PhysIQ/1.0"})
            try:
                with urllib.request.urlopen(req0, timeout=8) as r0:
                    gdata = json.loads(r0.read())
                articles = gdata.get("articles",[])[:5]
                if articles:
                    return "📰 **Latest News:**\n" + "\n".join([f"• {a['title']} — _{a.get('source',{{}}).get('name','')}_" for a in articles])
            except: pass
        # Free RSS fallback
        rss_url = f"https://news.google.com/rss/search?q={urllib.parse.quote(query)}&hl=en-IN&gl=IN&ceid=IN:en"
        req = urllib.request.Request(rss_url, headers={"User-Agent":"PhysIQ/1.0"})
        with urllib.request.urlopen(req, timeout=8) as r:
            xml = r.read().decode("utf-8","ignore")
        import re as _re
        titles = _re.findall(r'<title><!\[CDATA\[(.*?)\]\]></title>', xml)[1:6]
        if not titles:
            titles = _re.findall(r'<title>(.*?)</title>', xml)[1:6]
        if titles:
            return "📰 **Latest News:**\n" + "\n".join([f"• {t}" for t in titles])
    except: pass
    return None

def fetch_weather(city="Delhi"):
    """Free weather via Open-Meteo (no API key ever needed!)."""
    try:
        import urllib.request, json, urllib.parse
        # Step 1: geocode
        gc_url = f"https://geocoding-api.open-meteo.com/v1/search?name={urllib.parse.quote(city)}&count=1&language=en&format=json"
        req = urllib.request.Request(gc_url, headers={"User-Agent":"PhysIQ/1.0"})
        with urllib.request.urlopen(req, timeout=6) as r:
            gc = json.loads(r.read())
        if not gc.get("results"):
            return None
        loc = gc["results"][0]
        lat, lon = loc["latitude"], loc["longitude"]
        name = loc.get("name", city)
        # Step 2: weather
        wx_url = (f"https://api.open-meteo.com/v1/forecast?latitude={lat}&longitude={lon}"
                  f"&current=temperature_2m,weathercode,windspeed_10m,relativehumidity_2m"
                  f"&daily=temperature_2m_max,temperature_2m_min,precipitation_sum"
                  f"&timezone=auto&forecast_days=3")
        req2 = urllib.request.Request(wx_url, headers={"User-Agent":"PhysIQ/1.0"})
        with urllib.request.urlopen(req2, timeout=6) as r2:
            wx = json.loads(r2.read())
        curr = wx.get("current",{})
        temp = curr.get("temperature_2m","?")
        wind = curr.get("windspeed_10m","?")
        hum  = curr.get("relativehumidity_2m","?")
        code = curr.get("weathercode",0)
        icons = {0:"☀️",1:"🌤",2:"⛅",3:"☁️",45:"🌫",51:"🌦",61:"🌧",71:"❄️",80:"🌧",95:"⛈"}
        icon = icons.get(code, "🌡️")
        daily = wx.get("daily",{})
        forecasts = ""
        if daily.get("temperature_2m_max"):
            for i,d in enumerate(daily.get("time",[])[:3]):
                mx = daily["temperature_2m_max"][i]
                mn = daily["temperature_2m_min"][i]
                pr = daily.get("precipitation_sum",[0]*3)[i]
                forecasts += f"\n• {d}: {mn}°C – {mx}°C {'🌧' if pr>0 else '☀️'}"
        return (f"🌤️ **Weather in {name}:**\n"
                f"{icon} **{temp}°C** | 💨 Wind: {wind} km/h | 💧 Humidity: {hum}%\n"
                f"\n**3-Day Forecast:**{forecasts}")
    except Exception as e:
        return None

def fetch_stock(symbol="NIFTY"):
    """Free stock data via Yahoo Finance API (no key needed)."""
    try:
        import urllib.request, json, urllib.parse
        # Map common names to tickers
        ticker_map = {
            "nifty":"^NSEI","sensex":"^BSESN","nse":"^NSEI","bse":"^BSESN",
            "bitcoin":"BTC-USD","btc":"BTC-USD","ethereum":"ETH-USD","eth":"ETH-USD",
            "dow jones":"^DJI","s&p":"^GSPC","nasdaq":"^IXIC",
            "apple":"AAPL","google":"GOOGL","microsoft":"MSFT","tesla":"TSLA",
            "reliance":"RELIANCE.NS","tcs":"TCS.NS","infosys":"INFY.NS",
        }
        sym_key = symbol.lower().strip()
        ticker = ticker_map.get(sym_key, symbol.upper())
        url = f"https://query1.finance.yahoo.com/v8/finance/chart/{urllib.parse.quote(ticker)}?interval=1d&range=2d"
        req = urllib.request.Request(url, headers={
            "User-Agent":"Mozilla/5.0","Accept":"application/json"
        })
        with urllib.request.urlopen(req, timeout=8) as r:
            data = json.loads(r.read())
        result = data.get("chart",{}).get("result",[])
        if result:
            meta = result[0].get("meta",{})
            price = meta.get("regularMarketPrice","?")
            prev  = meta.get("chartPreviousClose", price)
            curr_name = meta.get("shortName", ticker)
            try:
                chg = float(price) - float(prev)
                pct = (chg/float(prev))*100
                arrow = "📈" if chg>=0 else "📉"
                return (f"💹 **{curr_name}** ({ticker})\n"
                        f"{arrow} **{price}** | Change: {chg:+.2f} ({pct:+.2f}%)\n"
                        f"Prev close: {prev}")
            except:
                return f"💹 **{ticker}**: {price}"
    except Exception as e:
        return None

def fetch_live_data(question: str, dtype: str):
    """Route to correct live data fetcher."""
    q = question.lower()
    if dtype == "weather":
        # Extract city name
        import re as _re
        cities = _re.findall(r'weather (?:in |at |for )?([A-Za-z\s]+?)(?:\?|$|\.)', question, _re.I)
        city = cities[0].strip() if cities else "Delhi"
        return fetch_weather(city)
    elif dtype == "stock":
        # Extract symbol
        import re as _re
        words = question.split()
        for w in words:
            r = fetch_stock(w)
            if r: return r
        return fetch_stock("NIFTY")
    elif dtype == "news":
        # Extract topic
        for trigger in LIVE_DATA_TRIGGERS["news"]:
            q2 = q.replace(trigger,"").strip()
        topic = q2 or "India technology science"
        return fetch_news(topic)
    elif dtype == "cricket":
        return fetch_news("cricket score live")
    return None


# ──────────────────────────────────────────────────────────────
# SECTION 4: APP INTEGRATIONS
# Deep-links and export helpers for external apps
# ──────────────────────────────────────────────────────────────

def show_app_integration_buttons(content_text: str, content_type: str = "text"):
    """Show export buttons for Microsoft Office, VS Code, etc."""
    import base64, urllib.parse

    st.markdown("**📤 Export to:**")
    cols = st.columns(6)

    # ── Word (.docx) ─────────────────────────────────────────
    with cols[0]:
        try:
            from docx import Document as DocxDoc
            import io
            doc = DocxDoc()
            doc.add_heading("PhysIQ Export", 0)
            doc.add_paragraph(content_text)
            buf = io.BytesIO()
            doc.save(buf)
            st.download_button("📝 Word", buf.getvalue(),
                file_name="physiq_export.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True)
        except:
            if st.button("📝 Word", use_container_width=True):
                st.info("Install: pip install python-docx")

    # ── Excel (.xlsx) ────────────────────────────────────────
    with cols[1]:
        try:
            import openpyxl, io
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "PhysIQ"
            for i, line in enumerate(content_text.split('\n'), 1):
                ws.cell(row=i, column=1, value=line)
            buf = io.BytesIO()
            wb.save(buf)
            st.download_button("📊 Excel", buf.getvalue(),
                file_name="physiq_export.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True)
        except:
            if st.button("📊 Excel", use_container_width=True):
                st.info("Install: pip install openpyxl")

    # ── PowerPoint (.pptx) ───────────────────────────────────
    with cols[2]:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            import io
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "PhysIQ Export"
            tf = slide.placeholders[1].text_frame
            tf.text = content_text[:500]
            buf = io.BytesIO()
            prs.save(buf)
            st.download_button("🎞️ PPT", buf.getvalue(),
                file_name="physiq_export.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True)
        except:
            if st.button("🎞️ PPT", use_container_width=True):
                st.info("Install: pip install python-pptx")

    # ── VS Code (open as file) ────────────────────────────────
    with cols[3]:
        # Download .py file — user can open in VS Code/PyCharm
        ext = ".py" if "python" in content_type.lower() else ".txt"
        st.download_button(f"💻 VS Code{ext}", content_text.encode(),
            file_name=f"physiq_code{ext}", mime="text/plain",
            use_container_width=True)

    # ── Programiz (web-based IDE link) ───────────────────────
    with cols[4]:
        encoded = urllib.parse.quote(content_text[:2000])
        programiz_url = f"https://www.programiz.com/python-programming/online-compiler/"
        if st.button("🟢 Programiz", use_container_width=True):
            st.markdown(f"[Open Programiz IDE →]({programiz_url})", unsafe_allow_html=True)
            st.code(content_text[:200], language="python")

    # ── Plain Text (.txt) ─────────────────────────────────────
    with cols[5]:
        st.download_button("📄 .txt", content_text.encode(),
            file_name="physiq_export.txt", mime="text/plain",
            use_container_width=True)


# ──────────────────────────────────────────────────────────────
# SECTION 5: NEW PLUGIN HTML GENERATORS
# ──────────────────────────────────────────────────────────────

def plugin_hypothesis_generator(topic=""):
    return f"""<!DOCTYPE html><html><head>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;font-family:'Segoe UI',sans-serif;padding:14px;color:#e6edf3}}
.wrap{{background:#161b22;border:1px solid #06d6a0;border-radius:14px;padding:16px}}
.title{{color:#06d6a0;font-weight:700;font-size:13px;margin-bottom:4px}}
.sub{{color:#8b949e;font-size:11px;margin-bottom:12px}}
.row{{display:flex;gap:8px;margin-bottom:10px}}
.inp{{flex:1;background:#0d1117;border:1px solid #30363d;border-radius:8px;padding:8px 12px;color:#e6edf3;font-size:12px}}
.btn{{background:#06d6a0;color:#000;border:none;border-radius:8px;padding:8px 16px;cursor:pointer;font-weight:700;font-size:12px}}
.hyp-card{{background:#0d1117;border:1px solid #30363d;border-radius:10px;padding:12px;margin-bottom:8px}}
.hyp-type{{font-size:10px;font-weight:700;letter-spacing:1px;margin-bottom:4px}}
.hyp-text{{font-size:12px;color:#c9d1d9;line-height:1.6}}
.hyp-test{{font-size:11px;color:#8b949e;margin-top:4px;font-style:italic}}
.chips{{display:flex;flex-wrap:wrap;gap:6px;margin-bottom:12px}}
.chip{{background:#21262d;border:1px solid #30363d;border-radius:12px;padding:4px 10px;font-size:11px;color:#8b949e;cursor:pointer}}
.chip:hover{{border-color:#06d6a0;color:#06d6a0}}
</style></head><body>
<div class="wrap">
<div class="title">🔬 Hypothesis Generator</div>
<div class="sub">Generate testable scientific hypotheses from any topic</div>
<div class="chips">
  <div class="chip" onclick="setT('Effect of temperature on enzyme activity')">🧪 Enzyme activity</div>
  <div class="chip" onclick="setT('Relationship between light intensity and photosynthesis')">🌱 Photosynthesis</div>
  <div class="chip" onclick="setT('Impact of gravity on projectile motion')">🎯 Projectile motion</div>
  <div class="chip" onclick="setT('Effect of concentration on reaction rate')">⚗️ Reaction rate</div>
</div>
<div class="row">
  <input class="inp" id="topicIn" placeholder="Enter a scientific topic..." value="{topic}">
  <button class="btn" onclick="gen()">Generate</button>
</div>
<div id="results"></div>
</div>
<script>
function setT(v){{document.getElementById('topicIn').value=v;}}
async function gen(){{
  const topic=document.getElementById('topicIn').value.trim();
  if(!topic)return;
  const res=document.getElementById('results');
  res.innerHTML='<div style="color:#8b949e;font-size:12px;padding:8px">🔬 Generating hypotheses...</div>';
  window.parent.postMessage({{type:'plugin_ai_call',id:'hyp',
    system:`You are a scientific research expert. Generate 4 testable scientific hypotheses for the given topic.
For each hypothesis return a JSON array with objects having:
- "type": "Null Hypothesis" or "Alternative Hypothesis" or "Directional" or "Non-directional"
- "hypothesis": the exact hypothesis statement (If-Then format)
- "test": how to test this hypothesis
- "variables": "Independent: X | Dependent: Y | Controlled: Z"
Return ONLY a valid JSON array, no markdown.`,
    user:'Topic: '+topic+'\n\nGenerate 4 scientific hypotheses:',max_tokens:600}}
  }},'*');
}}
window.addEventListener('message',e=>{{
  if(e.data?.type==='plugin_ai_response'&&e.data?.id==='hyp'){{
    const raw=e.data.result;
    const res=document.getElementById('results');
    try{{
      const s=raw.indexOf('['),en=raw.lastIndexOf(']')+1;
      const data=JSON.parse(raw.slice(s,en));
      const colors={{
        'Null Hypothesis':'#f85149','Alternative Hypothesis':'#06d6a0',
        'Directional':'#58a6ff','Non-directional':'#ffd166'
      }};
      res.innerHTML=data.map(h=>`
        <div class="hyp-card">
          <div class="hyp-type" style="color:${{colors[h.type]||'#06d6a0'}}">${{h.type?.toUpperCase()}}</div>
          <div class="hyp-text">${{h.hypothesis}}</div>
          <div class="hyp-test">🧪 Test: ${{h.test}}</div>
          <div class="hyp-test">📊 Variables: ${{h.variables}}</div>
        </div>`).join('');
    }}catch(e){{res.innerHTML='<div style="color:#f85149;font-size:12px;padding:8px">Parse error — try again</div>';}}
    window.parent.postMessage({{type:'plugin_result',plugin:'hypothesis_generator',result:e.data.result}},'*');
  }}
}});
if(document.getElementById('topicIn').value)setTimeout(gen,300);
</script></body></html>"""


def plugin_context_rewriter(query=""):
    return f"""<!DOCTYPE html><html><head>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;font-family:'Segoe UI',sans-serif;padding:14px;color:#e6edf3}}
.wrap{{background:#161b22;border:1px solid #a78bfa;border-radius:14px;padding:16px}}
.title{{color:#a78bfa;font-weight:700;font-size:13px;margin-bottom:4px}}
.sub{{color:#8b949e;font-size:11px;margin-bottom:12px}}
textarea{{width:100%;background:#0d1117;border:1px solid #30363d;border-radius:8px;
  padding:10px;color:#e6edf3;font-size:12px;height:80px;resize:vertical;margin-bottom:8px}}
.btn{{background:#a78bfa;color:#000;border:none;border-radius:8px;padding:8px 18px;cursor:pointer;font-weight:700;font-size:12px;margin-bottom:12px}}
.versions{{display:grid;grid-template-columns:repeat(2,1fr);gap:10px}}
.ver{{background:#0d1117;border:1px solid #30363d;border-radius:10px;padding:12px}}
.ver-title{{font-size:10px;font-weight:700;letter-spacing:1px;margin-bottom:6px}}
.ver-text{{font-size:12px;color:#c9d1d9;line-height:1.6}}
.use-btn{{margin-top:8px;background:#21262d;border:1px solid #30363d;color:#a78bfa;
  border-radius:6px;padding:4px 10px;font-size:11px;cursor:pointer;width:100%}}
.use-btn:hover{{background:#a78bfa;color:#000}}
</style></head><body>
<div class="wrap">
<div class="title">✏️ Context Rewriter</div>
<div class="sub">Rewrite your query to get better AI answers — 4 optimised versions</div>
<textarea id="queryIn" placeholder="Enter your question or query here...">{query}</textarea>
<button class="btn" onclick="rewrite()">✨ Rewrite for Better AI Understanding</button>
<div class="versions" id="vers"></div>
</div>
<script>
async function rewrite(){{
  const q=document.getElementById('queryIn').value.trim();if(!q)return;
  document.getElementById('vers').innerHTML='<div style="color:#8b949e;font-size:12px;padding:8px;grid-column:span 2">✨ Rewriting...</div>';
  window.parent.postMessage({{type:'plugin_ai_call',id:'ctx',
    system:`You are an expert prompt engineer. Rewrite the user query 4 different ways to get the best AI response.
Return a JSON array with 4 objects, each having:
- "style": "Specific & Precise" | "Step-by-Step" | "With Context" | "Socratic"
- "rewritten": the rewritten query
- "why": one sentence why this version works better
Return ONLY valid JSON array.`,
    user:'Original query: "'+q+'"\n\nRewrite in 4 styles:',max_tokens:500}}
  }},'*');
}}
window.addEventListener('message',e=>{{
  if(e.data?.type==='plugin_ai_response'&&e.data?.id==='ctx'){{
    const raw=e.data.result;const vers=document.getElementById('vers');
    try{{
      const s=raw.indexOf('['),en=raw.lastIndexOf(']')+1;
      const data=JSON.parse(raw.slice(s,en));
      const colors=['#58a6ff','#06d6a0','#ffd166','#f471b5'];
      vers.innerHTML=data.map((v,i)=>`
        <div class="ver">
          <div class="ver-title" style="color:${{colors[i]}}">${{v.style?.toUpperCase()}}</div>
          <div class="ver-text">${{v.rewritten}}</div>
          <div style="color:#8b949e;font-size:10px;margin-top:4px">${{v.why}}</div>
          <button class="use-btn" onclick="useQuery('${{v.rewritten?.replace(/'/g,"\\'")}}')">Use this version →</button>
        </div>`).join('');
    }}catch(err){{vers.innerHTML='<div style="color:#f85149;font-size:12px;padding:8px;grid-column:span 2">Parse error</div>';}}
  }}
}});
function useQuery(q){{
  const u=new URL(window.parent.location.href);
  u.searchParams.set('vq',encodeURIComponent(q));
  window.parent.history.replaceState({{}},'',u.toString());
  setTimeout(()=>{{window.parent.location.href=u.toString();}},100);
}}
if(document.getElementById('queryIn').value)setTimeout(rewrite,300);
</script></body></html>"""


def plugin_multi_step_planner(goal=""):
    return f"""<!DOCTYPE html><html><head>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;font-family:'Segoe UI',sans-serif;padding:14px;color:#e6edf3}}
.wrap{{background:#161b22;border:1px solid #ffd166;border-radius:14px;padding:16px}}
.title{{color:#ffd166;font-weight:700;font-size:13px;margin-bottom:4px}}
.sub{{color:#8b949e;font-size:11px;margin-bottom:12px}}
.row{{display:flex;gap:8px;margin-bottom:10px}}
.inp{{flex:1;background:#0d1117;border:1px solid #30363d;border-radius:8px;padding:8px 12px;color:#e6edf3;font-size:12px}}
.btn{{background:#ffd166;color:#000;border:none;border-radius:8px;padding:8px 16px;cursor:pointer;font-weight:700;font-size:12px}}
.timeline{{margin-top:10px;position:relative}}
.step{{display:flex;gap:12px;margin-bottom:12px;align-items:flex-start}}
.step-num{{width:28px;height:28px;border-radius:50%;background:#ffd166;color:#000;
  font-weight:700;font-size:12px;display:flex;align-items:center;justify-content:center;flex-shrink:0}}
.step-body{{flex:1;background:#0d1117;border:1px solid #30363d;border-radius:8px;padding:10px}}
.step-title{{font-weight:600;font-size:12px;color:#ffd166;margin-bottom:4px}}
.step-desc{{font-size:11px;color:#c9d1d9;line-height:1.6}}
.step-time{{font-size:10px;color:#8b949e;margin-top:4px}}
.done-btn{{background:none;border:1px solid #30363d;color:#3fb950;border-radius:4px;
  padding:2px 8px;font-size:10px;cursor:pointer;margin-top:6px}}
.done-btn:hover{{background:#3fb950;color:#000}}
.progress-bar{{background:#21262d;border-radius:4px;height:8px;margin-bottom:10px;overflow:hidden}}
.progress-fill{{height:100%;background:linear-gradient(90deg,#ffd166,#ff8c42);border-radius:4px;transition:width .3s}}
</style></head><body>
<div class="wrap">
<div class="title">📋 Multi-Step Planner</div>
<div class="sub">Break any goal into a clear, actionable step-by-step plan</div>
<div class="row">
  <input class="inp" id="goalIn" placeholder="Enter your goal e.g. 'Learn quantum mechanics from scratch'" value="{goal}">
  <button class="btn" onclick="plan()">📋 Plan It</button>
</div>
<div class="progress-bar"><div class="progress-fill" id="prog" style="width:0%"></div></div>
<div class="timeline" id="timeline"></div>
</div>
<script>
let total=0,done=0;
function markDone(btn,i){{done++;btn.textContent='✅ Done';btn.disabled=true;
  document.getElementById('prog').style.width=(done/total*100)+'%';}}
async function plan(){{
  const goal=document.getElementById('goalIn').value.trim();if(!goal)return;
  const tl=document.getElementById('timeline');
  tl.innerHTML='<div style="color:#8b949e;font-size:12px">📋 Creating plan...</div>';
  done=0;
  window.parent.postMessage({{type:'plugin_ai_call',id:'plan',
    system:`You are an expert life and study coach. Break the goal into a detailed step-by-step plan.
Return ONLY a JSON array of steps (6-10 steps), each with:
- "title": step title (short)
- "description": what to do in detail
- "duration": estimated time (e.g. "2 hours", "1 week")
- "difficulty": "Easy" | "Medium" | "Hard"
- "resources": helpful resources or tips (one line)
Return ONLY valid JSON array.`,
    user:'Goal: "'+goal+'"\n\nCreate a detailed action plan:',max_tokens:800}}
  }},'*');
}}
window.addEventListener('message',e=>{{
  if(e.data?.type==='plugin_ai_response'&&e.data?.id==='plan'){{
    const raw=e.data.result;const tl=document.getElementById('timeline');
    try{{
      const s=raw.indexOf('['),en=raw.lastIndexOf(']')+1;
      const data=JSON.parse(raw.slice(s,en));
      total=data.length;
      const dclr={{Easy:'#3fb950',Medium:'#ffd166',Hard:'#f85149'}};
      tl.innerHTML=data.map((step,i)=>`
        <div class="step">
          <div class="step-num">${{i+1}}</div>
          <div class="step-body">
            <div class="step-title">${{step.title}}</div>
            <div class="step-desc">${{step.description}}</div>
            <div class="step-time">⏱ ${{step.duration}} &nbsp;|&nbsp; <span style="color:${{dclr[step.difficulty]||'#ffd166'}}">${{step.difficulty}}</span></div>
            ${{step.resources?`<div class="step-time">💡 ${{step.resources}}</div>`:''}}
            <button class="done-btn" onclick="markDone(this,${{i}})">Mark Done</button>
          </div>
        </div>`).join('');
    }}catch(err){{tl.innerHTML='<div style="color:#f85149;font-size:12px">Parse error</div>';}}
    window.parent.postMessage({{type:'plugin_result',plugin:'multi_step_planner',result:e.data.result}},'*');
  }}
}});
if(document.getElementById('goalIn').value)setTimeout(plan,300);
</script></body></html>"""


def plugin_idea_generator(domain=""):
    return f"""<!DOCTYPE html><html><head>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;font-family:'Segoe UI',sans-serif;padding:14px;color:#e6edf3}}
.wrap{{background:#161b22;border:1px solid #f471b5;border-radius:14px;padding:16px}}
.title{{color:#f471b5;font-weight:700;font-size:13px;margin-bottom:4px}}
.sub{{color:#8b949e;font-size:11px;margin-bottom:12px}}
.row{{display:flex;gap:8px;margin-bottom:10px}}
.inp{{flex:1;background:#0d1117;border:1px solid #30363d;border-radius:8px;padding:8px 12px;color:#e6edf3;font-size:12px}}
.btn{{background:#f471b5;color:#000;border:none;border-radius:8px;padding:8px 16px;cursor:pointer;font-weight:700;font-size:12px}}
.types{{display:flex;gap:6px;flex-wrap:wrap;margin-bottom:10px}}
.type-btn{{background:#21262d;border:1px solid #30363d;border-radius:6px;padding:5px 12px;
  font-size:11px;color:#8b949e;cursor:pointer}}
.type-btn.active{{background:#f471b5;color:#000;border-color:#f471b5}}
.idea-grid{{display:grid;grid-template-columns:repeat(2,1fr);gap:10px;margin-top:10px}}
.idea-card{{background:#0d1117;border:1px solid #30363d;border-radius:10px;padding:12px;
  transition:border-color .2s;cursor:pointer}}
.idea-card:hover{{border-color:#f471b5}}
.idea-name{{font-size:12px;font-weight:700;color:#f471b5;margin-bottom:4px}}
.idea-desc{{font-size:11px;color:#c9d1d9;line-height:1.5;margin-bottom:6px}}
.idea-tags{{display:flex;flex-wrap:wrap;gap:4px}}
.tag{{background:#21262d;border-radius:10px;padding:2px 8px;font-size:10px;color:#8b949e}}
.score{{font-size:10px;margin-top:4px}}
</style></head><body>
<div class="wrap">
<div class="title">💡 Idea Generator</div>
<div class="sub">Generate startup, project, or research ideas with market analysis</div>
<div class="types">
  <div class="type-btn active" id="t-startup" onclick="setType('startup',this)">🚀 Startup</div>
  <div class="type-btn" id="t-project" onclick="setType('project',this)">🛠 Project</div>
  <div class="type-btn" id="t-research" onclick="setType('research',this)">🔬 Research</div>
  <div class="type-btn" id="t-app" onclick="setType('app',this)">📱 App</div>
</div>
<div class="row">
  <input class="inp" id="domainIn" placeholder="Enter domain e.g. 'Education', 'Healthcare', 'Physics'" value="{domain}">
  <button class="btn" onclick="generate()">💡 Generate</button>
</div>
<div class="idea-grid" id="ideas"></div>
</div>
<script>
let ideaType='startup';
function setType(t,el){{ideaType=t;document.querySelectorAll('.type-btn').forEach(b=>b.classList.remove('active'));el.classList.add('active');}}
async function generate(){{
  const domain=document.getElementById('domainIn').value.trim();if(!domain)return;
  const ig=document.getElementById('ideas');
  ig.innerHTML='<div style="color:#8b949e;font-size:12px;padding:8px;grid-column:span 2">💡 Generating ideas...</div>';
  window.parent.postMessage({{type:'plugin_ai_call',id:'idea',
    system:`You are a serial entrepreneur and innovation expert. Generate 6 creative ${{ideaType}} ideas for the given domain.
Return ONLY a JSON array with 6 objects, each having:
- "name": catchy idea name
- "description": 2-3 sentence description of the idea
- "problem": what problem it solves
- "target": who is the target user/customer
- "revenue": how it makes money (one line)
- "difficulty": "Easy" | "Medium" | "Hard" to build
- "market": market size estimate
- "tags": array of 3 keyword tags
Return ONLY valid JSON array.`,
    user:'Domain: "'+domain+'" | Type: '+ideaType+'\n\nGenerate 6 innovative ideas:',max_tokens:1200}}
  }},'*');
}}
window.addEventListener('message',e=>{{
  if(e.data?.type==='plugin_ai_response'&&e.data?.id==='idea'){{
    const raw=e.data.result;const ig=document.getElementById('ideas');
    try{{
      const s=raw.indexOf('['),en=raw.lastIndexOf(']')+1;
      const data=JSON.parse(raw.slice(s,en));
      const dclr={{Easy:'#3fb950',Medium:'#ffd166',Hard:'#f85149'}};
      ig.innerHTML=data.map(d=>`
        <div class="idea-card">
          <div class="idea-name">${{d.name}}</div>
          <div class="idea-desc">${{d.description}}</div>
          <div style="font-size:10px;color:#58a6ff;margin-bottom:4px">🎯 ${{d.target}}</div>
          <div style="font-size:10px;color:#06d6a0">💰 ${{d.revenue}}</div>
          <div style="font-size:10px;color:#8b949e;margin-top:2px">📊 Market: ${{d.market}}</div>
          <div class="score">⚙️ Difficulty: <span style="color:${{dclr[d.difficulty]||'#ffd166'}}">${{d.difficulty}}</span></div>
          <div class="idea-tags">${{(d.tags||[]).map(t=>`<span class="tag">#${{t}}</span>`).join('')}}</div>
        </div>`).join('');
    }}catch(err){{ig.innerHTML='<div style="color:#f85149;font-size:12px;padding:8px;grid-column:span 2">Parse error</div>';}}
    window.parent.postMessage({{type:'plugin_result',plugin:'idea_generator',result:e.data.result}},'*');
  }}
}});
if(document.getElementById('domainIn').value)setTimeout(generate,300);
</script></body></html>"""




# ══════════════════════════════════════════════════════════════
#  MAIN APP
# ══════════════════════════════════════════════════════════════
def show_app():
    user_name = st.session_state.user.user_metadata.get("full_name", st.session_state.user.email)

    # ── Load backend ───────────────────────────────────────────
    if not st.session_state.backend_ready:
        if not os.path.exists("./physics_index"):
            st.warning("⚠️ physics_index not found locally. The chatbot will work without RAG context.")
            # Create a dummy vectorstore so app doesn't crash
            st.session_state.vs = None
            st.session_state.Document = None
            st.session_state.backend_ready = True
            st.rerun()
        st.info("⏳ Loading knowledge base... (~30 seconds first time)")
        bar = st.progress(0, text="Starting...")
        try:
            bar.progress(20, text="Loading embeddings...")
            vs, Document = load_backend()
            bar.progress(60, text="Loading your learned answers...")
            learned = load_learned()
            if learned:
                docs = [Document(
                    page_content=f"Q: {e['question']}\nA: {e['answer']}",
                    metadata={"source": "user_learned"}
                ) for e in learned]
                vs.add_documents(docs)
            bar.progress(100, text="Ready!")
            st.session_state.vs       = vs
            st.session_state.Document = Document
            st.session_state.backend_ready = True
            st.rerun()
        except Exception as e:
            st.error(f"❌ {e}")
            st.stop()

    vs       = st.session_state.vs
    Document = st.session_state.Document

    if st.session_state.get("show_profile"):
        show_profile_page()
        return

    # ── SIDEBAR ────────────────────────────────────────────────
    with st.sidebar:
        # Theme toggle
        col_name, col_theme = st.columns([3,1])
        with col_name:
            st.markdown(f"**👋 {user_name}**")
            st.caption(st.session_state.user.email)
        with col_theme:
            if st.button("🌙" if st.session_state.dark_mode else "☀️", help="Toggle theme"):
                st.session_state.dark_mode = not st.session_state.dark_mode
                st.rerun()

        col_prof, col_out = st.columns(2)
        with col_prof:
            if st.button("👤 My Profile", use_container_width=True):
                st.session_state.show_profile = True
                # Auto-load saved profile
                if not st.session_state.user_profile:
                    saved = load_user_profile()
                    if saved:
                        st.session_state.user_profile = saved
                st.rerun()
        with col_out:
            pass
        if st.button("🚪 Sign Out", use_container_width=True):
            supabase.auth.sign_out()
            clear_session_from_js()
            for k in ["user","access_token","messages","pending_feedback","backend_ready","vs","Document"]:
                st.session_state[k] = None if k in ["user","access_token"] else \
                                      [] if k=="messages" else \
                                      False if k=="backend_ready" else None
            st.session_state.show_landing = True
            st.rerun()

        st.divider()

        # ── Numerical Solver Panel ─────────────────────────────
        # ── Creative Mode Toggle ───────────────────────────────
        st.markdown('<div class="section-label">🎛️ Mode Select</div>', unsafe_allow_html=True)
        voice_on = st.toggle("🎙️ Voice Mode", value=st.session_state.voice_mode,
            help="Speak your question, hear the answer")
        if voice_on != st.session_state.voice_mode:
            st.session_state.voice_mode = voice_on
            st.rerun()
        creative_on = st.toggle("✍️ Creative/English Mode", value=st.session_state.creative_mode,
            help="Essays, debates, letters, stories, poems")
        if creative_on != st.session_state.creative_mode:
            st.session_state.creative_mode = creative_on
            if creative_on: st.session_state.coding_mode = False
            st.rerun()

        coding_on = st.toggle("💻 Coding Mode", value=st.session_state.coding_mode,
            help="Python, Java, C++, C#, Roblox Lua, JS and more")
        if coding_on != st.session_state.coding_mode:
            st.session_state.coding_mode = coding_on
            if coding_on: st.session_state.creative_mode = False
            st.rerun()

        web_on = st.toggle("🌐 Web Search", value=st.session_state.get("web_search_mode",False),
            help="Auto-searches Wikipedia, arXiv, DuckDuckGo for extra info")
        if web_on != st.session_state.get("web_search_mode",False):
            st.session_state.web_search_mode = web_on
            st.rerun()

        if st.session_state.creative_mode:
            st.success("✍️ English/Creative Mode ON")
        elif st.session_state.coding_mode:
            st.success("💻 Coding Mode ON — Up to 2000 lines!")
        else:
            st.info("⚛️ Science Mode")
        if st.session_state.get("web_search_mode"):
            st.success("🌐 Web Search ON — Wikipedia, arXiv, DuckDuckGo active")

        st.divider()
        st.markdown('<div class="section-label">🔢 Numerical Solver</div>', unsafe_allow_html=True)
        st.caption("Paste any numerical problem for full step-by-step solution")
        num_problem = st.text_area("Problem:", placeholder="e.g. A ball is thrown with velocity 20 m/s at 30°. Find max height.", height=100, key="num_input", label_visibility="collapsed")
        if st.button("Solve Step-by-Step →", use_container_width=True):
            if num_problem.strip():
                with st.spinner("🔢 Solving..."):
                    result = solve_numerical(num_problem.strip(), vs)
                    st.session_state.solver_result = result
                    st.rerun()

        if st.session_state.solver_result:
            st.markdown('<div class="section-label">📋 Solution</div>', unsafe_allow_html=True)
            with st.expander("View Full Solution", expanded=True):
                st.markdown(st.session_state.solver_result)
            if st.button("Clear Solution"):
                st.session_state.solver_result = None
                st.rerun()

        st.divider()

        # ── Learned answers ────────────────────────────────────
        st.markdown('<div class="section-label">🧠 Learned Answers</div>', unsafe_allow_html=True)
        learned = load_learned()
        if not learned:
            st.caption("Nothing saved yet.")
        else:
            st.caption(f"✅ {len(learned)} answers saved")
            for e in learned[-3:]:
                with st.expander(f"{e['question'][:35]}..."):
                    st.write(e["answer"][:150])

        st.divider()

        # ── Past conversations ─────────────────────────────────
        st.markdown('<div class="section-label">💬 History</div>', unsafe_allow_html=True)
        past = load_past_conversations()
        if not past:
            st.caption("No history yet.")
        else:
            st.caption(f"💬 {len(past)} messages saved")
            user_msgs = [p for p in past if p["role"] == "user"]
            for p in user_msgs[-3:]:
                with st.expander(f"📅 {p['created_at'][:10]}: {p['content'][:30]}..."):
                    st.write(p["content"])

        st.divider()
        # ── Plugin Store ─────────────────────────────────────
        st.divider()
        st.markdown('<div class="section-label">🔌 Plugins</div>', unsafe_allow_html=True)
        if st.button("🔌 Plugin Store", use_container_width=True, key="open_plugin_store"):
            st.session_state.plugin_store_open = not st.session_state.get("plugin_store_open", False)
            st.rerun()
        if st.session_state.get("plugin_store_open"):
            plugin_catalog = dict(PLUGINS)
            plugin_catalog.update(st.session_state.get("custom_plugins", {}))
            for key, p in plugin_catalog.items():
                if st.button(f"{p['icon']} {p['name']}", key=f"sidebar_plug_{key}", use_container_width=True):
                    st.session_state[f"plugin_open_{key}"] = not st.session_state.get(f"plugin_open_{key}", False)
                    st.rerun()

        # ── PDF Reader ──────────────────────────────────────
        st.divider()
        st.markdown('<div class="section-label">📄 PDF Reader</div>', unsafe_allow_html=True)
        uploaded_pdf = st.file_uploader("Upload a PDF", type="pdf", key="pdf_upload",
            label_visibility="collapsed")
        if uploaded_pdf:
            if st.button("📖 Summarise PDF", use_container_width=True):
                with st.spinner("📖 Reading PDF..."):
                    pdf_text = extract_pdf_text(uploaded_pdf)
                if pdf_text:
                    summary = summarise_pdf(pdf_text)
                    if summary:
                        st.session_state.pdf_text = pdf_text
                        st.session_state.pdf_name = uploaded_pdf.name
                        reply = f"📄 **Summary of {uploaded_pdf.name}:**\n\n{summary}"
                        st.session_state.messages.append({"role":"assistant","content":reply})
                        st.rerun()
                else:
                    st.error("Could not read PDF. Make sure it contains text (not scanned images).")
            if st.session_state.get("pdf_text"):
                st.success(f"✅ {st.session_state.get('pdf_name','PDF')} loaded — ask questions about it!")

        # Past mistakes
        past_mistakes = load_quiz_mistakes()
        if past_mistakes:
            st.divider()
            st.markdown('<div class="section-label">📊 Quiz Mistakes</div>', unsafe_allow_html=True)
            st.caption(f"📝 {len(past_mistakes)} past mistakes recorded")
            if st.button("🎓 Get Lesson on Mistakes", use_container_width=True, key="lesson_btn"):
                with st.spinner("Generating lesson..."):
                    lesson = generate_mistake_lesson(past_mistakes)
                if lesson:
                    st.session_state.messages.append({"role":"assistant","content":"**📚 Personalised Lesson based on your quiz history:**\n\n"+lesson})
                    st.rerun()

        st.divider()
        c1, c2 = st.columns(2)
        with c1:
            if st.button("🗑️ Chat", use_container_width=True):
                st.session_state.messages = []
                st.session_state.pending_feedback = None
                st.session_state.simplify_target  = None
                st.rerun()
        with c2:
            if st.button("🗑️ Data", use_container_width=True):
                delete_all_data()
                st.toast("All data deleted.")
                st.rerun()

    # ── MAIN AREA ──────────────────────────────────────────────
    col_title, col_tog = st.columns([5,1])
    with col_title:
        st.markdown("## ⚛️ PhysIQ")
        st.caption("Physics · Chemistry · English Writing · Class 10 to College")
    with col_tog:
        st.write("")  # spacer

    if st.session_state.get("voice_mode"):
        st.markdown("**🎙️ Voice Panel**")
        show_voice_ui()
        st.caption("Speak into the mic, then send the transcript into the chat.")
        st.divider()

    # ── Chat history ───────────────────────────────────────────
    for i, msg in enumerate(st.session_state.messages):
        with st.chat_message(msg["role"]):
            if msg.get("is_image") and msg.get("b64"):
                try:
                    import base64
                    img_bytes = base64.b64decode(msg["b64"])
                    st.image(img_bytes, use_container_width=True)
                except Exception:
                    st.warning("This saved image could not be displayed.")
            st.write(msg["content"])

            if msg["role"] == "assistant":
                # Confidence badge
                conf = msg.get("confidence","")
                cls  = msg.get("conf_class","conf-med")
                if conf:
                    st.markdown(f'<span class="{cls}">📊 {conf}</span>', unsafe_allow_html=True)

                # Explain Simply button (only on last assistant message)
                if i == len(st.session_state.messages) - 1:
                    b1, b2, b3 = st.columns([1, 1, 3])
                    with b1:
                        if st.button("🗣️ Explain Simply", key=f"simple_{i}"):
                            st.session_state.simplify_target = {
                                "answer": msg["content"],
                                "question": st.session_state.messages[i-1]["content"] if i > 0 else ""
                            }
                            st.session_state.show_animator = False
                            st.rerun()
                    with b2:
                        if st.button("🎬 Animated Teaching", key=f"anim_{i}"):
                            q = st.session_state.messages[i-1]["content"] if i > 0 else ""
                            with st.spinner("🎬 Generating animation..."):
                                anim = generate_animation_script(q, msg["content"])
                                if anim:
                                    st.session_state.animation_data = anim
                                    st.session_state.show_animator = True
                                    st.session_state.simplify_target = None
                                else:
                                    st.warning("Could not generate animation for this topic.")
                            st.rerun()

            # Show simple explanation if triggered for this message
            if msg["role"] == "assistant" and i == len(st.session_state.messages) - 1:
                if st.session_state.simplify_target:
                    st.divider()
                    with st.spinner("🗣️ Making it simpler..."):
                        simple = explain_simply(
                            st.session_state.simplify_target["answer"],
                            st.session_state.simplify_target["question"]
                        )
                    if simple:
                        st.markdown("**🗣️ Simple Explanation:**")
                        st.info(simple)
                    col_ok, col_cl = st.columns([1,4])
                    with col_ok:
                        if st.button("✅ Got it!", key="got_it"):
                            st.session_state.simplify_target = None
                            st.rerun()

    # ── Plugin Panels ──────────────────────────────────────────
    plugin_catalog = dict(PLUGINS)
    plugin_catalog.update(st.session_state.get("custom_plugins", {}))
    active_plugins = [k for k in plugin_catalog if st.session_state.get(f"plugin_open_{k}")]
    if active_plugins:
        st.markdown("---")
        st.markdown("### 🔌 Active Plugins")
        for pk in active_plugins:
            plugin_meta = plugin_catalog[pk]
            with st.expander(f"{plugin_meta['icon']} {plugin_meta['name']}", expanded=True):
                col_close, _ = st.columns([1,5])
                with col_close:
                    if st.button("✕ Close", key=f"close_plugin_{pk}"):
                        st.session_state[f"plugin_open_{pk}"] = False
                        st.rerun()
                render_plugin(pk)

    # ── Quiz Panel ─────────────────────────────────────────────
    if st.session_state.get("quiz_active") and st.session_state.get("quiz_state"):
        col_qt, col_qc = st.columns([5,1])
        with col_qt:
            st.markdown("**🎯 Quiz Connector**")
        with col_qc:
            if st.button("✕ Exit Quiz", key="exit_quiz"):
                st.session_state.quiz_active = False
                st.session_state.quiz_state = None
                st.rerun()
        show_quiz_ui()

    # ── Animated Teaching Panel ───────────────────────────────
    if st.session_state.get("show_animator") and st.session_state.get("animation_data"):
        st.divider()
        col_at, col_cl = st.columns([5,1])
        with col_at:
            atype = st.session_state.animation_data.get("type","").upper()
            st.markdown(f"**🎬 Animated Teaching — {atype} Animation**")
        with col_cl:
            if st.button("✕ Close", key="close_anim"):
                st.session_state.show_animator = False
                st.rerun()
        if not os.path.exists("animator.html"):
            st.error("❌ animator.html not found! Make sure it is in your physics-chatbot folder.")
        else:
            show_animator_panel(st.session_state.animation_data)

    # ── Feedback ───────────────────────────────────────────────
    if st.session_state.pending_feedback and not st.session_state.simplify_target:
        pf = st.session_state.pending_feedback
        st.divider()
        st.markdown("**📝 Was this answer correct?**")
        c1, c2, c3 = st.columns(3)

        with c1:
            if st.button("✅ Yes — Save it"):
                vs.add_documents([Document(
                    page_content=f"Q: {pf['question']}\nA: {pf['answer']}",
                    metadata={"source": "user_learned"}
                )])
                save_learned(pf["question"], pf["answer"])
                st.session_state.pending_feedback = None
                st.toast("✅ Saved to your knowledge base!")
                st.rerun()
        with c2:
            if st.button("❌ Correct it"):
                st.session_state.pending_feedback["correcting"] = True
                st.rerun()
        with c3:
            if st.button("⏭️ Skip"):
                st.session_state.pending_feedback = None
                st.rerun()

        if pf.get("correcting"):
            correction = st.text_area("✏️ Type the correct answer:", height=100)
            if st.button("💾 Save Correction"):
                if correction.strip():
                    vs.add_documents([Document(
                        page_content=f"Q: {pf['question']}\nA: {correction.strip()}",
                        metadata={"source": "user_learned"}
                    )])
                    save_learned(pf["question"], correction.strip())
                    st.session_state.pending_feedback = None
                    st.toast("🧠 Correction saved!")
                    st.rerun()

    # ── Voice button (above chat bar) ──────────────────────────
    st.components.v1.html("""
<style>
*{margin:0;padding:0;box-sizing:border-box}
body{background:transparent;padding:2px 0}
.vrow{display:flex;align-items:center;gap:8px;justify-content:flex-end;padding:0 6px}
.mb{width:40px;height:40px;border-radius:50%;border:none;cursor:pointer;
  background:linear-gradient(135deg,#1f6feb,#388bfd);color:#fff;
  font-size:17px;display:flex;align-items:center;justify-content:center;
  box-shadow:0 2px 8px rgba(88,166,255,0.3);transition:all .2s}
.mb:hover{transform:scale(1.1)}
.mb.L{background:linear-gradient(135deg,#b91c1c,#ef4444);animation:pu 1s infinite}
.mb.D{background:linear-gradient(135deg,#15803d,#22c55e)}
@keyframes pu{0%{box-shadow:0 0 0 0 rgba(239,68,68,.5)}70%{box-shadow:0 0 0 12px rgba(239,68,68,0)}100%{box-shadow:0 0 0 0 rgba(239,68,68,0)}}}
.pill{background:#161b22;border:1px solid #30363d;border-radius:20px;padding:5px 12px;
  color:#e6edf3;font-size:12px;font-family:'Segoe UI',sans-serif;max-width:320px;
  white-space:nowrap;overflow:hidden;text-overflow:ellipsis;display:none}
.pill.v{display:block}
.sb{background:#238636;color:#fff;border:none;border-radius:8px;padding:5px 12px;
  font-size:12px;cursor:pointer;display:none;font-family:'Segoe UI',sans-serif}
.sb.v{display:block}
.sb:hover{background:#2ea043}
.tb{position:absolute;bottom:-3px;left:0;right:0;height:3px;background:rgba(88,166,255,.15);border-radius:2px;display:none}
.tf{height:100%;background:#58a6ff;border-radius:2px;transition:width .1s linear}
</style>
<div class="vrow">
  <div class="pill" id="pi"></div>
  <button class="sb" id="sb">➤ Send</button>
  <div style="position:relative">
    <button class="mb" id="mb" title="Click to speak">🎙️</button>
    <div class="tb" id="tb"><div class="tf" id="tf" style="width:100%"></div></div>
  </div>
</div>
<script>
const mb=document.getElementById('mb');
const pi=document.getElementById('pi');
const sb=document.getElementById('sb');
const tb=document.getElementById('tb');
const tf=document.getElementById('tf');
let R=null,on=false,txt='',sTimer=null,iTimer=null;
const SILS=2200;

function go(){
  const SR=window.SpeechRecognition||window.webkitSpeechRecognition;
  if(!SR){pi.textContent='⚠️ Use Chrome/Edge';pi.classList.add('v');return;}
  R=new SR();R.continuous=true;R.interimResults=true;R.lang='en-US';
  R.onstart=()=>{on=true;txt='';mb.className='mb L';mb.textContent='🔴';
    pi.textContent='Listening…';pi.classList.add('v');sb.classList.remove('v');
    tb.style.display='none';};
  R.onresult=(e)=>{
    let fin='',int='';
    for(let i=e.resultIndex;i<e.results.length;i++){
      if(e.results[i].isFinal)fin+=e.results[i][0].transcript;
      else int+=e.results[i][0].transcript;
    }
    if(fin)txt+=fin;
    pi.textContent='"'+(txt||int)+'"';
    // Silence countdown
    clearTimeout(sTimer);clearInterval(iTimer);
    let p=100;tf.style.width='100%';tb.style.display='block';
    iTimer=setInterval(()=>{p-=(100/SILS)*100;tf.style.width=Math.max(0,p)+'%';if(p<=0)clearInterval(iTimer);},100);
    sTimer=setTimeout(()=>{if(txt.trim())send();else R.stop();},SILS);
  };
  R.onend=()=>{on=false;tb.style.display='none';
    if(txt.trim()){mb.className='mb D';mb.textContent='✅';sb.classList.add('v');}
    else{mb.className='mb';mb.textContent='🎙️';pi.classList.remove('v');}};
  R.onerror=(e)=>{on=false;mb.className='mb';mb.textContent='🎙️';
    const m={'not-allowed':'🔒 Allow mic','no-speech':'No speech detected'};
    pi.textContent=m[e.error]||'Error: '+e.error;setTimeout(()=>pi.classList.remove('v'),2500);};
  R.start();
}

function send(){
  clearTimeout(sTimer);clearInterval(iTimer);
  if(R&&on)R.stop();
  if(!txt.trim())return;
  const q=txt.trim();
  // Write to URL so Streamlit can read it
  const u=new URL(window.parent.location.href);
  u.searchParams.set('vq',encodeURIComponent(q));
  window.parent.history.replaceState({},'',u.toString());
  pi.textContent='⏳ Sending…';sb.classList.remove('v');
  mb.className='mb';mb.textContent='🎙️';txt='';
  // Trigger page refresh
  setTimeout(()=>{window.parent.location.href=u.toString();},200);
}

mb.onclick=()=>{if(on){clearTimeout(sTimer);send();}else go();};
sb.onclick=()=>send();

// TTS reply
const syn=window.speechSynthesis;let vv=[];
function lv(){vv=syn.getVoices();}syn.onvoiceschanged=lv;lv();
window.addEventListener('message',e=>{
  if(e.data&&e.data.type==='physiq_tts'){
    syn.cancel();
    const u2=new SpeechSynthesisUtterance(e.data.text);
    u2.rate=1.0;u2.pitch=1.0;
    const pref=['Google UK English Female','Microsoft Libby','Samantha','Google US English'];
    for(const p of pref){const v=vv.find(x=>x.name===p);if(v){u2.voice=v;break;}}
    syn.speak(u2);
  }
});
</script>
""", height=56, scrolling=False)

    # Read voice question from URL query param
    _vq = st.query_params.get("vq","")
    if _vq and _vq != st.session_state.get("_last_vq",""):
        st.session_state._voice_question = _vq
        st.session_state._last_vq = _vq
        try: del st.query_params["vq"]
        except: pass

    # ── Chat input ─────────────────────────────────────────────
    voice_q = st.session_state.pop("_voice_question", None)
    if voice_q:
        question = voice_q
    else:
        question = st.chat_input("Ask about Physics, Chemistry, or English Writing...")

    if question:
        st.session_state.simplify_target = None

        # ── Voice question flag ──────────────────────────────
        is_voice_q = st.session_state.get("_is_voice_question", False)
        if is_voice_q:
            st.session_state._is_voice_question = False

        # ── Auto-activate plugins based on question ───────────
        suggested = ai_select_plugins(question, vs)
        for pk in suggested:
            if not st.session_state.get(f"plugin_open_{pk}"):
                st.session_state[f"plugin_open_{pk}"] = True
            if pk == "multi_agent":
                st.session_state.multi_agent_input = question

        # ── Auto-update profile every 20 messages ───────────
        if len(st.session_state.messages) % 20 == 0 and len(st.session_state.messages) > 0:
            import threading
            def bg_profile():
                try:
                    past = load_past_conversations()
                    mistakes = load_quiz_mistakes()
                    p = analyse_user_profile(st.session_state.messages, past, mistakes)
                    if p:
                        st.session_state.user_profile = p
                        save_user_profile(p)
                except: pass
            threading.Thread(target=bg_profile, daemon=True).start()

        # ── PDF question ─────────────────────────────────────
        if st.session_state.get("pdf_text"):
            pdf_keywords = ["in the pdf","from the pdf","the document","the pdf","page","according to","it says","uploaded"]
            if any(k in question.lower() for k in pdf_keywords):
                st.session_state.messages.append({"role":"user","content":question})
                save_message("user", question)
                with st.chat_message("user"):
                    st.write(question)
                with st.chat_message("assistant"):
                    with st.spinner("📄 Reading PDF to answer..."):
                        pdf_ans = answer_from_pdf(st.session_state.pdf_text, question)
                    if pdf_ans:
                        st.markdown(pdf_ans)
                        st.session_state.messages.append({"role":"assistant","content":pdf_ans,
                            "confidence":"📄 PDF Answer","conf_class":"conf-high"})
                        save_message("assistant", pdf_ans, "pdf")
                st.rerun()

        # ── Image generation request ─────────────────────────
        if is_image_request(question) or st.session_state.get("_regen_prompt"):
            regen = st.session_state.get("_regen_prompt")
            regen_style = st.session_state.get("_regen_style", "general")
            regen_q = st.session_state.get("_regen_question", question)
            if regen:
                st.session_state._regen_prompt = None
                st.session_state._regen_style = None
                st.session_state._regen_question = None
                # Re-generate with same prompt
                act_question = regen_q
                st.session_state.messages.append({"role":"user","content":f"[Regenerate] {regen_q}"})
                save_message("user", f"[Regenerate] {regen_q}")
                with st.chat_message("user"):
                    st.write(f"🔄 Regenerating: {regen_q}")
                with st.chat_message("assistant"):
                    style, w, h = detect_image_style(regen_q)
                    final_p = enhance_prompt_for_style(regen, style)
                    with st.spinner("🖼️ Regenerating image..."):
                        b64, err = generate_image_hf(final_p, w, h)
                    if b64:
                        show_generated_image(b64, final_p, style, regen_q)
                        reply = f"🎨 Regenerated! Style: **{style.upper()}** | Resolution: {w}×{h}px"
                        st.session_state.messages.append({"role":"assistant","content":reply,"is_image":True,"b64":b64})
                        save_message("assistant", reply)
                    else:
                        st.error(f"❌ {err}")
                st.rerun()
            else:
                st.session_state.messages.append({"role":"user","content":question})
                save_message("user", question)
                with st.chat_message("user"):
                    st.write(question)
                with st.chat_message("assistant"):
                    b64, final_prompt = handle_image_request(question, vs)
                    if b64:
                        show_generated_image(b64, final_prompt, detect_image_style(question)[0], question)
                        _, w, h = detect_image_style(question)
                        style = detect_image_style(question)[0]
                        reply = f"🎨 Generated! Style: **{style.upper()}** | Resolution: {w}×{h}px using FLUX.1"
                        st.session_state.messages.append({
                            "role":"assistant","content":reply,
                            "is_image":True,"b64":b64,"prompt":final_prompt
                        })
                        save_message("assistant", reply)
                st.rerun()

        # ── Detect quiz request ───────────────────────────────
        if is_quiz_request(question):
            st.session_state.messages.append({"role": "user", "content": question})
            save_message("user", question)
            with st.chat_message("user"):
                st.write(question)
            q_type, num_q = extract_quiz_topic(question)
            topic = question.lower()
            for word in QUIZ_TRIGGER_WORDS:
                topic = topic.replace(word,"").strip()
            topic = topic.strip(" on about:,") or "Physics and Chemistry"
            with st.chat_message("assistant"):
                with st.spinner(f"📝 Generating {num_q}-question quiz on {topic}..."):
                    questions = generate_quiz(topic, q_type, num_q)
                if questions:
                    reply = f"I have generated a {num_q}-question quiz on {topic}! The quiz is open below. Good luck!"
                    st.write(reply)
                    st.session_state.messages.append({"role":"assistant","content":reply})
                    save_message("assistant", reply)
                    st.session_state.quiz_state = {
                        "questions": questions, "current_idx": 0, "score": 0,
                        "answers": {}, "skipped": [], "phase": "quiz",
                        "topic": topic, "last_correct": None
                    }
                    st.session_state.quiz_active = True
                else:
                    st.error("Could not generate quiz. Please try again with a specific topic.")
            st.rerun()
        else:
            # ── Coding request ────────────────────────────────
            if is_coding_request(question):
                lang = detect_language(question)
                st.session_state.messages.append({"role": "user", "content": question})
                save_message("user", question)
                with st.chat_message("user"):
                    st.write(question)
                with st.chat_message("assistant"):
                    with st.spinner(f"💻 Writing {lang} code (up to 2000 lines)..."):
                        code_answer = generate_code(question, vs, lang)
                    if code_answer:
                        st.markdown(code_answer)
                        # ── Code PDF download + app integrations ──
                        col_pdf, col_int = st.columns([1,3])
                        with col_pdf:
                            show_code_download_button(code_answer, lang)
                        with col_int:
                            with st.expander("📤 Export to apps"):
                                show_app_integration_buttons(code_answer, lang)
                        st.session_state.messages.append({
                            "role":"assistant","content":code_answer,
                            "confidence":"💻 Code","conf_class":"conf-high"
                        })
                        save_message("assistant", code_answer, "code")
                        st.session_state.pending_feedback = {"question":question,"answer":code_answer}
                st.rerun()
            else:
                # ── Normal question ─────────────────────────────
                st.session_state.messages.append({"role": "user", "content": question})
        save_message("user", question)

        with st.chat_message("user"):
            st.write(question)

        history_text = "\n".join([
            f"{'Student' if m['role']=='user' else 'Tutor'}: {m['content']}"
            for m in st.session_state.messages[-6:]
        ])

        # ── Live data check ──────────────────────────────────
        live_type = needs_live_data(question)
        if live_type:
            live_result = fetch_live_data(question, live_type)
            if live_result:
                st.session_state.messages.append({"role":"user","content":question})
                save_message("user",question)
                with st.chat_message("user"): st.write(question)
                with st.chat_message("assistant"):
                    st.markdown(live_result)
                    # Also ask AI to explain
                    with st.spinner("🤖 Adding AI context..."):
                        ai_ctx = call_hf(
                            "You are a helpful tutor. Briefly explain or add context to this live data in 2-3 sentences.",
                            f"Data: {live_result}. User asked: {question}. Add brief context:",
                            max_tokens=200
                        )
                    if ai_ctx: st.markdown(ai_ctx)
                    reply = live_result + ("\n\n" + ai_ctx if ai_ctx else "")
                    st.session_state.messages.append({"role":"assistant","content":reply})
                    save_message("assistant",reply,"live_data")
                st.rerun()

        with st.chat_message("assistant"):
            use_web = st.session_state.get("web_search_mode", False)
            with st.spinner("🔁 Thinking..." + (" + 🌐 searching web..." if use_web else "")):
                if use_web:
                    answer, sources = answer_with_web_search(question, vs, history_text)
                    conf_label = "🌐 Web + AI"
                    conf_class = "conf-high"
                    if not answer:
                        answer, conf_label, conf_class = ask_question(question, vs, history_text)
                        sources = []
                else:
                    answer, conf_label, conf_class = ask_question(question, vs, history_text)
                    sources = []
                    # Auto web search on low confidence
                    if conf_class == "conf-low" and needs_web_search(question, conf_class):
                        with st.spinner("🌐 Confidence low — checking web..."):
                            web_r = web_search_all(question)
                            if web_r:
                                web_ctx = format_web_results(web_r)
                                sources = list(web_r.keys())
                                # Re-ask with web context
                                sys2 = "You are an expert tutor. Use the web results to improve your answer."
                                usr2 = f"Web results:{web_ctx}\nQuestion:{question}\nImprove this answer:{answer}"
                                better = call_hf(sys2, usr2, max_tokens=1200)
                                if better: answer = better; conf_label = "🌐 Web-enhanced"
                st.write(answer)
                if sources:
                    st.caption(f"📡 Sources: {', '.join(sources)}")
                st.markdown(f'<span class="{conf_class}">📊 {conf_label}</span>', unsafe_allow_html=True)

                st.session_state.messages.append({
                    "role": "assistant", "content": answer,
                    "confidence": conf_label, "conf_class": conf_class
                })
                save_message("assistant", answer, conf_label)
                st.session_state.pending_feedback = {"question": question, "answer": answer}

                # TTS voice reply if voice mode on
                if st.session_state.get("voice_mode"):
                    with st.spinner("🔊 Preparing voice..."):
                        short = call_hf(
                            "Convert this answer to a friendly spoken response under 50 words. No symbols, no markdown, natural speech only.",
                            f"Q:{question}\nA:{answer[:600]}\nSpoken version:",
                            max_tokens=120
                        )
                    if short:
                        st.components.v1.html(f"""<script>
setTimeout(()=>{{window.parent.postMessage({{type:'physiq_tts',text:{repr(short)}}}, '*');}},400);
</script>""", height=0)

        st.rerun()

# ══════════════════════════════════════════════════════════════
#  ROUTER
# ══════════════════════════════════════════════════════════════
if st.session_state.user is None or st.session_state.show_landing:
    show_landing()
else:
    show_app()
